"""
문서 보안 해제 저장 도구
- PowerPoint, Excel, 한글 문서를 COM으로 읽어서 새 파일로 저장
- DRM이 저장을 차단해도 내용을 추출하여 새 문서 생성
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import tempfile
import shutil
from pathlib import Path

# COM 관련
try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

# PPT 생성용
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Excel 생성용
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


class DocumentExtractor:
    """문서 추출기 메인 클래스"""

    def __init__(self):
        self.root = tk.Tk()
        self.root.title("문서 보안 해제 저장 도구")
        self.root.geometry("500x400")
        self.root.resizable(False, False)

        # 상태 변수
        self.current_doc_name = tk.StringVar(value="감지 중...")
        self.current_doc_type = tk.StringVar(value="-")
        self.save_path = tk.StringVar(value="")
        self.status_text = tk.StringVar(value="프로그램 시작됨")
        self.progress_var = tk.DoubleVar(value=0)

        self.setup_ui()
        self.detect_open_document()

    def setup_ui(self):
        """UI 구성"""
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # 제목
        title_label = ttk.Label(main_frame, text="문서 보안 해제 저장 도구",
                                font=("맑은 고딕", 16, "bold"))
        title_label.pack(pady=(0, 20))

        # 문서 정보 프레임
        info_frame = ttk.LabelFrame(main_frame, text="현재 열린 문서", padding="10")
        info_frame.pack(fill=tk.X, pady=(0, 15))

        # 문서명
        doc_frame = ttk.Frame(info_frame)
        doc_frame.pack(fill=tk.X, pady=2)
        ttk.Label(doc_frame, text="문서명:", width=10).pack(side=tk.LEFT)
        ttk.Label(doc_frame, textvariable=self.current_doc_name,
                  font=("맑은 고딕", 10)).pack(side=tk.LEFT, fill=tk.X, expand=True)

        # 문서 종류
        type_frame = ttk.Frame(info_frame)
        type_frame.pack(fill=tk.X, pady=2)
        ttk.Label(type_frame, text="종류:", width=10).pack(side=tk.LEFT)
        ttk.Label(type_frame, textvariable=self.current_doc_type,
                  font=("맑은 고딕", 10, "bold")).pack(side=tk.LEFT)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=self.detect_open_document).pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = ttk.LabelFrame(main_frame, text="저장 설정", padding="10")
        path_frame.pack(fill=tk.X, pady=(0, 15))

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Label(path_inner, text="저장 경로:").pack(side=tk.LEFT)
        ttk.Entry(path_inner, textvariable=self.save_path, width=35).pack(side=tk.LEFT, padx=(5, 5))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_save_path).pack(side=tk.LEFT)

        # 저장 버튼
        self.save_button = ttk.Button(main_frame, text="새 파일로 저장",
                                       command=self.start_extraction,
                                       style="Accent.TButton")
        self.save_button.pack(pady=15)

        # 진행바
        self.progress = ttk.Progressbar(main_frame, variable=self.progress_var,
                                         maximum=100, length=400)
        self.progress.pack(pady=(0, 10))

        # 상태 표시
        status_frame = ttk.Frame(main_frame)
        status_frame.pack(fill=tk.X)
        ttk.Label(status_frame, text="상태:").pack(side=tk.LEFT)
        ttk.Label(status_frame, textvariable=self.status_text,
                  font=("맑은 고딕", 9)).pack(side=tk.LEFT, padx=(5, 0))

        # 스타일 설정
        style = ttk.Style()
        style.configure("Accent.TButton", font=("맑은 고딕", 11, "bold"))

    def browse_save_path(self):
        """저장 경로 선택"""
        doc_type = self.current_doc_type.get()

        if doc_type == "PowerPoint":
            filetypes = [("PowerPoint 파일", "*.pptx")]
            default_ext = ".pptx"
        elif doc_type == "Excel":
            filetypes = [("Excel 파일", "*.xlsx")]
            default_ext = ".xlsx"
        elif doc_type == "한글":
            filetypes = [("한글 파일", "*.hwp"), ("한글 2014+", "*.hwpx")]
            default_ext = ".hwp"
        else:
            filetypes = [("모든 파일", "*.*")]
            default_ext = ""

        path = filedialog.asksaveasfilename(
            defaultextension=default_ext,
            filetypes=filetypes,
            title="저장할 위치 선택"
        )
        if path:
            self.save_path.set(path)

    def detect_open_document(self):
        """열려있는 문서 감지"""
        self.status_text.set("문서 감지 중...")
        self.current_doc_name.set("감지 중...")
        self.current_doc_type.set("-")

        # 백그라운드에서 감지
        thread = threading.Thread(target=self._detect_documents)
        thread.daemon = True
        thread.start()

    def _detect_documents(self):
        """문서 감지 (백그라운드)"""
        pythoncom.CoInitialize()

        detected = False

        # PowerPoint 확인
        try:
            ppt = win32com.client.GetObject(Class="PowerPoint.Application")
            if ppt.Presentations.Count > 0:
                presentation = ppt.ActivePresentation
                name = presentation.Name
                self.root.after(0, lambda: self.current_doc_name.set(name))
                self.root.after(0, lambda: self.current_doc_type.set("PowerPoint"))
                detected = True
        except:
            pass

        # Excel 확인
        if not detected:
            try:
                excel = win32com.client.GetObject(Class="Excel.Application")
                if excel.Workbooks.Count > 0:
                    workbook = excel.ActiveWorkbook
                    name = workbook.Name
                    self.root.after(0, lambda: self.current_doc_name.set(name))
                    self.root.after(0, lambda: self.current_doc_type.set("Excel"))
                    detected = True
            except:
                pass

        # 한글 확인
        if not detected:
            try:
                hwp = win32com.client.GetObject(Class="HWPFrame.HwpObject")
                # 한글은 파일명 가져오는 방식이 다름
                path = hwp.Path
                name = os.path.basename(path) if path else "제목 없음"
                self.root.after(0, lambda: self.current_doc_name.set(name))
                self.root.after(0, lambda: self.current_doc_type.set("한글"))
                detected = True
            except:
                pass

        if not detected:
            self.root.after(0, lambda: self.current_doc_name.set("열린 문서 없음"))
            self.root.after(0, lambda: self.current_doc_type.set("-"))
            self.root.after(0, lambda: self.status_text.set("문서를 먼저 열어주세요"))
        else:
            self.root.after(0, lambda: self.status_text.set("문서 감지 완료"))

        pythoncom.CoUninitialize()

    def start_extraction(self):
        """추출 시작"""
        if not self.save_path.get():
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        doc_type = self.current_doc_type.get()
        if doc_type == "-":
            messagebox.showwarning("경고", "열린 문서가 없습니다.")
            return

        self.save_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(target=self._extract_and_save)
        thread.daemon = True
        thread.start()

    def _extract_and_save(self):
        """추출 및 저장 (백그라운드)"""
        pythoncom.CoInitialize()

        try:
            doc_type = self.current_doc_type.get()
            save_path = self.save_path.get()

            if doc_type == "PowerPoint":
                self._extract_powerpoint(save_path)
            elif doc_type == "Excel":
                self._extract_excel(save_path)
            elif doc_type == "한글":
                self._extract_hwp(save_path)

            self.root.after(0, lambda: messagebox.showinfo("완료", f"저장 완료!\n{save_path}"))

        except Exception as e:
            self.root.after(0, lambda: messagebox.showerror("오류", f"저장 중 오류 발생:\n{str(e)}"))

        finally:
            self.root.after(0, lambda: self.save_button.config(state=tk.NORMAL))
            self.root.after(0, lambda: self.progress_var.set(0))
            pythoncom.CoUninitialize()

    def _extract_powerpoint(self, save_path):
        """PowerPoint 추출"""
        self.root.after(0, lambda: self.status_text.set("PowerPoint 추출 중..."))

        ppt_app = win32com.client.GetObject(Class="PowerPoint.Application")
        presentation = ppt_app.ActivePresentation

        # 방법 1: SaveCopyAs 시도 (가장 깔끔)
        try:
            self.root.after(0, lambda: self.status_text.set("SaveCopyAs 시도 중..."))
            presentation.SaveCopyAs(save_path)
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("SaveCopyAs 성공!"))
            return
        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"SaveCopyAs 실패, 수동 추출 시도..."))

        # 방법 2: Export 후 재구성
        try:
            self.root.after(0, lambda: self.status_text.set("슬라이드 Export 시도 중..."))
            temp_dir = tempfile.mkdtemp()

            total_slides = presentation.Slides.Count

            # 새 프레젠테이션 생성
            new_prs = Presentation()
            # 슬라이드 크기 설정 (EMU 단위)
            new_prs.slide_width = Emu(presentation.PageSetup.SlideWidth * 914400 / 72)
            new_prs.slide_height = Emu(presentation.PageSetup.SlideHeight * 914400 / 72)

            for i in range(1, total_slides + 1):
                progress = (i / total_slides) * 100
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                self.root.after(0, lambda n=i, t=total_slides: self.status_text.set(f"슬라이드 {n}/{t} 처리 중..."))

                slide = presentation.Slides(i)

                # 이미지로 내보내기
                img_path = os.path.join(temp_dir, f"slide_{i}.png")
                slide.Export(img_path, "PNG", 1920, 1080)

                # 새 슬라이드에 이미지 추가
                blank_layout = new_prs.slide_layouts[6]  # 빈 레이아웃
                new_slide = new_prs.slides.add_slide(blank_layout)

                # 이미지를 슬라이드 전체에 삽입
                new_slide.shapes.add_picture(
                    img_path,
                    Emu(0), Emu(0),
                    new_prs.slide_width, new_prs.slide_height
                )

            new_prs.save(save_path)

            # 임시 파일 정리
            shutil.rmtree(temp_dir)

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("Export 방식으로 저장 완료!"))
            return

        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"Export 실패: {str(e)}"))
            raise

    def _extract_excel(self, save_path):
        """Excel 추출"""
        self.root.after(0, lambda: self.status_text.set("Excel 추출 중..."))

        excel_app = win32com.client.GetObject(Class="Excel.Application")
        workbook = excel_app.ActiveWorkbook

        # 방법 1: SaveCopyAs 시도
        try:
            self.root.after(0, lambda: self.status_text.set("SaveCopyAs 시도 중..."))
            workbook.SaveCopyAs(save_path)
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("SaveCopyAs 성공!"))
            return
        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"SaveCopyAs 실패, 수동 추출 시도..."))

        # 방법 2: 셀 데이터 직접 복사
        try:
            new_wb = Workbook()
            new_wb.remove(new_wb.active)  # 기본 시트 제거

            total_sheets = workbook.Sheets.Count

            for sheet_idx in range(1, total_sheets + 1):
                sheet = workbook.Sheets(sheet_idx)
                sheet_name = sheet.Name

                self.root.after(0, lambda n=sheet_name: self.status_text.set(f"시트 '{n}' 처리 중..."))

                new_sheet = new_wb.create_sheet(title=sheet_name)

                # 사용 범위 확인
                used_range = sheet.UsedRange
                if used_range is None:
                    continue

                rows = used_range.Rows.Count
                cols = used_range.Columns.Count
                start_row = used_range.Row
                start_col = used_range.Column

                # 데이터 복사
                for r in range(rows):
                    progress = ((sheet_idx - 1) / total_sheets + (r / rows) / total_sheets) * 100
                    self.root.after(0, lambda p=progress: self.progress_var.set(p))

                    for c in range(cols):
                        try:
                            cell = sheet.Cells(start_row + r, start_col + c)
                            value = cell.Value

                            if value is not None:
                                new_cell = new_sheet.cell(row=r+1, column=c+1, value=value)

                                # 서식 복사 시도
                                try:
                                    font = cell.Font
                                    new_cell.font = Font(
                                        name=font.Name,
                                        size=font.Size,
                                        bold=font.Bold,
                                        italic=font.Italic
                                    )
                                except:
                                    pass
                        except:
                            continue

                # 열 너비 조정
                for c in range(1, cols + 1):
                    try:
                        width = sheet.Columns(start_col + c - 1).ColumnWidth
                        new_sheet.column_dimensions[get_column_letter(c)].width = width
                    except:
                        pass

            new_wb.save(save_path)
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("수동 추출 방식으로 저장 완료!"))

        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"추출 실패: {str(e)}"))
            raise

    def _extract_hwp(self, save_path):
        """한글 추출"""
        self.root.after(0, lambda: self.status_text.set("한글 문서 추출 중..."))

        try:
            hwp = win32com.client.GetObject(Class="HWPFrame.HwpObject")
        except:
            # GetObject 실패시 새로 연결 시도
            hwp = win32com.client.Dispatch("HWPFrame.HwpObject")

        # 방법 1: SaveAs 시도
        try:
            self.root.after(0, lambda: self.status_text.set("SaveAs 시도 중..."))

            # 한글 저장 형식
            if save_path.endswith('.hwpx'):
                hwp.SaveAs(save_path, "HWPX")
            else:
                hwp.SaveAs(save_path, "HWP")

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("SaveAs 성공!"))
            return
        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"SaveAs 실패: {str(e)}"))

        # 방법 2: 클립보드를 통한 복사 시도
        try:
            self.root.after(0, lambda: self.status_text.set("클립보드 복사 시도 중..."))

            # 전체 선택
            hwp.HAction.Run("SelectAll")
            # 복사
            hwp.HAction.Run("Copy")

            # 새 문서 생성
            hwp.HAction.Run("FileNew")
            # 붙여넣기
            hwp.HAction.Run("Paste")

            # 저장
            if save_path.endswith('.hwpx'):
                hwp.SaveAs(save_path, "HWPX")
            else:
                hwp.SaveAs(save_path, "HWP")

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("클립보드 방식으로 저장 완료!"))

        except Exception as e:
            self.root.after(0, lambda: self.status_text.set(f"추출 실패: {str(e)}"))
            raise

    def run(self):
        """프로그램 실행"""
        self.root.mainloop()


def check_dependencies():
    """의존성 확인"""
    missing = []

    if not HAS_WIN32COM:
        missing.append("pywin32")
    if not HAS_PPTX:
        missing.append("python-pptx")
    if not HAS_OPENPYXL:
        missing.append("openpyxl")

    if missing:
        msg = f"필요한 패키지가 없습니다:\n{', '.join(missing)}\n\n"
        msg += "설치 명령어:\npip install " + " ".join(missing)

        root = tk.Tk()
        root.withdraw()
        messagebox.showerror("의존성 오류", msg)
        sys.exit(1)


if __name__ == "__main__":
    check_dependencies()
    app = DocumentExtractor()
    app.run()
