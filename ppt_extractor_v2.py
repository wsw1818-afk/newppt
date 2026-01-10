"""
PPT 보안 해제 도구 v2
- COM으로 데이터 읽기
- python-pptx로 직접 파일 생성 (DRM 우회)
- 바탕화면에 상세 로그 저장
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import tempfile
import shutil
import datetime
import traceback
import io

# python-pptx 관련
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# COM 관련
try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False


class Logger:
    """바탕화면에 로그 파일 저장"""

    def __init__(self):
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        self.log_path = os.path.join(desktop, f"PPTExtractor_Log_{timestamp}.txt")
        self.log_file = open(self.log_path, "w", encoding="utf-8")
        self.log(f"=== PPT Extractor v2 로그 시작 ===")
        self.log(f"로그 파일: {self.log_path}")
        self.log(f"시작 시간: {datetime.datetime.now()}")
        self.log(f"python-pptx 사용: DRM 우회 모드")
        self.log("")

    def log(self, message):
        timestamp = datetime.datetime.now().strftime("%H:%M:%S.%f")[:-3]
        line = f"[{timestamp}] {message}"
        print(line)
        self.log_file.write(line + "\n")
        self.log_file.flush()

    def error(self, message, exception=None):
        self.log(f"[ERROR] {message}")
        if exception:
            self.log(f"[ERROR] 예외 타입: {type(exception).__name__}")
            self.log(f"[ERROR] 예외 메시지: {str(exception)}")
            self.log(f"[ERROR] 상세 트레이스백:")
            for line in traceback.format_exc().split("\n"):
                self.log(f"        {line}")

    def close(self):
        self.log("")
        self.log(f"=== 로그 종료: {datetime.datetime.now()} ===")
        self.log_file.close()


class PPTExtractorV2:
    """PPT 추출기 v2 - python-pptx 사용"""

    def __init__(self):
        self.logger = Logger()
        self.logger.log("PPTExtractor v2 초기화 시작")

        self.root = tk.Tk()
        self.root.title("PPT 보안 해제 도구 v2")
        self.root.geometry("550x500")
        self.root.resizable(False, False)

        # 상태 변수
        self.current_doc_name = tk.StringVar(value="감지 중...")
        self.slide_count = tk.StringVar(value="-")
        self.save_path = tk.StringVar(value="")
        self.status_text = tk.StringVar(value="프로그램 시작됨")
        self.progress_var = tk.DoubleVar(value=0)

        self.setup_ui()
        self.detect_open_ppt()

        self.logger.log("PPTExtractor v2 초기화 완료")

    def setup_ui(self):
        """UI 구성"""
        self.logger.log("UI 구성 시작")

        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # 제목
        title_label = ttk.Label(main_frame, text="PPT 보안 해제 도구 v2",
                                font=("맑은 고딕", 16, "bold"))
        title_label.pack(pady=(0, 10))

        # 설명
        desc_label = ttk.Label(main_frame,
                               text="COM으로 데이터 읽기 + python-pptx로 직접 파일 생성\n(PowerPoint 저장 기능 우회로 DRM 완전 우회)",
                               font=("맑은 고딕", 9), justify=tk.CENTER)
        desc_label.pack(pady=(0, 15))

        # 문서 정보 프레임
        info_frame = ttk.LabelFrame(main_frame, text="현재 열린 PPT", padding="10")
        info_frame.pack(fill=tk.X, pady=(0, 10))

        # 문서명
        doc_frame = ttk.Frame(info_frame)
        doc_frame.pack(fill=tk.X, pady=2)
        ttk.Label(doc_frame, text="문서명:", width=12).pack(side=tk.LEFT)
        ttk.Label(doc_frame, textvariable=self.current_doc_name,
                  font=("맑은 고딕", 10)).pack(side=tk.LEFT, fill=tk.X, expand=True)

        # 슬라이드 수
        slide_frame = ttk.Frame(info_frame)
        slide_frame.pack(fill=tk.X, pady=2)
        ttk.Label(slide_frame, text="슬라이드 수:", width=12).pack(side=tk.LEFT)
        ttk.Label(slide_frame, textvariable=self.slide_count,
                  font=("맑은 고딕", 10, "bold")).pack(side=tk.LEFT)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=self.detect_open_ppt).pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = ttk.LabelFrame(main_frame, text="새 파일 저장 위치", padding="10")
        path_frame.pack(fill=tk.X, pady=(0, 15))

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.save_path, width=45).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_save_path).pack(side=tk.LEFT)

        # 옵션 프레임
        option_frame = ttk.LabelFrame(main_frame, text="추출 옵션", padding="10")
        option_frame.pack(fill=tk.X, pady=(0, 15))

        self.copy_shapes = tk.BooleanVar(value=True)
        self.copy_texts = tk.BooleanVar(value=True)
        self.copy_images = tk.BooleanVar(value=True)

        ttk.Checkbutton(option_frame, text="도형 복사", variable=self.copy_shapes).pack(anchor=tk.W)
        ttk.Checkbutton(option_frame, text="텍스트 복사", variable=self.copy_texts).pack(anchor=tk.W)
        ttk.Checkbutton(option_frame, text="이미지 복사 (임시 파일로 추출)", variable=self.copy_images).pack(anchor=tk.W)

        # 추출 버튼
        self.extract_button = ttk.Button(main_frame, text="새 PPT로 추출 (DRM 우회)",
                                          command=self.start_extraction,
                                          style="Accent.TButton")
        self.extract_button.pack(pady=10)

        # 진행바
        self.progress = ttk.Progressbar(main_frame, variable=self.progress_var,
                                         maximum=100, length=450)
        self.progress.pack(pady=(0, 10))

        # 상태 표시
        status_frame = ttk.Frame(main_frame)
        status_frame.pack(fill=tk.X)
        ttk.Label(status_frame, text="상태:").pack(side=tk.LEFT)
        ttk.Label(status_frame, textvariable=self.status_text,
                  font=("맑은 고딕", 9)).pack(side=tk.LEFT, padx=(5, 0))

        # 스타일 설정
        style = ttk.Style()
        style.configure("Accent.TButton", font=("맑은 고딕", 11, "bold"))

        self.logger.log("UI 구성 완료")

    def browse_save_path(self):
        """저장 경로 선택"""
        self.logger.log("저장 경로 선택 대화상자 열기")

        doc_name = self.current_doc_name.get()
        if doc_name and doc_name != "감지 중..." and doc_name != "열린 PPT 없음":
            default_name = os.path.splitext(doc_name)[0] + "_복사본.pptx"
        else:
            default_name = "새문서.pptx"

        path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint 파일", "*.pptx")],
            initialfile=default_name,
            title="저장할 위치 선택"
        )
        if path:
            self.save_path.set(path)
            self.logger.log(f"저장 경로 선택됨: {path}")
        else:
            self.logger.log("저장 경로 선택 취소됨")

    def detect_open_ppt(self):
        """열려있는 PPT 감지"""
        self.logger.log("PPT 감지 시작")
        self.status_text.set("PPT 감지 중...")
        self.current_doc_name.set("감지 중...")
        self.slide_count.set("-")

        thread = threading.Thread(target=self._detect_ppt)
        thread.daemon = True
        thread.start()

    def _detect_ppt(self):
        """PPT 감지 (백그라운드)"""
        self.logger.log("백그라운드 PPT 감지 스레드 시작")
        pythoncom.CoInitialize()

        try:
            self.logger.log("PowerPoint.Application 객체 연결 시도")
            ppt = win32com.client.GetObject(Class="PowerPoint.Application")
            self.logger.log(f"PowerPoint 연결 성공, 열린 프레젠테이션 수: {ppt.Presentations.Count}")

            if ppt.Presentations.Count > 0:
                presentation = ppt.ActivePresentation
                name = presentation.Name
                count = presentation.Slides.Count

                self.logger.log(f"활성 프레젠테이션: {name}")
                self.logger.log(f"슬라이드 수: {count}")
                self.logger.log(f"파일 경로: {presentation.FullName}")

                self.root.after(0, lambda: self.current_doc_name.set(name))
                self.root.after(0, lambda: self.slide_count.set(f"{count}장"))
                self.root.after(0, lambda: self.status_text.set("PPT 감지 완료"))
            else:
                self.logger.log("열린 프레젠테이션 없음")
                self.root.after(0, lambda: self.current_doc_name.set("열린 PPT 없음"))
                self.root.after(0, lambda: self.slide_count.set("-"))
                self.root.after(0, lambda: self.status_text.set("PPT를 먼저 열어주세요"))

        except Exception as e:
            self.logger.error("PPT 감지 실패", e)
            self.root.after(0, lambda: self.current_doc_name.set("열린 PPT 없음"))
            self.root.after(0, lambda: self.slide_count.set("-"))
            self.root.after(0, lambda: self.status_text.set(f"감지 실패: {str(e)[:30]}"))

        pythoncom.CoUninitialize()
        self.logger.log("백그라운드 PPT 감지 스레드 종료")

    def start_extraction(self):
        """추출 시작"""
        self.logger.log("추출 시작 버튼 클릭")

        if not self.save_path.get():
            self.logger.log("저장 경로 미선택 - 경고 표시")
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        if self.current_doc_name.get() == "열린 PPT 없음":
            self.logger.log("열린 PPT 없음 - 경고 표시")
            messagebox.showwarning("경고", "열린 PPT가 없습니다.")
            return

        self.logger.log(f"추출 시작: 저장 경로 = {self.save_path.get()}")
        self.extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(target=self._extract_ppt)
        thread.daemon = True
        thread.start()

    def _extract_ppt(self):
        """PPT 추출 (백그라운드) - python-pptx 사용"""
        self.logger.log("=== PPT 추출 프로세스 시작 (python-pptx 모드) ===")
        pythoncom.CoInitialize()

        temp_dir = None

        try:
            save_path = self.save_path.get()
            self.logger.log(f"저장 경로: {save_path}")

            # 원본 PPT 가져오기 (COM)
            self.logger.log("원본 PPT 연결 시도...")
            self.root.after(0, lambda: self.status_text.set("원본 PPT 연결 중..."))

            ppt_app = win32com.client.GetObject(Class="PowerPoint.Application")
            self.logger.log("PowerPoint.Application 연결 성공")

            source_pres = ppt_app.ActivePresentation
            self.logger.log(f"원본 프레젠테이션: {source_pres.Name}")

            # python-pptx로 새 프레젠테이션 생성
            self.logger.log("python-pptx로 새 프레젠테이션 생성...")
            self.root.after(0, lambda: self.status_text.set("새 PPT 생성 중 (python-pptx)..."))
            self.root.after(0, lambda: self.progress_var.set(5))

            new_pres = Presentation()
            self.logger.log("새 Presentation 객체 생성 성공")

            # 슬라이드 크기 설정
            slide_width = source_pres.PageSetup.SlideWidth
            slide_height = source_pres.PageSetup.SlideHeight
            self.logger.log(f"원본 슬라이드 크기: {slide_width} x {slide_height} pt")

            # EMU 단위로 변환 (1pt = 12700 EMU)
            new_pres.slide_width = Emu(int(slide_width * 12700))
            new_pres.slide_height = Emu(int(slide_height * 12700))
            self.logger.log(f"슬라이드 크기 설정 완료")

            total_slides = source_pres.Slides.Count
            self.logger.log(f"총 슬라이드 수: {total_slides}")

            # 임시 디렉토리 생성 (이미지 추출용)
            temp_dir = tempfile.mkdtemp()
            self.logger.log(f"임시 디렉토리 생성: {temp_dir}")

            # 빈 레이아웃 가져오기
            blank_layout = new_pres.slide_layouts[6]  # 빈 레이아웃

            # 각 슬라이드 복사
            for i in range(1, total_slides + 1):
                self.logger.log(f"--- 슬라이드 {i}/{total_slides} 처리 시작 ---")

                progress = 5 + (i / total_slides) * 85
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                self.root.after(0, lambda n=i, t=total_slides: self.status_text.set(f"슬라이드 {n}/{t} 처리 중..."))

                source_slide = source_pres.Slides(i)
                self.logger.log(f"원본 슬라이드 {i} 가져옴")

                # 새 슬라이드 추가
                new_slide = new_pres.slides.add_slide(blank_layout)
                self.logger.log(f"새 슬라이드 {i} 추가 성공")

                # 도형들 복사
                shape_count = source_slide.Shapes.Count
                self.logger.log(f"슬라이드 {i}의 도형 수: {shape_count}")

                for j, shape in enumerate(source_slide.Shapes, 1):
                    try:
                        self.logger.log(f"  도형 {j}/{shape_count} 복사 시도: Type={shape.Type}, Name={shape.Name}")
                        self._copy_shape_to_pptx(shape, new_slide, temp_dir)
                        self.logger.log(f"  도형 {j}/{shape_count} 복사 성공")
                    except Exception as e:
                        self.logger.log(f"  도형 {j}/{shape_count} 복사 실패 (무시): {str(e)}")
                        continue

                self.logger.log(f"--- 슬라이드 {i}/{total_slides} 처리 완료 ---")

            # python-pptx로 직접 저장 (DRM 우회!)
            self.logger.log("=== python-pptx로 직접 저장 (DRM 우회) ===")
            self.root.after(0, lambda: self.status_text.set("파일 저장 중 (DRM 우회)..."))
            self.root.after(0, lambda: self.progress_var.set(95))

            self.logger.log(f"저장 시도: {save_path}")
            new_pres.save(save_path)
            self.logger.log("저장 완료!")

            # 파일 존재 확인
            if os.path.exists(save_path):
                file_size = os.path.getsize(save_path)
                self.logger.log(f"저장된 파일 확인: {save_path} (크기: {file_size} bytes)")
            else:
                self.logger.log(f"[경고] 저장된 파일을 찾을 수 없음: {save_path}")

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("추출 완료! (DRM 우회 성공)"))
            self.root.after(0, lambda: messagebox.showinfo("완료",
                f"PPT 추출 완료! (DRM 우회)\n{save_path}\n\n"
                f"총 {total_slides}장의 슬라이드가 복사되었습니다.\n\n"
                f"로그 파일: {self.logger.log_path}"))

            self.logger.log("=== PPT 추출 프로세스 성공적으로 완료 ===")

        except Exception as e:
            self.logger.error("PPT 추출 중 오류 발생", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {str(e)[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류",
                f"추출 중 오류 발생:\n{str(e)}\n\n"
                f"상세 로그: {self.logger.log_path}"))

        finally:
            # 임시 파일 정리
            if temp_dir and os.path.exists(temp_dir):
                try:
                    self.logger.log(f"임시 디렉토리 정리: {temp_dir}")
                    shutil.rmtree(temp_dir)
                    self.logger.log("임시 디렉토리 정리 완료")
                except Exception as e:
                    self.logger.error("임시 디렉토리 정리 실패", e)

            self.root.after(0, lambda: self.extract_button.config(state=tk.NORMAL))
            pythoncom.CoUninitialize()
            self.logger.log("COM 해제 완료")

    def _copy_shape_to_pptx(self, source_shape, target_slide, temp_dir):
        """COM 도형을 python-pptx 슬라이드로 복사"""
        shape_type = source_shape.Type
        shape_name = source_shape.Name

        self.logger.log(f"    _copy_shape_to_pptx: Type={shape_type}, Name={shape_name}")

        try:
            # 위치 및 크기 (pt -> EMU 변환)
            left = Emu(int(source_shape.Left * 12700))
            top = Emu(int(source_shape.Top * 12700))
            width = Emu(int(source_shape.Width * 12700))
            height = Emu(int(source_shape.Height * 12700))

            self.logger.log(f"    위치: Left={source_shape.Left}, Top={source_shape.Top}, Width={source_shape.Width}, Height={source_shape.Height}")

            # 그림인 경우 - 이미지로 추출 후 삽입
            if shape_type in [13, 11]:  # msoPicture, msoLinkedPicture
                self.logger.log(f"    이미지 타입 감지 (Type={shape_type})")
                if self.copy_images.get():
                    img_path = os.path.join(temp_dir, f"img_{id(source_shape)}.png")
                    self.logger.log(f"    이미지 Export 시도: {img_path}")
                    source_shape.Export(img_path, 2)  # ppShapeFormatPNG = 2
                    self.logger.log(f"    이미지 Export 성공")

                    target_slide.shapes.add_picture(img_path, left, top, width, height)
                    self.logger.log(f"    add_picture 성공")
                return

            # 텍스트 박스인 경우
            if shape_type == 17:  # msoTextBox
                self.logger.log(f"    텍스트박스 타입 감지")
                if self.copy_texts.get():
                    textbox = target_slide.shapes.add_textbox(left, top, width, height)
                    self.logger.log(f"    add_textbox 성공")

                    if source_shape.HasTextFrame:
                        if source_shape.TextFrame.HasText:
                            text = source_shape.TextFrame.TextRange.Text
                            self.logger.log(f"    텍스트: '{text[:30]}...' (길이: {len(text)})")

                            tf = textbox.text_frame
                            tf.word_wrap = True
                            p = tf.paragraphs[0]
                            p.text = text

                            # 폰트 설정 시도
                            try:
                                src_font = source_shape.TextFrame.TextRange.Font
                                p.font.size = Pt(src_font.Size)
                                p.font.bold = src_font.Bold
                                p.font.italic = src_font.Italic

                                # RGB 색상 변환
                                rgb = src_font.Color.RGB
                                r = rgb & 0xFF
                                g = (rgb >> 8) & 0xFF
                                b = (rgb >> 16) & 0xFF
                                p.font.color.rgb = RGBColor(r, g, b)
                                self.logger.log(f"    폰트 설정 성공")
                            except Exception as e:
                                self.logger.log(f"    폰트 설정 실패 (무시): {str(e)}")
                return

            # AutoShape, Freeform 등 - 이미지로 추출 후 삽입
            if shape_type in [1, 5, 14]:  # AutoShape, Freeform, Placeholder
                self.logger.log(f"    도형 타입 감지 (Type={shape_type}) - 이미지로 변환")
                if self.copy_shapes.get():
                    # 텍스트가 있으면 텍스트 박스로 생성
                    has_text = False
                    text = ""
                    try:
                        if source_shape.HasTextFrame:
                            if source_shape.TextFrame.HasText:
                                text = source_shape.TextFrame.TextRange.Text
                                has_text = True
                    except:
                        pass

                    if has_text and text.strip():
                        self.logger.log(f"    텍스트 포함 도형 - 텍스트박스로 생성")
                        textbox = target_slide.shapes.add_textbox(left, top, width, height)
                        tf = textbox.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = text
                        self.logger.log(f"    텍스트 복사 완료: '{text[:30]}...'")
                    else:
                        # 이미지로 추출
                        try:
                            img_path = os.path.join(temp_dir, f"shape_{id(source_shape)}.png")
                            self.logger.log(f"    도형 Export 시도: {img_path}")
                            source_shape.Export(img_path, 2)
                            target_slide.shapes.add_picture(img_path, left, top, width, height)
                            self.logger.log(f"    도형 이미지로 추가 성공")
                        except Exception as e:
                            self.logger.log(f"    도형 이미지 변환 실패 (건너뜀): {str(e)}")
                return

            # 테이블인 경우 - 이미지로 추출
            if shape_type == 19:  # msoTable
                self.logger.log(f"    테이블 타입 감지 - 이미지로 변환")
                if self.copy_shapes.get():
                    try:
                        img_path = os.path.join(temp_dir, f"table_{id(source_shape)}.png")
                        source_shape.Export(img_path, 2)
                        target_slide.shapes.add_picture(img_path, left, top, width, height)
                        self.logger.log(f"    테이블 이미지로 추가 성공")
                    except Exception as e:
                        self.logger.log(f"    테이블 이미지 변환 실패 (건너뜀): {str(e)}")
                return

            # Connector (선) - 이미지로 추출 불가, 건너뛰기
            if shape_type == 9:  # msoConnector
                self.logger.log(f"    Connector(선) 타입 - python-pptx에서 복잡, 건너뜀")
                return

            # 그룹인 경우 - 내부 도형 개별 처리
            if shape_type == 6:  # msoGroup
                self.logger.log(f"    그룹 타입 감지, 내부 도형 수: {source_shape.GroupItems.Count}")
                for item in source_shape.GroupItems:
                    try:
                        self._copy_shape_to_pptx(item, target_slide, temp_dir)
                    except Exception as e:
                        self.logger.log(f"    그룹 내 도형 복사 실패 (무시): {str(e)}")
                return

            # 기타 도형 - 이미지로 추출 시도
            self.logger.log(f"    기타 도형 타입 (Type={shape_type}) - 이미지로 변환 시도")
            if self.copy_shapes.get():
                try:
                    img_path = os.path.join(temp_dir, f"other_{id(source_shape)}.png")
                    source_shape.Export(img_path, 2)
                    target_slide.shapes.add_picture(img_path, left, top, width, height)
                    self.logger.log(f"    기타 도형 이미지로 추가 성공")
                except Exception as e:
                    self.logger.log(f"    기타 도형 이미지 변환 실패 (건너뜀): {str(e)}")

        except Exception as e:
            self.logger.log(f"    도형 복사 실패 (건너뜀): {str(e)}")

    def run(self):
        """프로그램 실행"""
        self.logger.log("메인 루프 시작")
        self.root.mainloop()
        self.logger.log("메인 루프 종료")
        self.logger.close()


def check_dependencies():
    """의존성 확인"""
    errors = []

    if not HAS_WIN32COM:
        errors.append("pywin32 패키지 필요 (pip install pywin32)")

    if not HAS_PPTX:
        errors.append("python-pptx 패키지 필요 (pip install python-pptx)")

    if errors:
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror("의존성 오류",
            "다음 패키지가 필요합니다:\n\n" + "\n".join(errors))
        sys.exit(1)


if __name__ == "__main__":
    check_dependencies()
    app = PPTExtractorV2()
    app.run()
