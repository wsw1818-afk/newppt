"""
문서 추출 도구 v3
- PPT, Excel, Word, 메모장 지원
- COM으로 데이터 읽기
- python-pptx, openpyxl, python-docx로 직접 파일 생성
- 슬라이드 전체 이미지 캡처 + 도형 속성 직접 재생성
- 바탕화면에 상세 로그 저장
"""

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import queue
import tempfile
import shutil
import datetime
import traceback
import io
import time
import zipfile
import base64
import ctypes
from ctypes import wintypes

APP_BUILD_ID = "2026-06-11-hwp-ui"

try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
    HAS_TKINTERDND = True
except ImportError:
    DND_FILES = None
    TkinterDnD = None
    HAS_TKINTERDND = False

# python-pptx 관련
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# openpyxl 관련 (Excel)
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
    from openpyxl.utils import get_column_letter
    from openpyxl.drawing.image import Image as OpenpyxlImage
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# python-docx 관련 (Word)
try:
    from docx import Document as DocxDocument
    from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor, Emu as DocxEmu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PDF 보안 해제 관련 (pypdf)
try:
    from pypdf import PdfReader, PdfWriter
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

# COM 관련
try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False


class Logger:
    """바탕화면에 로그 파일 저장"""

    def __init__(self):
        self._lock = threading.Lock()
        self._line_count = 0
        self._closed = False
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        self.log_path = os.path.join(desktop, f"DocExtractor_Log_{timestamp}.txt")
        self.log_file = open(self.log_path, "w", encoding="utf-8")
        self.log(f"=== 문서 추출기 v3 로그 시작 ===")
        self.log(f"로그 파일: {self.log_path}")
        self.log(f"시작 시간: {datetime.datetime.now()}")
        self.log(f"빌드 ID: {APP_BUILD_ID}")
        self.log(f"실행 파일: {sys.executable}")
        self.log(f"지원 문서: PPT, Excel, Word, 메모장")
        self.log("")

    def log(self, message):
        timestamp = datetime.datetime.now().strftime("%H:%M:%S.%f")[:-3]
        line = f"[{timestamp}] {message}"
        with self._lock:
            if self._closed:
                return
            try:
                if sys.stdout:
                    print(line)
            except Exception:
                pass
            try:
                self.log_file.write(line + "\n")
                self._line_count += 1
                self.log_file.flush()
            except Exception:
                pass

    def error(self, message, exception=None):
        self.log(f"[ERROR] {message}")
        if exception:
            self.log(f"[ERROR] 예외 타입: {type(exception).__name__}")
            self.log(f"[ERROR] 예외 메시지: {str(exception)}")
            self.log(f"[ERROR] 상세 트레이스백:")
            for line in traceback.format_exc().split("\n"):
                self.log(f"        {line}")

    def close(self):
        self.log("")
        self.log(f"=== 로그 종료: {datetime.datetime.now()} ===")
        with self._lock:
            try:
                self.log_file.flush()
                self.log_file.close()
            finally:
                self._closed = True


# AutoShape Type 매핑 (COM -> python-pptx)
# COM AutoShapeType 상수: https://docs.microsoft.com/en-us/office/vba/api/office.msoautoshapetype
AUTOSHAPE_MAPPING = {
    # 기본 도형
    1: MSO_SHAPE.RECTANGLE,           # 사각형
    2: MSO_SHAPE.PARALLELOGRAM,       # 평행사변형
    3: MSO_SHAPE.TRAPEZOID,           # 사다리꼴
    4: MSO_SHAPE.DIAMOND,             # 마름모
    5: MSO_SHAPE.ROUNDED_RECTANGLE,   # 둥근 사각형
    6: MSO_SHAPE.OCTAGON,             # 팔각형
    7: MSO_SHAPE.ISOSCELES_TRIANGLE,  # 이등변 삼각형 (별칭)
    8: MSO_SHAPE.RIGHT_TRIANGLE,      # 직각 삼각형 (별칭)
    9: MSO_SHAPE.ISOSCELES_TRIANGLE,  # 이등변 삼각형
    10: MSO_SHAPE.RIGHT_TRIANGLE,     # 직각 삼각형
    11: MSO_SHAPE.OVAL,               # 타원/원
    # 화살표/기타 기본값
    13: MSO_SHAPE.RIGHT_ARROW,        # 오른쪽 화살표
    14: MSO_SHAPE.LEFT_ARROW,         # 왼쪽 화살표
    15: MSO_SHAPE.UP_ARROW,           # 위쪽 화살표
    16: MSO_SHAPE.DOWN_ARROW,         # 아래쪽 화살표
    17: MSO_SHAPE.LEFT_RIGHT_ARROW,   # 좌우 화살표
    18: MSO_SHAPE.UP_DOWN_ARROW,      # 상하 화살표
    19: MSO_SHAPE.QUAD_ARROW,         # 4방향 화살표
    20: MSO_SHAPE.CHEVRON,            # 갈매기형
    21: MSO_SHAPE.NOTCHED_RIGHT_ARROW,# 홈이 있는 화살표
    22: MSO_SHAPE.PENTAGON,           # 오각형 (집 모양)
    23: MSO_SHAPE.CHEVRON,            # 갈매기형 (별칭)

    # 별
    12: MSO_SHAPE.STAR_5_POINT,       # 5각 별
    37: MSO_SHAPE.STAR_6_POINT,       # 6각 별
    38: MSO_SHAPE.STAR_8_POINT,       # 8각 별
    39: MSO_SHAPE.STAR_16_POINT,      # 16각 별
    40: MSO_SHAPE.STAR_24_POINT,      # 24각 별
    41: MSO_SHAPE.STAR_32_POINT,      # 32각 별

    # 블록 화살표
    24: MSO_SHAPE.RIGHT_ARROW_CALLOUT,    # 설명선 화살표
    25: MSO_SHAPE.LEFT_ARROW_CALLOUT,     # 설명선 화살표
    26: MSO_SHAPE.UP_ARROW_CALLOUT,       # 설명선 화살표
    27: MSO_SHAPE.DOWN_ARROW_CALLOUT,     # 설명선 화살표
    28: MSO_SHAPE.LEFT_RIGHT_ARROW_CALLOUT,
    29: MSO_SHAPE.UP_DOWN_ARROW_CALLOUT,

    # 설명선/말풍선
    30: MSO_SHAPE.ROUNDED_RECTANGLE,  # 둥근 사각형 설명선
    31: MSO_SHAPE.OVAL_CALLOUT,       # 타원 설명선
    32: MSO_SHAPE.CLOUD_CALLOUT,      # 구름 설명선

    # 기호
    33: MSO_SHAPE.HEART,              # 하트
    34: MSO_SHAPE.LIGHTNING_BOLT,     # 번개
    35: MSO_SHAPE.SUN,                # 태양
    36: MSO_SHAPE.MOON,               # 달

    # 도형 (추가)
    42: MSO_SHAPE.FOLDED_CORNER,      # 접힌 모서리
    43: MSO_SHAPE.SMILEY_FACE,        # 스마일
    44: MSO_SHAPE.NO_SYMBOL,          # 금지 표시
    45: MSO_SHAPE.BLOCK_ARC,          # 호
    46: MSO_SHAPE.DONUT,              # 도넛
    47: MSO_SHAPE.RECTANGLE,           # 기울어진 텍스트 (TEXT_SLANT 미지원 → 폴백)
    48: MSO_SHAPE.RECTANGLE,           # 아치형 텍스트 (TEXT_ARCH_DOWN_CURVE 미지원 → 폴백)
}


class DocumentExtractorV3:
    """문서 추출기 v3 - PPT, Excel, Word, 메모장 지원"""

    EXCEL_VALUE_CELL_LIMIT = 500_000
    EXCEL_FORMAT_CELL_LIMIT = 50_000
    EXCEL_ROW_HEIGHT_COPY_LIMIT = 5_000
    EXCEL_COLUMN_WIDTH_COPY_LIMIT = 1_024
    EXCEL_OBJECT_RETRY_COUNT = 2
    EXCEL_OBJECT_RETRY_DELAY = 0.08
    PPT_CLIPBOARD_RETRY_COUNT = 2
    PPT_CLIPBOARD_RETRY_DELAY = 0.08

    def _create_root_window(self):
        if HAS_TKINTERDND:
            try:
                root = TkinterDnD.Tk()
                self.dnd_available = True
                return root
            except Exception as error:
                self.dnd_init_error = str(error)
                self.logger.log(f"드래그앤드롭 루트 초기화 실패: {self.dnd_init_error[:120]}")
        self.dnd_available = False
        return tk.Tk()

    def __init__(self):
        self.logger = Logger()
        self.logger.log("DocumentExtractor v3 초기화 시작")
        self.dnd_available = False
        self.dnd_init_error = ""

        self.root = self._create_root_window()
        self.root.title("문서 추출 도구 v3")
        self.root.geometry("900x660")
        self.root.resizable(False, False)

        # 상태 변수 (공통)
        self.status_text = tk.StringVar(value="프로그램 시작됨")
        self.progress_var = tk.DoubleVar(value=0)

        # PPT 상태 변수
        self.ppt_doc_name = tk.StringVar(value="감지 중...")
        self.ppt_slide_count = tk.StringVar(value="-")
        self.ppt_save_path = tk.StringVar(value="")
        self.ppt_source_path = tk.StringVar(value="")
        self.ppt_list = []
        self.selected_ppt_index = tk.IntVar(value=0)
        self.ppt_input_mode = "open"

        # Excel 상태 변수
        self.excel_doc_name = tk.StringVar(value="감지 중...")
        self.excel_sheet_count = tk.StringVar(value="-")
        self.excel_save_path = tk.StringVar(value="")
        self.excel_source_path = tk.StringVar(value="")
        self.excel_list = []
        self.selected_excel_index = tk.IntVar(value=0)
        self.excel_input_mode = "open"

        # 한글 상태 변수
        self.hwp_doc_name = tk.StringVar(value="감지 중...")
        self.hwp_save_path = tk.StringVar(value="")
        self.hwp_source_path = tk.StringVar(value="")
        self.hwp_list = []
        self.selected_hwp_index = tk.IntVar(value=0)

        # Word 상태 변수
        self.word_doc_name = tk.StringVar(value="감지 중...")
        self.word_page_count = tk.StringVar(value="-")
        self.word_save_path = tk.StringVar(value="")
        self.word_source_path = tk.StringVar(value="")
        self.word_list = []
        self.selected_word_index = tk.IntVar(value=0)
        self.word_input_mode = "open"

        # 메모장 상태 변수
        self.notepad_doc_name = tk.StringVar(value="감지 중...")
        self.notepad_save_path = tk.StringVar(value="")
        self.notepad_source_path = tk.StringVar(value="")
        self.notepad_list = []
        self.notepad_input_mode = "open"

        # PDF 보안 해제 상태
        self.pdf_source_path = tk.StringVar(value="")
        self.pdf_save_path = tk.StringVar(value="")

        # 일괄 변환 상태
        self.batch_files = []
        self.batch_output_dir = tk.StringVar(value="")
        self.batch_status_text = tk.StringVar(value="파일을 추가하고 출력 폴더를 선택하세요.")

        # 일괄 변환 Office 인스턴스 예열·재사용 (콜드 스타트 비용을 최초 1회로 축소)
        self._office_job_queue = queue.Queue()
        self._office_worker = None
        self._warm_office_apps = {}  # key("ppt"/"excel"/"word") -> COM app, 워커 스레드 소유

        # 탭 변경 추적 (중복 감지 방지)
        self.last_tab_index = -1
        self.tab_detected = [False, False, False, False, False, False, False]  # PPT, Excel, Word, 한글, 메모장, PDF, 일괄 변환
        self.current_doc_index = 0
        self.nav_buttons = []
        self._hwp_detecting = False

        self.setup_ui()
        if self.dnd_available:
            self._setup_drag_drop()
        else:
            reason = self.dnd_init_error or "tkinterdnd2 패키지가 없습니다."
            self.logger.log(f"드래그앤드롭 비활성화: {reason}")

        self.logger.log("DocumentExtractor v3 초기화 완료")

    def _configure_styles(self):
        """업무용 도구 화면에 맞춘 공통 ttk 스타일."""
        style = ttk.Style()
        try:
            style.theme_use("clam")
        except Exception:
            pass

        self.ui_colors = {
            "app_bg": "#edf1f6",
            "panel_bg": "#ffffff",
            "nav_bg": "#ffffff",
            "nav_border": "#d8e0ea",
            "nav_fg": "#2f3b4a",
            "nav_muted": "#667085",
            "nav_selected_bg": "#eaf2ff",
            "nav_selected_fg": "#1f5fbf",
            "text": "#1f2937",
            "muted": "#5d6b7c",
            "border": "#d9e0e8",
            "section_bg": "#ffffff",
            "accent": "#2f6fed",
            "accent_active": "#245fd1",
            "field_bg": "#ffffff",
        }
        c = self.ui_colors

        style.configure("App.TFrame", background=c["app_bg"])
        style.configure("TFrame", background=c["section_bg"])
        style.configure("Panel.TFrame", background="#ffffff")
        style.configure("Card.TFrame", background=c["section_bg"])
        style.configure("Footer.TFrame", background=c["app_bg"])
        style.configure("TLabel", background=c["section_bg"], foreground=c["text"], font=("맑은 고딕", 9))
        style.configure("Title.TLabel", background=c["app_bg"], foreground=c["text"], font=("맑은 고딕", 15, "bold"))
        style.configure("Subtitle.TLabel", background=c["app_bg"], foreground=c["muted"], font=("맑은 고딕", 9))
        style.configure("PanelTitle.TLabel", background=c["panel_bg"], foreground=c["text"], font=("맑은 고딕", 14, "bold"))
        style.configure("PanelSubtitle.TLabel", background=c["panel_bg"], foreground=c["muted"], font=("맑은 고딕", 9))
        style.configure("Section.TLabelframe", background=c["section_bg"], bordercolor=c["border"], relief=tk.SOLID)
        style.configure("Section.TLabelframe.Label", background=c["section_bg"], foreground=c["text"], font=("맑은 고딕", 9, "bold"))
        style.configure("TRadiobutton", background=c["section_bg"], foreground=c["text"], font=("맑은 고딕", 9))
        style.configure("TCheckbutton", background=c["section_bg"], foreground=c["text"], font=("맑은 고딕", 9))
        style.configure("TEntry", fieldbackground=c["field_bg"], bordercolor=c["border"], lightcolor=c["border"], darkcolor=c["border"])
        style.configure("TCombobox", fieldbackground=c["field_bg"], background=c["field_bg"], bordercolor=c["border"], arrowcolor=c["muted"])
        style.map(
            "TCombobox",
            fieldbackground=[("readonly", c["field_bg"])],
            background=[("readonly", c["field_bg"])],
            foreground=[("readonly", c["text"])],
        )
        style.configure("TButton", padding=(12, 5), font=("맑은 고딕", 9))
        style.map("TButton", background=[("active", "#e8edf5")])
        style.configure("Secondary.TButton", background="#f6f8fb", foreground=c["text"], bordercolor=c["border"], padding=(12, 5))
        style.configure("Accent.TButton", background=c["accent"], foreground="#ffffff", bordercolor=c["accent"], padding=(18, 9), font=("맑은 고딕", 10, "bold"))
        style.map("Accent.TButton", background=[("active", c["accent_active"])], foreground=[("active", "#ffffff")])
        style.configure("Horizontal.TProgressbar", background=c["accent"], troughcolor="#dfe5ee", bordercolor="#dfe5ee")

    def setup_ui(self):
        """UI 구성"""
        self.logger.log("UI 구성 시작")
        self._configure_styles()
        self.root.configure(bg=self.ui_colors["app_bg"])

        main_frame = ttk.Frame(self.root, padding=0, style="App.TFrame")
        main_frame.pack(fill=tk.BOTH, expand=True)

        header_frame = ttk.Frame(main_frame, padding=(22, 16, 22, 8), style="App.TFrame")
        header_frame.pack(fill=tk.X)

        ttk.Label(header_frame, text="문서 추출 도구 v3", style="Title.TLabel").pack(anchor=tk.W)
        ttk.Label(
            header_frame,
            text="PPT, Excel, Word, 메모장, PDF를 파일 또는 열린 문서에서 가져와 새 파일로 내보냅니다.",
            style="Subtitle.TLabel",
        ).pack(anchor=tk.W, pady=(2, 0))

        body_frame = ttk.Frame(main_frame, padding=(14, 2, 14, 10), style="App.TFrame")
        body_frame.pack(fill=tk.BOTH, expand=True)

        c = self.ui_colors
        sidebar = tk.Frame(
            body_frame,
            width=188,
            bg=c["nav_bg"],
            bd=0,
            highlightthickness=1,
            highlightbackground=c["nav_border"],
            highlightcolor=c["nav_border"],
        )
        sidebar.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 12))
        sidebar.pack_propagate(False)

        tk.Label(
            sidebar,
            text="문서 종류",
            bg=c["nav_bg"],
            fg=c["nav_muted"],
            font=("맑은 고딕", 9, "bold"),
            anchor=tk.W,
        ).pack(fill=tk.X, padx=16, pady=(16, 8))

        self.doc_views = [
            ("PowerPoint", "PPT", "슬라이드/도형 보존", self.detect_open_ppt),
            ("Excel", "XLS", "시트/도형 보존", self.detect_open_excel),
            ("Word", "DOC", "문서 구조 보존", self.detect_open_word),
            ("한글", "HWP", "문서 구조 보존", None),
            ("메모장", "TXT", "텍스트 추출", self.detect_open_notepad),
            ("PDF", "PDF", "보안 해제", None),
            ("일괄 변환", "ALL", "파일 묶음 처리", None),
        ]
        self.view_title_text = tk.StringVar(value=self.doc_views[0][0])
        self.view_summary_text = tk.StringVar(value=self.doc_views[0][2])

        for index, (title, badge, summary, _detect_fn) in enumerate(self.doc_views):
            item = tk.Frame(sidebar, bg=c["nav_bg"], bd=0, cursor="hand2")
            item.pack(fill=tk.X, padx=10, pady=3)
            badge_label = tk.Label(item, text=badge, width=4, bg="#f2f5f9", fg=c["nav_muted"],
                                   font=("맑은 고딕", 8, "bold"))
            badge_label.pack(side=tk.LEFT, padx=(8, 7), pady=9)
            text_box = tk.Frame(item, bg=c["nav_bg"])
            text_box.pack(side=tk.LEFT, fill=tk.X, expand=True, pady=7)
            title_label = tk.Label(text_box, text=title, bg=c["nav_bg"], fg=c["nav_fg"],
                                   font=("맑은 고딕", 9, "bold"), anchor=tk.W)
            title_label.pack(fill=tk.X)
            summary_label = tk.Label(text_box, text=summary, bg=c["nav_bg"], fg=c["nav_muted"],
                                     font=("맑은 고딕", 8), anchor=tk.W)
            summary_label.pack(fill=tk.X)

            for widget in (item, badge_label, text_box, title_label, summary_label):
                widget.bind("<Button-1>", lambda _event, i=index: self._select_document_view(i))
            self.nav_buttons.append((item, badge_label, text_box, title_label, summary_label))

        content_shell = tk.Frame(
            body_frame,
            bg=c["panel_bg"],
            bd=0,
            highlightthickness=1,
            highlightbackground=c["border"],
            highlightcolor=c["border"],
            padx=16,
            pady=14,
        )
        content_shell.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        title_row = ttk.Frame(content_shell, style="Panel.TFrame")
        title_row.pack(fill=tk.X, pady=(0, 10))
        ttk.Label(
            title_row,
            textvariable=self.view_title_text,
            style="PanelTitle.TLabel",
        ).pack(side=tk.LEFT)
        ttk.Label(
            title_row,
            textvariable=self.view_summary_text,
            style="PanelSubtitle.TLabel",
        ).pack(side=tk.LEFT, padx=(10, 0), pady=(3, 0))

        self.content_area = ttk.Frame(content_shell, style="Panel.TFrame")
        self.content_area.pack(fill=tk.BOTH, expand=True)
        self.content_area.grid_rowconfigure(0, weight=1)
        self.content_area.grid_columnconfigure(0, weight=1)

        # PPT 탭
        self.ppt_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.ppt_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_ppt_tab()

        # Excel 탭
        self.excel_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.excel_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_excel_tab()

        # Word 탭
        self.word_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.word_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_word_tab()

        # 한글 탭
        self.hwp_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.hwp_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_hwp_tab()

        # 메모장 탭
        self.notepad_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.notepad_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_notepad_tab()

        # 일괄 변환 탭
        self.batch_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.batch_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_batch_tab()

        # PDF 보안 해제 탭
        self.pdf_tab = ttk.Frame(self.content_area, style="Panel.TFrame")
        self.pdf_tab.grid(row=0, column=0, sticky="nsew")
        self._setup_pdf_tab()

        self.content_frames = [self.ppt_tab, self.excel_tab, self.word_tab, self.hwp_tab, self.notepad_tab, self.pdf_tab, self.batch_tab]
        self._select_document_view(0, detect=False)

        footer_frame = ttk.Frame(main_frame, padding=(12, 0, 12, 12), style="Footer.TFrame")
        footer_frame.pack(fill=tk.X)

        # 진행바 (공통)
        self.progress = ttk.Progressbar(footer_frame, variable=self.progress_var,
                                         maximum=100, length=550)
        self.progress.pack(fill=tk.X, pady=(0, 6))

        # 상태 표시 (공통)
        status_frame = ttk.Frame(footer_frame, style="Footer.TFrame")
        status_frame.pack(fill=tk.X)
        ttk.Label(status_frame, text="상태:").pack(side=tk.LEFT)
        ttk.Label(status_frame, textvariable=self.status_text,
                  font=("맑은 고딕", 9)).pack(side=tk.LEFT, padx=(5, 0))

        self.logger.log("UI 구성 완료")
        self.root.after(120, self._schedule_detect)

    def _create_section(self, parent, title):
        """기본 LabelFrame 대신 쓰는 밝은 카드형 섹션."""
        c = self.ui_colors
        outer = tk.Frame(
            parent,
            bg=c["section_bg"],
            bd=0,
            highlightthickness=1,
            highlightbackground=c["border"],
            highlightcolor=c["border"],
        )
        outer.pack(fill=tk.X, pady=(0, 10), padx=5)

        tk.Label(
            outer,
            text=title,
            bg=c["section_bg"],
            fg=c["text"],
            font=("맑은 고딕", 9, "bold"),
            anchor=tk.W,
        ).pack(fill=tk.X, padx=12, pady=(10, 4))

        body = ttk.Frame(outer, style="Card.TFrame")
        body.pack(fill=tk.X, padx=12, pady=(0, 12))
        return body

    def _connect_com_app(self, prog_id, display_name, allow_dispatch=True, use_get_active=True):
        """Office/HWP COM 애플리케이션 연결을 공통 처리한다."""
        try:
            app = win32com.client.GetObject(Class=prog_id)
            self.logger.log(f"{display_name} GetObject 연결 성공")
            return app, False
        except Exception as e1:
            self.logger.log(f"{display_name} GetObject 실패: {str(e1)[:50]}")

        if use_get_active:
            try:
                app = win32com.client.GetActiveObject(prog_id)
                self.logger.log(f"{display_name} GetActiveObject 연결 성공")
                return app, False
            except Exception as e2:
                self.logger.log(f"{display_name} GetActiveObject 실패: {str(e2)[:50]}")

        if allow_dispatch:
            try:
                app = win32com.client.Dispatch(prog_id)
                self.logger.log(f"{display_name} Dispatch 연결 성공")
                return app, True
            except Exception as e3:
                self.logger.log(f"{display_name} Dispatch 실패: {str(e3)[:50]}")

        raise Exception(f"{display_name}에 연결할 수 없습니다. {display_name}를 먼저 실행해주세요.")

    def _create_isolated_com_app(self, prog_id, display_name):
        """파일 직접/일괄 변환용 Office 인스턴스를 사용자 작업과 분리해서 만든다."""
        try:
            app = win32com.client.DispatchEx(prog_id)
            self.logger.log(f"{display_name} DispatchEx 격리 인스턴스 생성 성공")
            if prog_id == "PowerPoint.Application":
                try:
                    app.Visible = True
                    self.logger.log("PowerPoint 격리 인스턴스 표시 상태 설정")
                except Exception as visible_error:
                    self.logger.log(f"PowerPoint 표시 상태 설정 실패: {str(visible_error)[:60]}")
            return app, True
        except Exception as exc:
            self.logger.log(f"{display_name} DispatchEx 격리 인스턴스 생성 실패: {str(exc)[:80]}")
            raise Exception(
                f"{display_name} 변환용 격리 인스턴스를 만들 수 없습니다. "
                f"{display_name} 설치/보안 정책을 확인해 주세요."
            ) from exc

    def _set_office_display_alerts(self, app, value, label):
        """자동 변환 중 바꾼 DisplayAlerts 값을 나중에 되돌릴 수 있게 보관한다."""
        try:
            original = app.DisplayAlerts
            app.DisplayAlerts = value
            self.logger.log(f"{label} DisplayAlerts 변경: {original} -> {value}")
            return original
        except Exception as exc:
            self.logger.log(f"{label} DisplayAlerts 변경 실패: {str(exc)[:60]}")
            return None

    def _restore_office_display_alerts(self, app, original, label):
        if app is None or original is None:
            return
        try:
            app.DisplayAlerts = original
            self.logger.log(f"{label} DisplayAlerts 원복: {original}")
        except Exception as exc:
            self.logger.log(f"{label} DisplayAlerts 원복 실패: {str(exc)[:60]}")

    def _is_expected_app_not_running(self, exc, display_name):
        """사용자가 아직 앱을 열지 않은 정상 감지 실패인지 판단한다."""
        message = str(exc)
        return (
            f"{display_name}에 연결할 수 없습니다" in message
            and "먼저 실행" in message
        )

    def _get_ppt_app(self, allow_dispatch=True):
        return self._connect_com_app("PowerPoint.Application", "PowerPoint", allow_dispatch=allow_dispatch)

    def _get_excel_app(self, allow_dispatch=True):
        return self._connect_com_app("Excel.Application", "Excel", allow_dispatch=allow_dispatch)

    def _get_hwp_app(self, allow_dispatch=True):
        hwp, created = self._connect_com_app(
            "HWPFrame.HwpObject",
            "한글",
            allow_dispatch=allow_dispatch,
            use_get_active=False,
        )
        self._hwp_suppress_security_prompt(hwp)
        return hwp, created

    def _hwp_suppress_security_prompt(self, hwp):
        """한글 자동화 파일 접근 보안 팝업을 자동 수락 처리한다(변환 전 권한 묻는 대화상자 억제).

        FilePathCheckerModule을 등록하면 스크립트의 파일 접근에 매번 묻던 보안 승인 창이
        뜨지 않는다. 모듈 DLL이 없는 환경에서는 실패하지만 변환 자체에는 영향 없다.
        """
        try:
            hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
            self.logger.log("한글 보안 모듈 등록: 파일 접근 권한 팝업 자동 수락")
        except Exception as reg_error:
            self.logger.log(f"FilePathCheckerModule 등록 실패(무시): {str(reg_error)[:80]}")

    def _get_hwp_app_for_extraction(self):
        """추출용 HWP 연결. 새 빈 문서가 만들어지는 연결은 차단한다."""
        before_titles = self._list_hwp_window_titles()
        if not before_titles:
            raise Exception("열린 한글 문서가 없습니다. 한글에서 문서를 먼저 열어주세요.")

        try:
            return self._get_hwp_app(allow_dispatch=False)
        except Exception as active_error:
            self.logger.log(f"한글 활성 COM 연결 실패, Dispatch 연결 검증 시도: {str(active_error)[:80]}")

        hwp, created = self._get_hwp_app(allow_dispatch=True)
        after_titles = self._list_hwp_window_titles()
        if created and len(after_titles) > len(before_titles):
            try:
                hwp.Quit()
            except Exception:
                pass
            raise Exception(
                "한글 COM 연결 중 새 빈 문서가 생성되어 중단했습니다.\n"
                "기존 한글 문서를 닫았다가 다시 열거나, 프로그램과 한글을 같은 권한으로 실행해 주세요."
            )

        return hwp, created

    def _is_hwp_window_title(self, title):
        normalized = (title or "").strip()
        return (
            normalized.endswith(" - 한글")
            or normalized.endswith("- 한글")
            or " - 한글 " in normalized
        )

    def _get_visible_windows(self):
        """현재 보이는 최상위 창의 핸들/제목 목록을 가져온다."""
        windows = []
        user32 = ctypes.windll.user32
        enum_proc_type = ctypes.WINFUNCTYPE(wintypes.BOOL, wintypes.HWND, wintypes.LPARAM)

        def enum_proc(hwnd, lparam):
            try:
                if not user32.IsWindowVisible(hwnd):
                    return True
                length = user32.GetWindowTextLengthW(hwnd)
                if length <= 0:
                    return True
                buffer = ctypes.create_unicode_buffer(length + 1)
                user32.GetWindowTextW(hwnd, buffer, length + 1)
                title = buffer.value.strip()
                if title:
                    windows.append((int(hwnd), title))
            except Exception:
                pass
            return True

        user32.EnumWindows(enum_proc_type(enum_proc), 0)
        return windows

    def _get_visible_window_titles(self):
        """현재 보이는 최상위 창 제목 목록을 가져온다."""
        return [title for _, title in self._get_visible_windows()]

    def _list_hwp_windows(self):
        """한글 COM이 ROT에 없을 때 새 인스턴스를 만들지 않고 창 핸들만 감지한다."""
        return [
            (hwnd, title)
            for hwnd, title in self._get_visible_windows()
            if self._is_hwp_window_title(title)
        ]

    def _list_hwp_window_titles(self):
        """한글 COM이 ROT에 없을 때 새 인스턴스를 만들지 않고 창 제목만 감지한다."""
        return [title for _, title in self._list_hwp_windows()]

    def _get_window_class_name(self, hwnd):
        try:
            buffer = ctypes.create_unicode_buffer(256)
            ctypes.windll.user32.GetClassNameW(hwnd, buffer, 256)
            return buffer.value
        except Exception:
            return ""

    def _get_window_title(self, hwnd):
        try:
            user32 = ctypes.windll.user32
            title_length = user32.GetWindowTextLengthW(hwnd)
            title = ctypes.create_unicode_buffer(title_length + 1)
            user32.GetWindowTextW(hwnd, title, title_length + 1)
            return title.value
        except Exception:
            return ""

    def _is_notepad_window(self, hwnd):
        class_name = self._get_window_class_name(hwnd)
        if class_name == "Notepad":
            return True

        title = self._get_window_title(hwnd).strip().lower()
        return title.endswith(" - notepad") or title.endswith(" - 메모장")

    def _unpack_hwp_item(self, item):
        if not item:
            return "", "", 1, 0
        name = item[0] if len(item) > 0 else ""
        path = item[1] if len(item) > 1 else ""
        hwp_index = item[2] if len(item) > 2 else 1
        hwnd = item[3] if len(item) > 3 else 0
        return name, path, hwp_index, hwnd

    def _find_hwp_window_for_item(self, item):
        name, _path, _hwp_index, hwnd = self._unpack_hwp_item(item)
        user32 = ctypes.windll.user32
        if hwnd and user32.IsWindow(hwnd):
            return hwnd

        for candidate_hwnd, title in self._list_hwp_windows():
            if not name or name == title or name in title or title in name:
                return candidate_hwnd
        return 0

    def _set_clipboard_text(self, text):
        try:
            import win32clipboard
        except Exception as import_error:
            raise Exception(f"클립보드 접근에 필요한 pywin32 모듈을 불러오지 못했습니다: {import_error}")

        win32clipboard.OpenClipboard()
        try:
            win32clipboard.EmptyClipboard()
            win32clipboard.SetClipboardText(text, win32clipboard.CF_UNICODETEXT)
        finally:
            win32clipboard.CloseClipboard()

    def _send_vk(self, vk_code, delay=0.08):
        user32 = ctypes.windll.user32
        KEYEVENTF_KEYUP = 0x0002
        user32.keybd_event(vk_code, 0, 0, 0)
        time.sleep(delay)
        user32.keybd_event(vk_code, 0, KEYEVENTF_KEYUP, 0)
        time.sleep(delay)

    def _send_hotkey(self, *vk_codes):
        user32 = ctypes.windll.user32
        KEYEVENTF_KEYUP = 0x0002
        for vk_code in vk_codes:
            user32.keybd_event(vk_code, 0, 0, 0)
            time.sleep(0.03)
        for vk_code in reversed(vk_codes):
            user32.keybd_event(vk_code, 0, KEYEVENTF_KEYUP, 0)
            time.sleep(0.03)
        time.sleep(0.08)

    def _activate_window(self, hwnd):
        user32 = ctypes.windll.user32
        SW_RESTORE = 9
        user32.ShowWindow(hwnd, SW_RESTORE)
        user32.SetForegroundWindow(hwnd)
        try:
            title = self._get_window_title(hwnd)
            if title and HAS_WIN32COM:
                shell = win32com.client.Dispatch("WScript.Shell")
                shell.AppActivate(title)
        except Exception:
            pass
        time.sleep(0.4)

    def _is_window_enabled_visible(self, hwnd):
        user32 = ctypes.windll.user32
        try:
            return bool(user32.IsWindow(hwnd) and user32.IsWindowVisible(hwnd) and user32.IsWindowEnabled(hwnd))
        except Exception:
            return False

    def _get_window_rect_tuple(self, hwnd):
        rect = wintypes.RECT()
        try:
            ctypes.windll.user32.GetWindowRect(hwnd, ctypes.byref(rect))
            return (rect.left, rect.top, rect.right, rect.bottom)
        except Exception:
            return (0, 0, 0, 0)

    def _enum_descendant_windows(self, parent_hwnd):
        user32 = ctypes.windll.user32
        found = []
        queue = [parent_hwnd]
        seen = set()
        while queue:
            current = queue.pop(0)
            if current in seen:
                continue
            seen.add(current)

            child = 0
            while True:
                child = user32.FindWindowExW(current, child, None, None)
                if not child:
                    break
                found.append(int(child))
                queue.append(int(child))
        return found

    def _log_window_tree(self, hwnd, label, limit=80):
        try:
            title = self._get_window_title(hwnd)
            class_name = self._get_window_class_name(hwnd)
            self.logger.log(f"{label}: hwnd={hwnd}, class='{class_name}', title='{title}'")
            for idx, child in enumerate(self._enum_descendant_windows(hwnd)[:limit], start=1):
                child_title = self._get_window_title(child)
                child_class = self._get_window_class_name(child)
                rect = self._get_window_rect_tuple(child)
                self.logger.log(
                    f"  child {idx}: hwnd={child}, class='{child_class}', "
                    f"title='{child_title}', rect={rect}"
                )
        except Exception as log_error:
            self.logger.log(f"{label} 창 구조 로그 실패: {str(log_error)[:80]}")

    def _wait_for_save_dialog(self, owner_hwnd, timeout=2.5):
        user32 = ctypes.windll.user32
        deadline = time.time() + timeout
        while time.time() < deadline:
            candidates = [user32.GetForegroundWindow()]
            candidates.extend(hwnd for hwnd, _title in self._get_visible_windows())
            for hwnd in candidates:
                if not hwnd or hwnd == owner_hwnd:
                    continue
                title = self._get_window_title(hwnd)
                class_name = self._get_window_class_name(hwnd)
                title_lower = title.lower()
                if (
                    class_name == "#32770"
                    or "저장" in title
                    or "다른 이름" in title
                    or "save" in title_lower
                ):
                    return hwnd
            time.sleep(0.1)
        return 0

    def _find_save_dialog_edit(self, dialog_hwnd):
        user32 = ctypes.windll.user32
        candidates = []
        for child in self._enum_descendant_windows(dialog_hwnd):
            if not user32.IsWindow(child):
                continue
            class_name = self._get_window_class_name(child)
            if class_name not in ("Edit", "RichEdit20W", "RichEdit50W", "RICHEDIT50W"):
                continue
            left, top, right, bottom = self._get_window_rect_tuple(child)
            width = max(0, right - left)
            height = max(0, bottom - top)
            if width < 80 or height < 10:
                continue
            visible = bool(user32.IsWindowVisible(child))
            enabled = bool(user32.IsWindowEnabled(child))
            candidates.append(((visible, enabled, top, left, width), child, class_name, (left, top, right, bottom)))

        if not candidates:
            self.logger.log("한글 UI 저장: 파일명 입력칸 후보 없음")
            return 0
        candidates.sort(key=lambda item: item[0], reverse=True)
        summary = ", ".join(
            f"hwnd={child}/class={class_name}/rect={rect}"
            for _score, child, class_name, rect in candidates[:3]
        )
        self.logger.log(f"한글 UI 저장: 파일명 입력칸 후보 {summary}")
        return candidates[0][1]

    def _find_save_dialog_button(self, dialog_hwnd):
        button_candidates = []
        for child in self._enum_descendant_windows(dialog_hwnd):
            if not self._is_window_enabled_visible(child):
                continue
            class_name = self._get_window_class_name(child)
            if class_name != "Button":
                continue
            title = self._get_window_title(child).strip()
            normalized = title.replace("&", "").lower()
            if (
                "저장" in title
                or "save" in normalized
                or "확인" in title
                or normalized in ("ok", "yes")
            ):
                button_candidates.append(child)

        return button_candidates[0] if button_candidates else 0

    def _find_dialog_confirmation_button(self, dialog_hwnd):
        for child in self._enum_descendant_windows(dialog_hwnd):
            if not self._is_window_enabled_visible(child):
                continue
            if self._get_window_class_name(child) != "Button":
                continue
            title = self._get_window_title(child).strip()
            normalized = title.replace("&", "").lower()
            if "확인" in title or normalized in ("ok", "yes"):
                return child
        return 0

    def _submit_save_dialog_by_controls(self, dialog_hwnd, save_path):
        user32 = ctypes.windll.user32
        edit_hwnd = self._find_save_dialog_edit(dialog_hwnd)
        if not edit_hwnd:
            return False

        WM_SETTEXT = 0x000C
        BM_CLICK = 0x00F5
        self.logger.log(f"한글 UI 저장: 파일명 입력칸 감지 hwnd={edit_hwnd}")
        user32.SetForegroundWindow(dialog_hwnd)
        time.sleep(0.2)
        dialog_save_path = os.path.normpath(save_path)
        user32.SendMessageW(edit_hwnd, WM_SETTEXT, 0, dialog_save_path)
        time.sleep(0.2)
        entered = self._get_window_title(edit_hwnd)
        if entered:
            self.logger.log(f"한글 UI 저장: 파일명 입력값 확인='{entered}'")
        else:
            self.logger.log("한글 UI 저장: 파일명 입력값은 대화상자 제한으로 읽지 못했지만 저장을 계속 진행합니다")

        button_hwnd = self._find_save_dialog_button(dialog_hwnd)
        if button_hwnd:
            self.logger.log(f"한글 UI 저장: 저장 버튼 감지 hwnd={button_hwnd}")
            user32.SendMessageW(button_hwnd, BM_CLICK, 0, 0)
        else:
            self.logger.log("한글 UI 저장: 저장 버튼 감지 실패, Enter 전송")
            self._send_vk(0x0D)
        return True

    def _submit_save_dialog(self, dialog_hwnd, save_path):
        user32 = ctypes.windll.user32
        user32.SetForegroundWindow(dialog_hwnd)
        time.sleep(0.2)
        self._log_window_tree(dialog_hwnd, "한글 UI 저장 대화상자")
        if self._submit_save_dialog_by_controls(dialog_hwnd, save_path):
            return

        self.logger.log("한글 UI 저장: 컨트롤 직접 입력 실패, 키보드 입력 폴백")
        self._set_clipboard_text(os.path.normpath(save_path))
        self._send_hotkey(0x11, 0x41)  # Ctrl+A
        self._send_hotkey(0x11, 0x56)  # Ctrl+V
        self._send_vk(0x0D)  # Enter

    def _try_hwp_save_shortcut(self, hwnd, save_path, shortcut_name, shortcut_fn):
        self._activate_window(hwnd)
        shortcut_fn()
        dialog_hwnd = self._wait_for_save_dialog(hwnd)
        if not dialog_hwnd:
            self.logger.log(f"한글 UI 저장 단축키 {shortcut_name}: 저장 대화상자 감지 실패")
            return False

        self.logger.log(f"한글 UI 저장 단축키 {shortcut_name}: 저장 대화상자 감지 hwnd={dialog_hwnd}")
        self._submit_save_dialog(dialog_hwnd, save_path)
        return True

    def _confirm_hwp_save_dialogs(self, owner_hwnd):
        """저장 중 확인/경고 대화상자가 뜨면 기본 확인 버튼을 눌러 진행한다."""
        user32 = ctypes.windll.user32
        for hwnd, _title in self._get_visible_windows():
            if hwnd == owner_hwnd:
                continue
            class_name = self._get_window_class_name(hwnd)
            title = self._get_window_title(hwnd)
            if class_name != "#32770":
                continue
            if not any(token in title.lower() for token in ("한글", "hwp", "저장", "확인", "경고", "알림", "save")):
                continue
            button_hwnd = self._find_dialog_confirmation_button(hwnd)
            if button_hwnd:
                self.logger.log(f"한글 UI 저장 확인 대화상자 처리: title='{title}', button={button_hwnd}")
                user32.SendMessageW(button_hwnd, 0x00F5, 0, 0)  # BM_CLICK

    def _save_hwp_via_window(self, hwnd, save_path, save_format):
        """COM 연결이 막힌 한글 창을 UI 단축키로 다른 이름 저장한다."""
        user32 = ctypes.windll.user32
        if not hwnd or not user32.IsWindow(hwnd):
            raise Exception("한글 창 핸들을 찾을 수 없어 UI 저장을 진행할 수 없습니다.")

        target_ext = ".hwpx" if save_format == "hwpx" else ".hwp"
        if os.path.splitext(save_path)[1].lower() != target_ext:
            raise Exception(f"한글 저장 형식과 파일 확장자가 다릅니다. {target_ext} 파일로 저장해 주세요.")

        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        backup_path = None
        if os.path.exists(save_path):
            backup_path = self._add_filename_suffix(
                save_path,
                f".docextract_backup_{datetime.datetime.now().strftime('%Y%m%d%H%M%S')}"
            )
            os.replace(save_path, backup_path)

        dialog_save_path = os.path.normpath(save_path)
        self.logger.log(f"한글 COM 대체 UI 저장 시도: hwnd={hwnd}, path={save_path}, dialog_path={dialog_save_path}")

        try:
            shortcuts = [
                ("Alt+V", lambda: self._send_hotkey(0x12, 0x56)),
                ("Ctrl+Shift+S", lambda: self._send_hotkey(0x11, 0x10, 0x53)),
            ]
            for shortcut_name, shortcut_fn in shortcuts:
                if not self._try_hwp_save_shortcut(hwnd, dialog_save_path, shortcut_name, shortcut_fn):
                    continue

                deadline = time.time() + 20
                while time.time() < deadline:
                    if os.path.exists(save_path) and os.path.getsize(save_path) > 0:
                        header = self._read_header_hex(save_path, 8)
                        if header.startswith("53 43 44 53"):
                            self.logger.log("한글 UI 저장 결과가 회사 보안/DRM 컨테이너(SCDS)라서 실패 처리합니다")
                            try:
                                os.remove(save_path)
                            except Exception:
                                pass
                            raise Exception(self._hwp_drm_container_message(save_path))
                        if backup_path and os.path.exists(backup_path):
                            os.remove(backup_path)
                        return
                    self._confirm_hwp_save_dialogs(hwnd)
                    time.sleep(0.2)
                self.logger.log(f"한글 UI 저장 단축키 {shortcut_name}: 파일 생성 대기 시간 초과")
                self._send_vk(0x1B)  # Esc
                time.sleep(0.3)

            self._send_vk(0x1B)  # Esc
            raise Exception(
                "한글 UI 저장으로 결과 파일이 생성되지 않았습니다.\n"
                "한글 창이 앞에 떠 있는지, 저장 대화상자가 보안 프로그램에 의해 막히지 않았는지 확인해 주세요."
            )
        except Exception:
            if backup_path and os.path.exists(backup_path) and not os.path.exists(save_path):
                os.replace(backup_path, save_path)
            raise

    def _find_child_window_by_classes(self, parent_hwnd, class_names):
        """직계 자식뿐 아니라 중첩된 Win32 텍스트 컨트롤까지 찾는다."""
        user32 = ctypes.windll.user32
        wanted = set(class_names)
        queue = [parent_hwnd]
        seen = set()

        while queue:
            current = queue.pop(0)
            if current in seen:
                continue
            seen.add(current)

            child = 0
            while True:
                child = user32.FindWindowExW(current, child, None, None)
                if not child:
                    break
                class_name = self._get_window_class_name(child)
                if class_name in wanted:
                    return child
                queue.append(child)

        return 0

    def _clean_xml_text(self, text):
        """python-docx가 저장할 수 없는 Word/Win32 제어문자를 제거한다."""
        if not text:
            return ""

        cleaned = []
        removed = 0
        for char in text:
            code = ord(char)
            if char in ("\t", "\n", "\r") or (
                0x20 <= code <= 0xD7FF
                or 0xE000 <= code <= 0xFFFD
                or 0x10000 <= code <= 0x10FFFF
            ):
                cleaned.append(char)
            else:
                removed += 1

        if removed:
            self.logger.log(f"  XML 비호환 제어문자 제거: {removed}개")
        return "".join(cleaned)

    def _log_elapsed(self, label, start_time):
        elapsed = time.perf_counter() - start_time
        self.logger.log(f"{label}: {elapsed:.2f}초")
        return elapsed

    def _read_header_hex(self, path, size=16):
        try:
            with open(path, "rb") as f:
                return " ".join(f"{b:02X}" for b in f.read(size))
        except Exception:
            return "읽기 실패"

    def _hwp_drm_container_message(self, save_path):
        return (
            "한글 저장 결과가 회사 보안/DRM 컨테이너(SCDS)입니다.\n"
            "이 파일은 일반 HWP로 변환된 것이 아니어서 변환 실패로 처리했습니다.\n\n"
            "해결하려면 보안 프로그램의 공식 반출/해제 권한으로 원본을 먼저 일반 HWP/HWPX로 저장해야 합니다.\n"
            f"대상 경로: {save_path}"
        )

    def _add_filename_suffix(self, path, suffix):
        directory = os.path.dirname(path)
        stem, ext = os.path.splitext(os.path.basename(path))
        return os.path.join(directory, f"{stem}{suffix}{ext}")

    def _make_local_temp_path(self, suffix):
        """OneDrive/DRM 감시 폴더를 피해서 로컬 임시 파일 경로를 만든다."""
        fd, temp_path = tempfile.mkstemp(prefix="docextract_", suffix=suffix)
        os.close(fd)
        try:
            os.remove(temp_path)
        except Exception:
            pass
        return temp_path

    def _publish_verified_file(self, temp_path, save_path, label):
        """검증된 로컬 임시 파일을 최종 위치로 복사하고 최종본도 다시 검증한다."""
        self._validate_office_openxml(temp_path, label)
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)

        ext = os.path.splitext(save_path)[1] or ".tmp"
        fd, stage_path = tempfile.mkstemp(prefix=".docextract_", suffix=ext, dir=target_dir)
        os.close(fd)
        try:
            shutil.copyfile(temp_path, stage_path)
            os.replace(stage_path, save_path)
            self._validate_copied_file(temp_path, save_path, f"{label} 최종 파일")
        except Exception:
            if os.path.exists(stage_path):
                try:
                    os.remove(stage_path)
                except Exception:
                    pass
            raise

    def _publish_existing_verified_file(self, source_path, save_path, label):
        """이미 정상 Office 패키지인 파일은 Office를 열지 않고 그대로 복사한다."""
        if os.path.abspath(source_path).lower() == os.path.abspath(save_path).lower():
            raise Exception("원본과 같은 경로로는 복사할 수 없습니다.")
        self._validate_office_openxml(source_path, label, deep=False)
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)

        ext = os.path.splitext(save_path)[1] or ".tmp"
        fd, stage_path = tempfile.mkstemp(prefix=".docextract_", suffix=ext, dir=target_dir)
        os.close(fd)
        try:
            shutil.copyfile(source_path, stage_path)
            os.replace(stage_path, save_path)
            self._validate_copied_file(source_path, save_path, f"{label} 최종 파일")
            self.logger.log(f"{label} 파일 직접 복사 완료: {save_path}")
        except Exception:
            if os.path.exists(stage_path):
                try:
                    os.remove(stage_path)
                except Exception:
                    pass
            raise

    def _try_existing_office_file_copy(self, source_path, save_path, label):
        """정상 OpenXML 파일은 Office COM을 열기 전에 빠르게 검증 복사한다."""
        source_ext = os.path.splitext(source_path)[1].lower()
        target_ext = os.path.splitext(save_path)[1].lower()
        if source_ext != target_ext:
            raise Exception(f"원본 확장자({source_ext})와 저장 확장자({target_ext})가 달라 직접 복사를 건너뜁니다.")
        self._publish_existing_verified_file(source_path, save_path, label)

    def _validate_copied_file(self, source_path, save_path, label):
        """검증된 파일 복사 결과는 전체 ZIP 재검사 대신 크기와 헤더만 확인한다."""
        if not os.path.exists(save_path) or os.path.getsize(save_path) == 0:
            raise Exception(f"{label} 파일이 없거나 비어 있습니다.")
        source_size = os.path.getsize(source_path)
        target_size = os.path.getsize(save_path)
        if source_size != target_size:
            raise Exception(f"{label} 파일 크기가 원본과 다릅니다: {source_size} -> {target_size}")
        source_header = self._read_header_hex(source_path, 4)
        target_header = self._read_header_hex(save_path, 4)
        if source_header != target_header:
            raise Exception(f"{label} 파일 헤더가 원본과 다릅니다: {source_header} -> {target_header}")

    def _validate_office_openxml(self, path, label, deep=True):
        """확장자가 OpenXML이면 실제 ZIP 패키지인지 확인한다."""
        ext = os.path.splitext(path)[1].lower()
        required_members = {
            ".pptx": ["[Content_Types].xml", "ppt/presentation.xml"],
            ".pptm": ["[Content_Types].xml", "ppt/presentation.xml"],
            ".ppsx": ["[Content_Types].xml", "ppt/presentation.xml"],
            ".potx": ["[Content_Types].xml", "ppt/presentation.xml"],
            ".xlsx": ["[Content_Types].xml", "xl/workbook.xml"],
            ".xlsm": ["[Content_Types].xml", "xl/workbook.xml"],
            ".xlsb": ["[Content_Types].xml", "xl/workbook.bin"],
            ".docx": ["[Content_Types].xml", "word/document.xml"],
            ".docm": ["[Content_Types].xml", "word/document.xml"],
        }.get(ext)
        if not required_members:
            return

        if not os.path.exists(path) or os.path.getsize(path) == 0:
            raise Exception(f"{label} 저장 결과 파일이 없거나 비어 있습니다.")

        header = self._read_header_hex(path)
        if not zipfile.is_zipfile(path):
            hint = ""
            if header.startswith("53 43 44 53"):
                hint = " DRM/보안 컨테이너(SCDS)로 저장된 것으로 보입니다."
            raise Exception(
                f"{label} 저장 결과가 정상 {ext} 파일이 아닙니다. header={header}.{hint}"
            )

        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            missing = [member for member in required_members if member not in names]
            if missing:
                raise Exception(f"{label} 저장 결과에 필수 항목이 없습니다: {missing}")
            if deep:
                bad_member = archive.testzip()
                if bad_member:
                    raise Exception(f"{label} 저장 결과 ZIP 항목 손상: {bad_member}")

    def _run_with_heartbeat(self, label, func, interval=10, warn_after=30):
        """긴 COM 호출 동안 로그 파일이 멈춘 것처럼 보이지 않게 주기 로그를 남긴다."""
        stop_event = threading.Event()
        start_time = time.perf_counter()

        def heartbeat():
            warned = False
            while not stop_event.wait(interval):
                elapsed = time.perf_counter() - start_time
                self.logger.log(f"{label} 진행 중... {elapsed:.0f}초 경과")
                if not warned and elapsed >= warn_after:
                    self.logger.log(
                        f"{label} 지연 중: Office 저장 대화상자나 OneDrive 동기화 상태를 확인해 주세요."
                    )
                    warned = True

        thread = threading.Thread(target=heartbeat, daemon=True)
        thread.start()
        try:
            return func()
        finally:
            stop_event.set()

    def _save_native_copy(self, source_doc, save_path, label):
        """Office의 원본 복사 기능으로 서식/개체를 가장 정확하게 보존한다."""
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)

        ext = os.path.splitext(save_path)[1] or ".tmp"
        temp_path = self._make_local_temp_path(ext)
        try:
            self.logger.log(f"{label} 원본 복사 로컬 임시 저장: {temp_path}")
            self._run_with_heartbeat(
                f"{label} SaveCopyAs",
                lambda: source_doc.SaveCopyAs(temp_path),
            )
            self._publish_verified_file(temp_path, save_path, label)
            self.logger.log(f"{label} 원본 복사 저장 완료: {save_path}")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _get_app_process_id(self, app):
        """Office 애플리케이션의 최상위 창 기준 프로세스 ID를 가져온다."""
        hwnd = 0
        for attr in ("HWND", "Hwnd", "hwnd"):
            try:
                hwnd = int(getattr(app, attr))
                if hwnd:
                    break
            except Exception:
                hwnd = 0
        if not hwnd:
            return None

        process_id = wintypes.DWORD()
        try:
            ctypes.windll.user32.GetWindowThreadProcessId(wintypes.HWND(hwnd), ctypes.byref(process_id))
            return int(process_id.value) if process_id.value else None
        except Exception:
            return None

    def _close_office_modal_dialogs(self, app, label):
        """자동 변환 중 남은 Office 저장/경고 대화상자를 닫는다."""
        process_id = self._get_app_process_id(app)
        if not process_id:
            return 0

        user32 = ctypes.windll.user32
        WM_CLOSE = 0x0010
        closed = []
        enum_proc_type = ctypes.WINFUNCTYPE(wintypes.BOOL, wintypes.HWND, wintypes.LPARAM)

        def enum_proc(hwnd, lparam):
            try:
                if not user32.IsWindowVisible(hwnd):
                    return True
                pid = wintypes.DWORD()
                user32.GetWindowThreadProcessId(hwnd, ctypes.byref(pid))
                if int(pid.value) != process_id:
                    return True
                class_name = self._get_window_class_name(hwnd)
                if class_name != "#32770":
                    return True
                title = self._get_window_title(hwnd)
                user32.PostMessageW(hwnd, WM_CLOSE, 0, 0)
                closed.append(title or class_name)
            except Exception:
                pass
            return True

        user32.EnumWindows(enum_proc_type(enum_proc), 0)
        if closed:
            self.logger.log(f"{label} 남은 대화상자 닫기: {len(closed)}개 ({', '.join(closed[:3])})")
            time.sleep(0.2)
        return len(closed)

    def _get_powerpoint_clipboard_package(self):
        """PowerPoint slide copy clipboard package bytes."""
        try:
            import win32clipboard
        except Exception as import_error:
            raise Exception(f"win32clipboard import failed: {import_error}")

        package_format = win32clipboard.RegisterClipboardFormat("PowerPoint 14.0 Slides Package")
        last_error = None
        for _ in range(20):
            try:
                win32clipboard.OpenClipboard()
                try:
                    if win32clipboard.IsClipboardFormatAvailable(package_format):
                        data = win32clipboard.GetClipboardData(package_format)
                        if isinstance(data, bytes) and len(data) > 0:
                            return data
                        last_error = Exception("PowerPoint clipboard package is empty")
                    else:
                        last_error = Exception("PowerPoint clipboard package format is not available")
                finally:
                    win32clipboard.CloseClipboard()
            except Exception as error:
                last_error = error
            time.sleep(0.1)
        raise Exception(str(last_error))

    def _write_ppt_clipboard_package_as_pptx(self, package_data, temp_path, target_ext=".pptx"):
        """Convert the PowerPoint clipboard OPC package into a regular PPTX package."""
        if not zipfile.is_zipfile(io.BytesIO(package_data)):
            header = " ".join(f"{byte:02X}" for byte in package_data[:16])
            raise Exception(f"PowerPoint clipboard package is not ZIP. header={header}")

        macro_content_type = b"application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml"
        pptx_content_type = b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
        with zipfile.ZipFile(io.BytesIO(package_data), "r") as source_archive:
            names = set(source_archive.namelist())
            if "clipboard/presentation.xml" not in names:
                raise Exception("PowerPoint clipboard package has no presentation.xml")

            with zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target_archive:
                for item in source_archive.infolist():
                    if item.is_dir():
                        continue
                    name = item.filename
                    data = source_archive.read(name)
                    if name.startswith("clipboard/"):
                        name = "ppt/" + name[len("clipboard/"):]

                    data = data.replace(b"/clipboard/", b"/ppt/")
                    data = data.replace(b'Target="clipboard/', b'Target="ppt/')
                    data = data.replace(b"Target='clipboard/", b"Target='ppt/")
                    if target_ext.lower() != ".pptm":
                        data = data.replace(macro_content_type, pptx_content_type)
                    target_archive.writestr(name, data)

    def _save_ppt_clipboard_package_copy(self, source_pres, save_path):
        """Copy slides to clipboard and rebuild the PowerPoint slide package without SaveAs."""
        temp_path = self._make_local_temp_path(".pptx")
        try:
            total_slides = source_pres.Slides.Count
            slide_indices = tuple(range(1, total_slides + 1))
            self.logger.log("PPT 클립보드 슬라이드 패키지 복원 시도")

            last_error = None
            package_data = None
            for retry in range(1, self.PPT_CLIPBOARD_RETRY_COUNT + 2):
                try:
                    source_pres.Slides.Range(slide_indices).Copy()
                    time.sleep(max(self.PPT_CLIPBOARD_RETRY_DELAY, 0.2))
                    package_data = self._get_powerpoint_clipboard_package()
                    break
                except Exception as error:
                    last_error = error
                    if retry <= self.PPT_CLIPBOARD_RETRY_COUNT:
                        time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)

            if not package_data:
                raise Exception(f"PPT 클립보드 슬라이드 패키지를 가져오지 못했습니다: {last_error}")

            target_ext = os.path.splitext(save_path)[1].lower() or ".pptx"
            self._write_ppt_clipboard_package_as_pptx(package_data, temp_path, target_ext)
            self._publish_verified_file(temp_path, save_path, "PPT 클립보드 슬라이드 패키지")
            self.logger.log(f"PPT 클립보드 슬라이드 패키지 복원 완료: {save_path}")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _save_ppt_slide_clone(self, source_pres, save_path):
        """PowerPoint 내부 복사/붙여넣기로 슬라이드를 최대한 원본에 가깝게 복제한다."""
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)

        temp_path = self._make_local_temp_path(".pptx")

        app = source_pres.Application
        target_pres = None
        try:
            self.logger.log(f"PPT 슬라이드 복제 로컬 임시 저장: {temp_path}")
            target_pres = app.Presentations.Add()

            try:
                target_pres.PageSetup.SlideWidth = source_pres.PageSetup.SlideWidth
                target_pres.PageSetup.SlideHeight = source_pres.PageSetup.SlideHeight
            except Exception as size_error:
                self.logger.log(f"PPT 슬라이드 크기 복사 실패: {str(size_error)[:60]}")

            while target_pres.Slides.Count > 0:
                target_pres.Slides(1).Delete()

            total_slides = source_pres.Slides.Count
            slide_indices = tuple(range(1, total_slides + 1))

            def copy_slide_range(indices, insert_index):
                last_error = None
                for retry in range(1, self.PPT_CLIPBOARD_RETRY_COUNT + 2):
                    try:
                        source_pres.Slides.Range(indices).Copy()
                        time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
                        target_pres.Slides.Paste(insert_index)
                        return
                    except Exception as paste_error:
                        last_error = paste_error
                        if retry >= self.PPT_CLIPBOARD_RETRY_COUNT + 1:
                            raise
                        time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
                if last_error:
                    raise last_error

            try:
                copy_slide_range(slide_indices, 1)
                self.logger.log(f"  PPT 슬라이드 일괄 복제: {total_slides}/{total_slides}")
            except Exception as bulk_error:
                self.logger.log(f"PPT 슬라이드 일괄 복제 실패, 개별 복제 전환: {str(bulk_error)[:80]}")
                while target_pres.Slides.Count > 0:
                    target_pres.Slides(1).Delete()
                for slide_idx in range(1, total_slides + 1):
                    try:
                        copy_slide_range((slide_idx,), target_pres.Slides.Count + 1)
                    except Exception as paste_error:
                        raise Exception(f"슬라이드 {slide_idx} 복제 실패: {str(paste_error)[:80]}")

                    if slide_idx == 1 or slide_idx == total_slides or slide_idx % 5 == 0:
                        self.logger.log(f"  PPT 슬라이드 복제: {slide_idx}/{total_slides}")

            self._run_with_heartbeat(
                "PPT 슬라이드 복제 SaveAs",
                lambda: target_pres.SaveAs(temp_path, 24),
            )
            self._publish_verified_file(temp_path, save_path, "PPT 슬라이드 복제")
            self.logger.log(f"PPT 슬라이드 복제 저장 완료: {save_path}")
        finally:
            if target_pres is not None:
                try:
                    target_pres.Close()
                except Exception:
                    pass
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _capture_ppt_slide_image(self, source_slide, slide_idx, temp_dir, export_width, export_height):
        """PowerPoint Export가 막힌 문서도 슬라이드를 이미지로 확보한다."""
        img_path = os.path.join(temp_dir, f"visual_slide_{slide_idx}.png")
        export_error = None
        try:
            source_slide.Export(img_path, "PNG", export_width, export_height)
            if os.path.exists(img_path) and os.path.getsize(img_path) > 0:
                return img_path
            export_error = Exception("PowerPoint Export created an empty image")
        except Exception as error:
            export_error = error

        self.logger.log(
            f"  슬라이드 {slide_idx} Export 실패, 클립보드 이미지 복사 시도: {str(export_error)[:80]}"
        )
        clipboard_error = None
        for retry in range(1, self.PPT_CLIPBOARD_RETRY_COUNT + 2):
            try:
                source_slide.Copy()
                for _ in range(12):
                    time.sleep(max(self.PPT_CLIPBOARD_RETRY_DELAY, 0.1))
                    clipboard_path = self._get_image_from_clipboard(temp_dir)
                    if clipboard_path and os.path.exists(clipboard_path) and os.path.getsize(clipboard_path) > 0:
                        return clipboard_path
                clipboard_error = Exception("클립보드에서 사용 가능한 이미지를 찾지 못했습니다.")
            except Exception as error:
                clipboard_error = error
            if retry <= self.PPT_CLIPBOARD_RETRY_COUNT:
                time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)

        raise Exception(
            f"슬라이드 {slide_idx} 이미지 캡처 실패. "
            f"Export={str(export_error)[:80]}, Clipboard={str(clipboard_error)[:80]}"
        )

    def _save_ppt_visual_copy(self, source_pres, save_path):
        """각 슬라이드를 전체 이미지로 저장해 화면 배치/서식을 가장 정확하게 보존한다."""
        if not HAS_PPTX:
            raise Exception("python-pptx 패키지가 필요합니다. pip install python-pptx")

        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        temp_dir = tempfile.mkdtemp(prefix="docextract_visual_")
        temp_path = self._make_local_temp_path(".pptx")

        try:
            visual_pres = Presentation()
            slide_width = source_pres.PageSetup.SlideWidth
            slide_height = source_pres.PageSetup.SlideHeight
            visual_pres.slide_width = Emu(int(slide_width * 12700))
            visual_pres.slide_height = Emu(int(slide_height * 12700))
            blank_layout = visual_pres.slide_layouts[6]
            export_width = max(1, int(slide_width * 2))
            export_height = max(1, int(slide_height * 2))
            total_slides = source_pres.Slides.Count

            self.logger.log(f"PPT 화면 그대로 이미지 저장 시작: {save_path}")
            for slide_idx in range(1, total_slides + 1):
                img_path = self._capture_ppt_slide_image(
                    source_pres.Slides(slide_idx),
                    slide_idx,
                    temp_dir,
                    export_width,
                    export_height,
                )
                slide = visual_pres.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    img_path,
                    Emu(0),
                    Emu(0),
                    visual_pres.slide_width,
                    visual_pres.slide_height,
                )
                if slide_idx == 1 or slide_idx == total_slides or slide_idx % 5 == 0:
                    self.logger.log(f"  PPT 화면 그대로 이미지: {slide_idx}/{total_slides}")

            visual_pres.save(temp_path)
            self._publish_verified_file(temp_path, save_path, "PPT 화면 그대로")
            self.logger.log(f"PPT 화면 그대로 이미지 저장 완료: {save_path}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _try_save_ppt_visual_companion(self, source_pres, save_path):
        visual_path = self._add_filename_suffix(save_path, "_화면그대로")
        try:
            self._save_ppt_visual_copy(source_pres, visual_path)
            self.logger.log(f"PPT 화면 그대로 추가본 생성: {visual_path}")
            return visual_path
        except Exception as visual_error:
            self.logger.log(f"PPT 화면 그대로 추가본 생성 실패: {str(visual_error)[:100]}")
            return None

    def _save_hwp_document(self, hwp, save_path, save_format):
        """한글 버전별 SaveAs 인자 차이를 흡수한다."""
        hwp_format = "HWPX" if save_format == "hwpx" else "HWP"
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)

        try:
            hwp.HAction.GetDefault("FileSaveAs_S", hwp.HParameterSet.HFileOpenSave.HSet)
            file_open_save = hwp.HParameterSet.HFileOpenSave
            file_open_save.filename = save_path
            try:
                file_open_save.FileName = save_path
            except Exception:
                pass
            file_open_save.Format = hwp_format
            result = hwp.HAction.Execute("FileSaveAs_S", file_open_save.HSet)
            if result is False:
                raise Exception("FileSaveAs_S 액션이 실패를 반환했습니다.")
        except Exception as action_error:
            self.logger.log(f"한글 FileSaveAs_S 실패, SaveAs 재시도: {str(action_error)[:80]}")
            try:
                hwp.SaveAs(save_path, hwp_format, "")
            except Exception as saveas_error:
                self.logger.log(f"한글 SaveAs 3인자 실패, 2인자 재시도: {str(saveas_error)[:80]}")
                hwp.SaveAs(save_path, hwp_format)

        for _ in range(20):
            if os.path.exists(save_path) and os.path.getsize(save_path) > 0:
                header = self._read_header_hex(save_path, 8)
                if header.startswith("53 43 44 53"):
                    self.logger.log("한글 저장 결과가 회사 보안/DRM 컨테이너(SCDS)라서 실패 처리합니다")
                    try:
                        os.remove(save_path)
                    except Exception:
                        pass
                    raise Exception(self._hwp_drm_container_message(save_path))
                return
            time.sleep(0.1)
        raise Exception("한글 저장 후 결과 파일이 생성되지 않았거나 비어 있습니다.")

    def _hwp_extract_hwpml(self, hwp):
        """열린 한글 문서를 파일 저장 없이 HWPML2X(완전 구조 XML) 문자열로 추출한다.

        GetTextFile은 메모리의 복호화된 내용을 문자열로 돌려주므로, 파일 SaveAs 단계에서
        보안 래퍼가 SCDS 컨테이너로 다시 감싸는 것을 우회한다(Word의 WordOpenXML 재구성과 동일 원리).
        """
        try:
            xml = hwp.GetTextFile("HWPML2X", "")
        except Exception as get_error:
            raise Exception(f"HWPML2X 메모리 추출 실패: {str(get_error)[:120]}")
        if not xml or len(xml) < 50:
            raise Exception(f"HWPML2X 추출 결과가 비어 있습니다 (len={len(xml) if xml else 0})")
        self.logger.log(f"HWPML2X 메모리 추출 성공: {len(xml)}자")
        return xml

    def _hwp_rebuild_via_hwpml(self, hwp, save_path, save_format):
        """원본 메모리의 HWPML2X를 새 한글 문서로 재구성해 일반 HWP/HWPX로 저장한다.

        새 문서는 원본의 보안 컨텍스트가 없어 SaveAs가 SCDS로 감기지 않을 가능성이 높다.
        한컴 권고대로 Copy/Paste 대신 SetTextFile로 내용을 옮긴다.
        """
        xml = self._hwp_extract_hwpml(hwp)
        self.logger.log("HWPML2X 새 문서 재구성 시작")
        hwp.XHwpDocuments.Add(0)  # 0 = 새 창에 빈 문서, 자동 활성화
        try:
            hwp.SetTextFile(xml, "HWPML2X", "")
            self._save_hwp_document(hwp, save_path, save_format)  # SCDS 헤더 검증 포함
            self.logger.log(f"HWPML2X 재구성 저장 완료: {save_path}")
        finally:
            # 재구성용 새 문서만 닫고 원본은 보존한다.
            try:
                hwp.XHwpDocuments.Active_XHwpDocument.Close(0)
            except Exception:
                try:
                    hwp.HAction.Run("FileClose")
                except Exception:
                    pass

    def _hwp_save_hwpml_direct(self, hwp, save_path):
        """HWPML2X를 파이썬이 직접 파일로 기록한다(한글 SaveAs 미사용 → DRM 확실 우회).

        산출물은 한글에서 열리는 .hwpml(XML) 형식으로, 원본 구조를 보존한다.
        """
        xml = self._hwp_extract_hwpml(hwp)
        hwpml_path = os.path.splitext(save_path)[0] + ".hwpml"
        with open(hwpml_path, "w", encoding="utf-8") as hwpml_file:
            hwpml_file.write(xml)
        if not (os.path.exists(hwpml_path) and os.path.getsize(hwpml_path) > 0):
            raise Exception("HWPML 직접 저장 결과 파일이 비어 있습니다.")
        self.logger.log(f"HWPML 직접 저장 완료: {hwpml_path} ({os.path.getsize(hwpml_path)} bytes)")
        return hwpml_path

    def _hwp_extract_success(self, result_path, extract_start, method_label, extra=""):
        """한글 추출 성공 시 진행바/상태/완료 안내를 공통 처리한다."""
        message = f"한글 추출 완료 ({method_label})!\n{result_path}"
        if extra:
            message += f"\n\n{extra}"
        self.root.after(0, lambda: self.progress_var.set(100))
        self.root.after(0, lambda: self.status_text.set("한글 추출 완료!"))
        self.root.after(0, lambda: messagebox.showinfo("완료", message))
        self.logger.log(f"한글 추출 완료({method_label}): {result_path}")
        self._log_elapsed("한글 전체 추출 시간", extract_start)

    def _hwp_write_with_fallbacks(self, hwp, save_path, save_format):
        """연결된 한글 문서를 3단계 폴백으로 저장하고 (결과경로, 방법)을 반환한다.

        (1) 직접 SaveAs → (2) HWPML2X 메모리 재구성 → (3) HWPML2X 직접 .hwpml 기록.
        SaveAs는 DRM에서 SCDS로 감기므로 메모리 추출(GetTextFile) 경로가 본 우회책이다.
        UI 메시지 없이 저장만 하므로 탭/직접/일괄 변환이 공유한다. 3단계 모두 실패 시 예외.
        """
        # 방법 1: 직접 SaveAs (정상 HWP면 성공, DRM이면 SCDS 헤더로 실패 처리됨)
        try:
            self.logger.log(f"방법1 직접 SaveAs 시도: {save_path}")
            self._save_hwp_document(hwp, save_path, save_format)
            return save_path, "직접 저장"
        except Exception as save_error:
            self.logger.log(f"방법1 SaveAs 실패(보안 컨테이너 가능성): {str(save_error)[:150]}")

        # 방법 2: HWPML2X 메모리 재구성 → 일반 HWP (원본 구조 보존, 파일 저장 우회)
        try:
            self.logger.log("방법2 HWPML2X 메모리 재구성 시도")
            self._hwp_rebuild_via_hwpml(hwp, save_path, save_format)
            return save_path, "메모리 재구성"
        except Exception as rebuild_error:
            self.logger.log(f"방법2 HWPML2X 재구성 실패: {str(rebuild_error)[:150]}")

        # 방법 3: HWPML2X를 파이썬이 직접 .hwpml로 기록 (한글 SaveAs 미사용, 가장 확실한 우회)
        self.logger.log("방법3 HWPML 직접 저장 시도")
        hwpml_path = self._hwp_save_hwpml_direct(hwp, save_path)
        return hwpml_path, "HWPML 직접 저장"

    def _convert_hwp_file(self, source_path, save_path):
        """원본 HWP 파일을 새 한글 인스턴스로 열어 메모리 추출 후 저장한다(직접/일괄 변환용).

        UI 메시지 없이 동작하며, 완료·오류 안내는 호출자(직접/일괄 변환 흐름)가 담당한다.
        """
        if not os.path.exists(source_path):
            raise Exception(f"원본 파일을 찾을 수 없습니다: {source_path}")
        save_format = "hwpx" if os.path.splitext(save_path)[1].lower() == ".hwpx" else "hwp"
        hwp, created = self._get_hwp_app(allow_dispatch=True)
        self.logger.log(f"한글 파일 변환 시작(created={created}): {source_path}")
        try:
            try:
                opened = hwp.Open(source_path, "HWP", "forceopen:true")
            except Exception as open_error:
                self.logger.log(f"Open(3인자) 실패, 단순 Open 재시도: {str(open_error)[:100]}")
                opened = hwp.Open(source_path)
            if opened is False:
                raise Exception("한글이 원본 파일을 열지 못했습니다(보안 차단 가능성).")
            result_path, method = self._hwp_write_with_fallbacks(hwp, save_path, save_format)
            self.logger.log(f"한글 파일 변환 완료({method}): {result_path}")
        finally:
            if created:
                try:
                    hwp.Quit()
                except Exception:
                    pass

    def _hwp_save_with_fallbacks(self, hwp, save_path, save_format, extract_start):
        """탭 변환용: 3단계 폴백 저장 + 진행바/완료 안내."""
        self.root.after(0, lambda: self.status_text.set("메모리에서 원본 구조 추출 중..."))
        self.root.after(0, lambda: self.progress_var.set(60))
        try:
            result_path, method = self._hwp_write_with_fallbacks(hwp, save_path, save_format)
        except Exception as hwpml_error:
            self.logger.log(f"방법3 HWPML 직접 저장 실패: {str(hwpml_error)[:150]}")
            raise Exception(
                "한글 변환에 실패했습니다.\n\n"
                "직접 저장·메모리 재구성·HWPML 추출이 모두 막혔습니다.\n"
                "회사 보안(DRM)이 메모리 추출까지 차단하는 환경일 수 있습니다."
            )
        extra = ""
        if result_path.lower().endswith(".hwpml"):
            extra = "원본 구조를 보존한 .hwpml 파일입니다. 한글에서 열어 .hwp로 다시 저장할 수 있습니다."
        self._hwp_extract_success(result_path, extract_start, method, extra=extra)

    def _extract_hwp_from_file(self, source_path, save_path, save_format, extract_start):
        """원본 HWP 파일을 새 한글 인스턴스로 직접 열어 메모리에서 추출/재구성한다.

        열린 문서 COM 연결(ROT)이 막힌 환경 대응: Dispatch로 새 인스턴스를 만들고 파일을
        Open하면 보안 모듈이 사용자 권한으로 복호화하므로, GetTextFile로 메모리에서 내용을
        빼내 SaveAs(SCDS 재포장)를 우회한다.
        """
        if not os.path.exists(source_path):
            raise Exception(f"원본 파일을 찾을 수 없습니다: {source_path}")

        self.root.after(0, lambda: self.status_text.set("새 한글 인스턴스로 원본 여는 중..."))
        self.root.after(0, lambda: self.progress_var.set(25))
        hwp, created = self._get_hwp_app(allow_dispatch=True)
        self.logger.log(f"한글 인스턴스 확보(created={created}), 원본 Open 시도: {source_path}")
        try:
            try:
                opened = hwp.Open(source_path, "HWP", "forceopen:true")
            except Exception as open_error:
                self.logger.log(f"Open(3인자) 실패, 단순 Open 재시도: {str(open_error)[:100]}")
                opened = hwp.Open(source_path)
            if opened is False:
                raise Exception("한글이 원본 파일을 열지 못했습니다(보안 차단 가능성).")
            self.logger.log("원본 파일 Open 성공 → 메모리 추출 단계로")
            self.root.after(0, lambda: self.progress_var.set(45))
            self._hwp_save_with_fallbacks(hwp, save_path, save_format, extract_start)
        finally:
            if created:
                try:
                    hwp.Quit()
                except Exception:
                    pass

    def _copy_word_document_file(self, source_doc, save_path):
        """저장된 Word 원본 파일을 원본 상태 변경 없이 복사한다."""
        try:
            source_path = source_doc.FullName
        except Exception:
            source_path = ""

        if not source_path or not os.path.exists(source_path):
            raise Exception("Word 원본 파일 경로를 확인할 수 없습니다. 문서를 먼저 저장해 주세요.")

        source_ext = os.path.splitext(source_path)[1].lower()
        target_ext = os.path.splitext(save_path)[1].lower()
        if source_ext and target_ext and source_ext != target_ext:
            raise Exception(f"원본 확장자({source_ext})와 저장 확장자({target_ext})가 달라 파일 복사를 생략합니다.")

        try:
            if not source_doc.Saved:
                raise Exception("Word 문서에 저장되지 않은 변경사항이 있습니다. 저장 후 다시 시도해 주세요.")
        except Exception as saved_error:
            if "저장되지 않은 변경사항" in str(saved_error):
                raise

        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        if os.path.abspath(source_path).lower() == os.path.abspath(save_path).lower():
            raise Exception("원본과 같은 경로로는 복사할 수 없습니다.")

        shutil.copy2(source_path, save_path)
        if not os.path.exists(save_path) or os.path.getsize(save_path) <= 0:
            raise Exception("Word 파일 복사 후 결과 파일이 생성되지 않았거나 비어 있습니다.")
        self._validate_office_openxml(save_path, "Word")

    def _write_word_flat_opc_as_docx(self, flat_xml, temp_path):
        """WordOpenXML(Flat OPC)을 일반 DOCX ZIP 패키지로 변환한다."""
        try:
            from lxml import etree
        except Exception as import_error:
            raise Exception(f"WordOpenXML 변환에는 lxml 패키지가 필요합니다: {import_error}")

        pkg_ns = "http://schemas.microsoft.com/office/2006/xmlPackage"
        content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"

        if isinstance(flat_xml, str):
            root = etree.fromstring(flat_xml.encode("utf-8"))
        else:
            root = etree.fromstring(flat_xml)

        part_content_types = []
        written_names = set()
        with zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as archive:
            for part in root.xpath("//pkg:part", namespaces={"pkg": pkg_ns}):
                raw_name = part.get(f"{{{pkg_ns}}}name") or part.get("name")
                if not raw_name:
                    continue
                zip_name = raw_name.lstrip("/")
                content_type = part.get(f"{{{pkg_ns}}}contentType") or part.get("contentType")
                if content_type:
                    part_name = raw_name if raw_name.startswith("/") else f"/{raw_name}"
                    part_content_types.append((part_name, content_type))

                xml_data = part.find(f"{{{pkg_ns}}}xmlData")
                binary_data = part.find(f"{{{pkg_ns}}}binaryData")
                if xml_data is not None:
                    children = [child for child in xml_data if isinstance(child.tag, str)]
                    data = (
                        etree.tostring(children[0], encoding="UTF-8")
                        if children else b""
                    )
                elif binary_data is not None:
                    encoded = "".join("".join(binary_data.itertext()).split())
                    data = base64.b64decode(encoded) if encoded else b""
                else:
                    data = b""

                archive.writestr(zip_name, data)
                written_names.add(zip_name)

            if "[Content_Types].xml" not in written_names:
                types = etree.Element(f"{{{content_types_ns}}}Types", nsmap={None: content_types_ns})
                etree.SubElement(
                    types,
                    f"{{{content_types_ns}}}Default",
                    Extension="rels",
                    ContentType="application/vnd.openxmlformats-package.relationships+xml",
                )
                etree.SubElement(
                    types,
                    f"{{{content_types_ns}}}Default",
                    Extension="xml",
                    ContentType="application/xml",
                )
                for part_name, content_type in part_content_types:
                    etree.SubElement(
                        types,
                        f"{{{content_types_ns}}}Override",
                        PartName=part_name,
                        ContentType=content_type,
                    )
                archive.writestr(
                    "[Content_Types].xml",
                    etree.tostring(types, encoding="UTF-8", xml_declaration=True),
                )

    def _save_word_openxml_copy(self, source_doc, save_path):
        """DRM 컨테이너 파일 복사가 실패할 때 Word 내부 OOXML로 구조를 복원한다."""
        target_ext = os.path.splitext(save_path)[1].lower()
        if target_ext != ".docx":
            raise Exception("WordOpenXML 구조 복원은 .docx 저장만 지원합니다.")

        temp_path = self._make_local_temp_path(".docx")
        try:
            self.logger.log("Word WordOpenXML 구조 복원 시도")
            flat_xml = source_doc.WordOpenXML
            if not flat_xml:
                raise Exception("WordOpenXML 데이터가 비어 있습니다.")

            self._write_word_flat_opc_as_docx(flat_xml, temp_path)
            self._publish_verified_file(temp_path, save_path, "Word WordOpenXML")
            self.logger.log(f"Word WordOpenXML 구조 복원 완료: {save_path}")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _setup_ppt_tab(self):
        """PPT 탭 설정"""
        tab = self.ppt_tab

        # 문서 정보 프레임
        info_frame = self._create_section(tab, "PPT 입력 선택")

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=2)
        ttk.Label(source_inner, text="파일 선택:", width=12).pack(side=tk.LEFT)
        self.ppt_source_entry = ttk.Entry(source_inner, textvariable=self.ppt_source_path, width=45, state="readonly")
        self.ppt_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_ppt_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # PPT 선택 콤보박스
        select_frame = ttk.Frame(info_frame)
        select_frame.pack(fill=tk.X, pady=2)
        ttk.Label(select_frame, text="PPT 선택:", width=12).pack(side=tk.LEFT)
        self.ppt_combo = ttk.Combobox(select_frame, state="readonly", width=40)
        self.ppt_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.ppt_combo.bind("<<ComboboxSelected>>", self.on_ppt_selected)

        # 슬라이드 수
        slide_frame = ttk.Frame(info_frame)
        slide_frame.pack(fill=tk.X, pady=2)
        ttk.Label(slide_frame, text="슬라이드 수:", width=12).pack(side=tk.LEFT)
        ttk.Label(slide_frame, textvariable=self.ppt_slide_count,
                  font=("맑은 고딕", 10, "bold")).pack(side=tk.LEFT)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=lambda: self.detect_open_ppt(prefer_open=True),
                   style="Secondary.TButton").pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = self._create_section(tab, "새 파일 저장 위치")

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.ppt_save_path, width=45).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_ppt_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # 추출 모드 프레임
        mode_frame = self._create_section(tab, "추출 모드")

        self.ppt_extract_mode = tk.StringVar(value="native_copy")

        ttk.Radiobutton(mode_frame, text="원본 그대로 복사 (서식/넓이/높이/도형 보존)",
                        variable=self.ppt_extract_mode, value="native_copy").pack(anchor=tk.W)
        ttk.Radiobutton(mode_frame, text="하이브리드 (편집 가능, 느림: 도형 속성 재생성)",
                        variable=self.ppt_extract_mode, value="hybrid").pack(anchor=tk.W)
        ttk.Radiobutton(mode_frame, text="텍스트 중심 + 객체 보존 (도형/이미지는 그림으로 포함)",
                        variable=self.ppt_extract_mode, value="text_only").pack(anchor=tk.W)

        # 추출 버튼
        self.ppt_extract_button = ttk.Button(tab, text="새 PPT로 내보내기",
                                              command=self.start_ppt_extraction,
                                              style="Accent.TButton")
        self.ppt_extract_button.pack(pady=10)

    def _setup_excel_tab(self):
        """Excel 탭 설정"""
        tab = self.excel_tab

        # 문서 정보 프레임
        info_frame = self._create_section(tab, "Excel 입력 선택")

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=2)
        ttk.Label(source_inner, text="파일 선택:", width=12).pack(side=tk.LEFT)
        self.excel_source_entry = ttk.Entry(source_inner, textvariable=self.excel_source_path, width=45, state="readonly")
        self.excel_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_excel_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # Excel 선택 콤보박스
        select_frame = ttk.Frame(info_frame)
        select_frame.pack(fill=tk.X, pady=2)
        ttk.Label(select_frame, text="Excel 선택:", width=12).pack(side=tk.LEFT)
        self.excel_combo = ttk.Combobox(select_frame, state="readonly", width=40)
        self.excel_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.excel_combo.bind("<<ComboboxSelected>>", self.on_excel_selected)

        # 시트 수
        sheet_frame = ttk.Frame(info_frame)
        sheet_frame.pack(fill=tk.X, pady=2)
        ttk.Label(sheet_frame, text="시트 수:", width=12).pack(side=tk.LEFT)
        ttk.Label(sheet_frame, textvariable=self.excel_sheet_count,
                  font=("맑은 고딕", 10, "bold")).pack(side=tk.LEFT)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=lambda: self.detect_open_excel(prefer_open=True),
                   style="Secondary.TButton").pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = self._create_section(tab, "새 파일 저장 위치")

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.excel_save_path, width=45).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_excel_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # 추출 옵션 프레임
        option_frame = self._create_section(tab, "추출 옵션")

        self.excel_include_format = tk.BooleanVar(value=False)
        self.excel_include_formulas = tk.BooleanVar(value=False)
        self.excel_native_copy = tk.BooleanVar(value=True)

        ttk.Checkbutton(option_frame, text="원본 그대로 복사 우선 (서식/넓이/높이/도형 보존)",
                        variable=self.excel_native_copy).pack(anchor=tk.W)
        ttk.Checkbutton(option_frame, text="서식 포함 (느림: 글꼴, 색상, 행/열 크기)",
                        variable=self.excel_include_format).pack(anchor=tk.W)
        ttk.Checkbutton(option_frame, text="수식 대신 값만 저장",
                        variable=self.excel_include_formulas).pack(anchor=tk.W)

        # 추출 버튼
        self.excel_extract_button = ttk.Button(tab, text="새 Excel로 내보내기",
                                                command=self.start_excel_extraction,
                                                style="Accent.TButton")
        self.excel_extract_button.pack(pady=10)

    def _setup_hwp_tab(self):
        """한글 탭 설정"""
        tab = self.hwp_tab

        info_frame = self._create_section(tab, "원본 한글 파일 선택")
        ttk.Label(
            info_frame,
            text="원본 .hwp 파일을 고르면 새 한글 인스턴스로 열어 메모리에서 추출합니다(보안 PC 대응, 파일 저장 우회).\n"
                 "파일을 이 칸에 끌어다 놓아도 됩니다.",
            font=("맑은 고딕", 9), justify=tk.LEFT,
        ).pack(anchor=tk.W, pady=2)

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=5)
        ttk.Label(source_inner, text="한글 선택:", width=12).pack(side=tk.LEFT)
        self.hwp_source_entry = ttk.Entry(source_inner, textvariable=self.hwp_source_path,
                                          width=45, state="readonly")
        self.hwp_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_hwp_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        path_frame = self._create_section(tab, "새 파일 저장 위치")
        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.hwp_save_path, width=45).pack(
            side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_hwp_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        format_frame = self._create_section(tab, "저장 형식")
        self.hwp_save_format = tk.StringVar(value="hwp")
        format_inner = ttk.Frame(format_frame)
        format_inner.pack(fill=tk.X)
        ttk.Radiobutton(format_inner, text="HWP (한글 문서)",
                        variable=self.hwp_save_format, value="hwp").pack(side=tk.LEFT, padx=(0, 16))
        ttk.Radiobutton(format_inner, text="HWPX (한글 2014 이상)",
                        variable=self.hwp_save_format, value="hwpx").pack(side=tk.LEFT)

        self.hwp_extract_button = ttk.Button(tab, text="원본 파일 선택 후 변환하기",
                                             command=self.start_hwp_extraction,
                                             style="Accent.TButton")
        self.hwp_extract_button.pack(pady=10)

    def _on_tab_changed(self, event):
        """이전 Notebook 이벤트 호환용 감지 예약."""
        self._schedule_detect()

    def _select_document_view(self, index, detect=True):
        """좌측 사이드바 선택에 맞춰 작업 패널을 전환한다."""
        self.current_doc_index = index
        title, _badge, summary, _detect_fn = self.doc_views[index]
        self.view_title_text.set(title)
        self.view_summary_text.set(summary)

        for frame_index, frame in enumerate(getattr(self, "content_frames", [])):
            if frame_index == index:
                frame.tkraise()

        c = self.ui_colors
        for button_index, widgets in enumerate(self.nav_buttons):
            item, badge_label, text_box, title_label, summary_label = widgets
            if button_index == index:
                item.configure(bg=c["nav_selected_bg"])
                text_box.configure(bg=c["nav_selected_bg"])
                badge_label.configure(bg=c["accent"], fg="#ffffff")
                title_label.configure(bg=c["nav_selected_bg"], fg=c["nav_selected_fg"])
                summary_label.configure(bg=c["nav_selected_bg"], fg=c["nav_selected_fg"])
            else:
                item.configure(bg=c["nav_bg"])
                text_box.configure(bg=c["nav_bg"])
                badge_label.configure(bg="#f2f5f9", fg=c["nav_muted"])
                title_label.configure(bg=c["nav_bg"], fg=c["nav_fg"])
                summary_label.configure(bg=c["nav_bg"], fg=c["nav_muted"])

        if detect:
            self.status_text.set(f"{title} 감지 준비")
            self._schedule_detect()

    def _schedule_detect(self):
        """현재 선택된 문서 종류를 debounce 후 감지한다."""
        if hasattr(self, '_pending_detect') and self._pending_detect:
            self.root.after_cancel(self._pending_detect)
            self._pending_detect = None

        self._pending_detect = self.root.after(50, self._do_detect)

    def _do_detect(self):
        """실제 감지 실행"""
        self._pending_detect = None
        current_tab = self.current_doc_index

        # debounce가 중복 이벤트를 정리하므로 탭 이동 시마다 최신 상태를 다시 확인한다.
        self.tab_detected[current_tab] = True

        # 해당 탭 감지 실행 (doc_views 순서와 무관하게 감지 함수로 직접 분기)
        detect_fn = self.doc_views[current_tab][3]
        if detect_fn is not None:
            detect_fn()
        elif self.content_frames[current_tab] is self.hwp_tab:
            self.status_text.set("원본 한글 파일을 선택하세요")
        elif self.content_frames[current_tab] is self.pdf_tab:
            self.status_text.set("보안 해제할 PDF 파일을 선택하세요")
        elif self.content_frames[current_tab] is self.batch_tab:
            self.status_text.set("일괄 변환 파일을 추가하세요")

    def _make_unique_output_path_with_ext(self, output_dir, source_path, ext):
        stem = os.path.splitext(os.path.basename(source_path))[0]
        candidate = os.path.join(output_dir, f"{stem}_복사본{ext}")
        if not os.path.exists(candidate):
            return candidate
        for index in range(2, 1000):
            candidate = os.path.join(output_dir, f"{stem}_복사본_{index}{ext}")
            if not os.path.exists(candidate):
                return candidate
        raise Exception(f"출력 파일명을 만들 수 없습니다: {source_path}")

    def _default_direct_save_path(self, source_path, kind, preferred_ext=None):
        output_dir = os.path.dirname(os.path.abspath(source_path)) or os.getcwd()
        if preferred_ext:
            return self._make_unique_output_path_with_ext(output_dir, source_path, preferred_ext)
        return self._make_unique_output_path(output_dir, source_path, kind)

    def _apply_source_file_selection(self, kind, source_var, save_var, path, label, preferred_ext=None):
        if not path:
            return
        source_var.set(path)
        self._use_direct_file_input(kind, path, label)
        save_var.set(self._default_direct_save_path(path, kind, preferred_ext))
        self.status_text.set(f"{label} 파일 선택됨")

    def _parse_drop_paths(self, data):
        if not data:
            return []
        try:
            raw_items = self.root.tk.splitlist(data)
        except Exception:
            raw_items = [data]
        paths = []
        for item in raw_items:
            path = str(item).strip()
            if path.startswith("file:///"):
                path = path[8:]
            if path:
                paths.append(os.path.normpath(path))
        return paths

    def _expand_supported_drop_paths(self, paths):
        supported = []
        for raw_path in paths:
            path = os.path.abspath(raw_path)
            if os.path.isdir(path):
                for root_dir, _dirs, files in os.walk(path):
                    for filename in files:
                        file_path = os.path.join(root_dir, filename)
                        if self._batch_file_kind(file_path):
                            supported.append(file_path)
            elif os.path.isfile(path) and self._batch_file_kind(path):
                supported.append(path)
        return supported

    def _register_drop_target(self, widget, callback):
        if not HAS_TKINTERDND or widget is None:
            return
        try:
            widget.drop_target_register(DND_FILES)
            widget.dnd_bind("<<Drop>>", callback)
        except Exception as error:
            self.logger.log(f"드롭 대상 등록 실패: {widget} ({str(error)[:80]})")

    def _setup_drag_drop(self):
        self._register_drop_target(
            self.ppt_source_entry,
            lambda event: self._handle_direct_file_drop(
                event, "ppt", self.ppt_source_path, self.ppt_save_path, "PPT"
            ),
        )
        self._register_drop_target(
            self.excel_source_entry,
            lambda event: self._handle_direct_file_drop(
                event, "excel", self.excel_source_path, self.excel_save_path, "Excel"
            ),
        )
        self._register_drop_target(
            self.word_source_entry,
            lambda event: self._handle_direct_file_drop(
                event, "word", self.word_source_path, self.word_save_path, "Word"
            ),
        )
        self._register_drop_target(
            self.notepad_source_entry,
            lambda event: self._handle_direct_file_drop(
                event,
                "text",
                self.notepad_source_path,
                self.notepad_save_path,
                "TXT",
                ".docx" if self.notepad_save_format.get() == "docx" else ".txt",
            ),
        )
        self._register_drop_target(
            self.hwp_source_entry,
            lambda event: self._handle_direct_file_drop(
                event, "hwp", self.hwp_source_path, self.hwp_save_path, "한글"
            ),
        )
        for widget in (self.batch_tab, self.batch_file_listbox):
            self._register_drop_target(widget, self._handle_batch_file_drop)
        self.logger.log("드래그앤드롭 활성화")

    def _handle_direct_file_drop(self, event, kind, source_var, save_var, label, preferred_ext=None):
        paths = self._parse_drop_paths(getattr(event, "data", ""))
        files = self._expand_supported_drop_paths(paths)
        matching_files = [path for path in files if self._batch_file_kind(path) == kind]
        if not matching_files:
            self.status_text.set(f"{label} 지원 파일을 드롭해주세요")
            return
        selected_path = matching_files[0]
        self._apply_source_file_selection(kind, source_var, save_var, selected_path, label, preferred_ext)
        if len(matching_files) > 1:
            self.status_text.set(f"{label} 파일 1개 선택됨, 나머지는 일괄 변환에 드롭하세요")
        return "break"

    def _handle_batch_file_drop(self, event):
        paths = self._parse_drop_paths(getattr(event, "data", ""))
        added = self._add_batch_paths(paths)
        if added:
            self._select_document_view(self.content_frames.index(self.batch_tab), detect=False)
            self.status_text.set(f"일괄 변환 파일 {added}개 추가됨")
        else:
            self.batch_status_text.set("추가할 수 있는 파일이 없습니다.")
        return "break"

    def _set_input_mode(self, kind, mode):
        attr = {
            "ppt": "ppt_input_mode",
            "excel": "excel_input_mode",
            "word": "word_input_mode",
            "text": "notepad_input_mode",
        }.get(kind)
        if attr:
            setattr(self, attr, mode)

    def _is_direct_file_input_active(self, kind):
        attr = {
            "ppt": "ppt_input_mode",
            "excel": "excel_input_mode",
            "word": "word_input_mode",
            "text": "notepad_input_mode",
        }.get(kind)
        return bool(attr and getattr(self, attr, "open") == "file")

    def _show_direct_file_input(self, kind, source_path, label):
        file_name = os.path.basename(source_path) if source_path else f"{label} 파일 선택됨"
        if kind == "ppt":
            if hasattr(self, "ppt_combo"):
                self.ppt_combo.set("")
            self.selected_ppt_index.set(0)
            self.ppt_doc_name.set(file_name)
            self.ppt_slide_count.set("-")
        elif kind == "excel":
            if hasattr(self, "excel_combo"):
                self.excel_combo.set("")
            self.selected_excel_index.set(0)
            self.excel_doc_name.set(file_name)
            self.excel_sheet_count.set("-")
        elif kind == "word":
            if hasattr(self, "word_combo"):
                self.word_combo.set("")
            self.selected_word_index.set(0)
            self.word_doc_name.set(file_name)
            self.word_page_count.set("-")
        elif kind == "text":
            if hasattr(self, "notepad_combo"):
                self.notepad_combo.set("")
            self.notepad_doc_name.set(file_name)

    def _use_direct_file_input(self, kind, source_path, label):
        self._set_input_mode(kind, "file")
        self._show_direct_file_input(kind, source_path, label)
        self.logger.log(f"{label} 파일 직접 선택 사용: {source_path}")

    def _use_open_document_input(self, kind, source_var, label):
        self._set_input_mode(kind, "open")
        if source_var.get().strip():
            source_var.set("")
            self.logger.log(f"{label} 열린 문서 선택으로 파일 직접 선택 경로 초기화")

    def browse_ppt_source_path(self):
        path = filedialog.askopenfilename(
            title="변환할 PPT 파일 선택",
            filetypes=[
                ("PowerPoint", "*.ppt;*.pptx;*.pptm;*.ppsx;*.potx"),
                ("모든 파일", "*.*"),
            ],
        )
        self._apply_source_file_selection("ppt", self.ppt_source_path, self.ppt_save_path, path, "PPT")

    def browse_excel_source_path(self):
        path = filedialog.askopenfilename(
            title="변환할 Excel 파일 선택",
            filetypes=[
                ("Excel", "*.xls;*.xlsx;*.xlsm;*.xlsb"),
                ("모든 파일", "*.*"),
            ],
        )
        self._apply_source_file_selection("excel", self.excel_source_path, self.excel_save_path, path, "Excel")

    def browse_word_source_path(self):
        path = filedialog.askopenfilename(
            title="변환할 Word 파일 선택",
            filetypes=[
                ("Word", "*.doc;*.docx;*.docm"),
                ("모든 파일", "*.*"),
            ],
        )
        self._apply_source_file_selection("word", self.word_source_path, self.word_save_path, path, "Word")

    def browse_notepad_source_path(self):
        path = filedialog.askopenfilename(
            title="변환할 TXT 파일 선택",
            filetypes=[("텍스트", "*.txt"), ("모든 파일", "*.*")],
        )
        preferred_ext = ".docx" if self.notepad_save_format.get() == "docx" else ".txt"
        self._apply_source_file_selection("text", self.notepad_source_path, self.notepad_save_path, path, "TXT", preferred_ext)

    def _prepare_direct_file_conversion(self, kind, source_path, save_path):
        if not source_path:
            return None
        source_path = os.path.abspath(source_path)
        if not os.path.isfile(source_path):
            raise Exception(f"선택한 파일이 없습니다: {source_path}")
        detected_kind = self._batch_file_kind(source_path)
        if detected_kind != kind:
            raise Exception(f"선택한 파일 형식이 맞지 않습니다: {source_path}")
        if not save_path:
            save_path = self._default_direct_save_path(source_path, kind)
        save_path = os.path.abspath(save_path)
        if source_path.lower() == save_path.lower():
            raise Exception("원본 파일과 같은 경로로 저장할 수 없습니다.")
        return source_path, save_path

    def _read_text_file_for_conversion(self, source_path):
        last_error = None
        for encoding in ("utf-8-sig", "cp949", "utf-16", "utf-8"):
            try:
                with open(source_path, "r", encoding=encoding) as source_file:
                    return source_file.read()
            except Exception as error:
                last_error = error
        raise Exception(f"텍스트 파일 인코딩을 읽을 수 없습니다: {last_error}")

    def _convert_text_source_file(self, source_path, target_path):
        target_ext = os.path.splitext(target_path)[1].lower()
        target_dir = os.path.dirname(os.path.abspath(target_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        if target_ext == ".docx":
            if not HAS_DOCX:
                raise Exception("DOCX 저장에는 python-docx 패키지가 필요합니다.")
            text = self._read_text_file_for_conversion(source_path)
            doc = DocxDocument()
            for line in text.splitlines():
                doc.add_paragraph(self._clean_xml_text(line))
            doc.save(target_path)
            self._validate_office_openxml(target_path, "TXT 파일 DOCX 변환")
            return
        shutil.copy2(source_path, target_path)
        if not os.path.exists(target_path) or os.path.getsize(target_path) <= 0:
            raise Exception("TXT 복사 결과 파일이 없거나 비어 있습니다.")

    def _convert_direct_file(self, kind, source_path, save_path):
        if kind == "text":
            self._convert_text_source_file(source_path, save_path)
            return
        if kind == "pdf":
            self._convert_pdf_file(source_path, save_path)
            return
        if kind == "hwp":
            self._convert_hwp_file(source_path, save_path)
            return
        if kind in {"ppt", "excel", "word"}:
            label = {"ppt": "PPT", "excel": "Excel", "word": "Word"}[kind]
            try:
                self._try_existing_office_file_copy(source_path, save_path, label)
                return
            except Exception as direct_copy_error:
                self.logger.log(
                    f"{label} 직접 파일 복사 불가, Office 내부 복원 시도: {str(direct_copy_error)[:120]}"
                )

        if not HAS_WIN32COM:
            raise Exception("Office 파일 직접 변환에는 pywin32/win32com이 필요합니다.")

        # 일괄 변환과 동일한 예열·재사용 인스턴스를 사용한다.
        # (이 메서드는 Office 워커 스레드에서 실행되어 COM 아파트먼트가 일치한다.)
        if kind == "ppt":
            app = self._acquire_warm_office_app("ppt", "PowerPoint.Application", "PowerPoint")
            self._batch_convert_ppt_file(app, source_path, save_path, skip_direct=True)
        elif kind == "excel":
            app = self._acquire_warm_office_app("excel", "Excel.Application", "Excel")
            self._batch_convert_excel_file(app, source_path, save_path, skip_direct=True)
        elif kind == "word":
            app = self._acquire_warm_office_app("word", "Word.Application", "Word")
            self._batch_convert_word_file(app, source_path, save_path)
        else:
            raise Exception(f"지원하지 않는 직접 변환 형식입니다: {kind}")

    def _start_direct_file_conversion(self, kind, source_path, save_path, save_var, button, label):
        try:
            source_path, save_path = self._prepare_direct_file_conversion(kind, source_path, save_path)
        except Exception as error:
            messagebox.showwarning("경고", str(error))
            return True

        save_var.set(save_path)
        button.config(state=tk.DISABLED)
        self.progress_var.set(0)
        # 일괄 변환과 동일한 예열·재사용 워커에서 처리해 Office 콜드 스타트를 없앤다.
        self._ensure_office_worker()
        self._office_job_queue.put(
            lambda: self._extract_direct_file(kind, source_path, save_path, button, label)
        )
        return True

    def _extract_direct_file(self, kind, source_path, save_path, button, label):
        self.logger.log(f"=== {label} 파일 직접 변환 시작 ===")
        extract_start = time.perf_counter()
        try:
            self.root.after(0, lambda: self.status_text.set(f"{label} 파일 변환 중..."))
            self.root.after(0, lambda: self.progress_var.set(10))
            self._convert_direct_file(kind, source_path, save_path)
            self._log_elapsed(f"{label} 파일 직접 변환 시간", extract_start)
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set(f"{label} 파일 변환 완료!"))
            self.root.after(0, lambda: messagebox.showinfo("완료", f"{label} 파일 변환 완료!\n{save_path}"))
        except Exception as error:
            message = str(error)
            self.logger.error(f"{label} 파일 직접 변환 오류", error)
            self.root.after(0, lambda: self.status_text.set(f"오류: {message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"파일 변환 중 오류:\n{message}"))
        finally:
            self.root.after(0, lambda: button.config(state=tk.NORMAL))

    # ========== PPT 관련 메서드 ==========

    def browse_ppt_save_path(self):
        """PPT 저장 경로 선택"""
        self.logger.log("PPT 저장 경로 선택 대화상자 열기")

        direct_source = self.ppt_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("ppt"):
            doc_name = os.path.basename(direct_source)
        else:
            doc_name = self.ppt_doc_name.get()
        if doc_name and doc_name != "감지 중..." and doc_name != "열린 PPT 없음":
            src_ext = os.path.splitext(doc_name)[1] or ".pptx"
            default_ext = src_ext if src_ext.lower() in [".pptx", ".ppt", ".pptm"] else ".pptx"
            default_name = os.path.splitext(doc_name)[0] + "_복사본" + default_ext
        else:
            default_ext = ".pptx"
            default_name = "새문서.pptx"

        path = filedialog.asksaveasfilename(
            defaultextension=default_ext,
            filetypes=[("PowerPoint 파일", "*.pptx *.pptm *.ppt"), ("모든 파일", "*.*")],
            initialfile=default_name,
            title="저장할 위치 선택"
        )
        if path:
            self.ppt_save_path.set(path)
            self.logger.log(f"PPT 저장 경로 선택됨: {path}")

    def detect_open_ppt(self, prefer_open=False):
        """열려있는 PPT 감지"""
        if prefer_open:
            self._use_open_document_input("ppt", self.ppt_source_path, "PPT")
        self.logger.log("PPT 감지 시작")
        if self._is_direct_file_input_active("ppt"):
            self.status_text.set("PPT 감지 중... (파일 선택 유지)")
            self._show_direct_file_input("ppt", self.ppt_source_path.get().strip(), "PPT")
        else:
            self.status_text.set("PPT 감지 중...")
            self.ppt_doc_name.set("감지 중...")
            self.ppt_slide_count.set("-")

        thread = threading.Thread(target=self._detect_ppt)
        thread.daemon = True
        thread.start()

    def _detect_ppt(self):
        """PPT 감지 (백그라운드)"""
        self.logger.log("백그라운드 PPT 감지 스레드 시작")
        pythoncom.CoInitialize()

        try:
            ppt, _ = self._get_ppt_app(allow_dispatch=False)
            ppt_count = ppt.Presentations.Count
            self.logger.log(f"PowerPoint 연결 성공, 열린 프레젠테이션 수: {ppt_count}")

            if ppt_count > 0:
                ppt_names = []
                ppt_info = []

                for i in range(1, ppt_count + 1):
                    try:
                        presentation = ppt.Presentations(i)
                        name = presentation.Name
                        slide_count = presentation.Slides.Count
                        ppt_names.append(f"{name} ({slide_count}장)")
                        ppt_info.append((name, slide_count, i))
                        self.logger.log(f"  PPT {i}: {name}, {slide_count}장")
                    except Exception as e:
                        self.logger.log(f"  PPT {i} 정보 가져오기 실패: {str(e)}")

                self.ppt_list = ppt_info

                def update_combo():
                    self.ppt_combo['values'] = ppt_names
                    if self._is_direct_file_input_active("ppt"):
                        self._show_direct_file_input("ppt", self.ppt_source_path.get().strip(), "PPT")
                        self.status_text.set(f"PPT {len(ppt_names)}개 감지됨 (파일 선택 유지)")
                        return
                    if ppt_names:
                        self.ppt_combo.current(0)
                        self.selected_ppt_index.set(1)
                        self.ppt_doc_name.set(ppt_info[0][0])
                        self.ppt_slide_count.set(f"{ppt_info[0][1]}장")
                    self.status_text.set(f"PPT {len(ppt_names)}개 감지됨")

                self.root.after(0, update_combo)
            else:
                self.logger.log("열린 프레젠테이션 없음")
                self.ppt_list = []
                def clear_combo():
                    self.ppt_combo.set("")
                    self.ppt_combo['values'] = []
                    if self._is_direct_file_input_active("ppt"):
                        self._show_direct_file_input("ppt", self.ppt_source_path.get().strip(), "PPT")
                        self.status_text.set("PPT 파일 선택됨")
                        return
                    self.ppt_doc_name.set("열린 PPT 없음")
                    self.ppt_slide_count.set("-")
                    self.status_text.set("PPT를 먼저 열어주세요")
                self.root.after(0, clear_combo)

        except Exception as e:
            expected_not_running = self._is_expected_app_not_running(e, "PowerPoint")
            if expected_not_running:
                self.logger.log(f"PPT 감지: PowerPoint가 아직 실행 중이 아닙니다. ({str(e)[:80]})")
            else:
                self.logger.error("PPT 감지 실패", e)
            self.ppt_list = []
            err_msg = str(e)[:30]
            def show_error():
                self.ppt_combo.set("")
                if self._is_direct_file_input_active("ppt"):
                    self._show_direct_file_input("ppt", self.ppt_source_path.get().strip(), "PPT")
                    self.status_text.set("PPT 파일 선택됨")
                    return
                self.ppt_doc_name.set("열린 PPT 없음")
                self.ppt_slide_count.set("-")
                if expected_not_running:
                    self.status_text.set("PPT를 먼저 열어주세요")
                else:
                    self.status_text.set(f"PPT 감지 실패: {err_msg}")
            self.root.after(0, show_error)

        pythoncom.CoUninitialize()

    def on_ppt_selected(self, event):
        """PPT 콤보박스 선택 이벤트"""
        selected_idx = self.ppt_combo.current()
        if selected_idx >= 0 and selected_idx < len(self.ppt_list):
            self._use_open_document_input("ppt", self.ppt_source_path, "PPT")
            name, slide_count, ppt_index = self.ppt_list[selected_idx]
            self.selected_ppt_index.set(ppt_index)
            self.ppt_doc_name.set(name)
            self.ppt_slide_count.set(f"{slide_count}장")
            self.logger.log(f"PPT 선택: {name} (인덱스 {ppt_index})")

    def start_ppt_extraction(self):
        """PPT 추출 시작"""
        self.logger.log("PPT 추출 시작 버튼 클릭")

        save_path = self.ppt_save_path.get()
        mode = self.ppt_extract_mode.get()
        ppt_index = self.selected_ppt_index.get()
        self.logger.log(f"PPT 추출 설정: mode={mode}, index={ppt_index}, save_path={save_path}")

        direct_source = self.ppt_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("ppt"):
            if self._start_direct_file_conversion(
                "ppt", direct_source, save_path, self.ppt_save_path, self.ppt_extract_button, "PPT"
            ):
                return

        if mode == "image_only":
            messagebox.showwarning(
                "이미지 변환 비활성화",
                "텍스트를 이미지로 변환하는 방식은 사용할 수 없습니다.\n"
                "원본 그대로 복사 또는 텍스트 중심 + 객체 보존을 선택해주세요."
            )
            return

        if not save_path:
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        if self.ppt_doc_name.get() == "열린 PPT 없음" or not self.ppt_list:
            messagebox.showwarning("경고", "열린 PPT가 없습니다.")
            return

        self.ppt_extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(target=self._extract_ppt, args=(save_path, mode, ppt_index))
        thread.daemon = True
        thread.start()

    def _extract_ppt(self, save_path, mode, ppt_index):
        """PPT 추출 (백그라운드)"""
        self.logger.log("=== PPT 추출 프로세스 시작 ===")
        extract_start = time.perf_counter()
        pythoncom.CoInitialize()

        temp_dir = None
        ppt_app = None
        original_alerts = None

        try:
            self.root.after(0, lambda: self.status_text.set("원본 PPT 연결 중..."))

            self.logger.log("PowerPoint COM 연결 시도")
            ppt_app, _ = self._get_ppt_app()
            try:
                original_alerts = ppt_app.DisplayAlerts
                ppt_app.DisplayAlerts = 1  # ppAlertsNone
                self.logger.log("PowerPoint 경고창 표시 비활성화")
            except Exception as alerts_error:
                self.logger.log(f"PowerPoint 경고창 설정 실패: {str(alerts_error)[:60]}")
            ppt_count = ppt_app.Presentations.Count
            if ppt_count == 0:
                raise Exception("열린 PowerPoint 문서가 없습니다. PowerPoint에서 문서를 먼저 열어주세요.")

            if ppt_index > 0 and ppt_index <= ppt_count:
                source_pres = ppt_app.Presentations(ppt_index)
            else:
                source_pres = ppt_app.ActivePresentation

            self.logger.log(f"원본 프레젠테이션: {source_pres.Name}")

            if mode == "native_copy":
                try:
                    self.root.after(0, lambda: self.status_text.set("원본 PPT 그대로 복사 중..."))
                    self.root.after(0, lambda: self.progress_var.set(20))
                    total_slides = source_pres.Slides.Count
                    self._save_native_copy(source_pres, save_path, "PPT")
                    self._log_elapsed("PPT 원본 복사 시간", extract_start)
                    self.root.after(0, lambda: self.progress_var.set(100))
                    self.root.after(0, lambda: self.status_text.set("PPT 원본 복사 완료!"))
                    self.root.after(0, lambda: messagebox.showinfo("완료",
                        f"PPT 원본 복사 완료!\n{save_path}\n\n총 {total_slides}장"))
                    return
                except Exception as copy_error:
                    self.logger.log(
                        f"PPT 원본 복사 결과 검증 실패, 원본 구조 복원 경로로 전환: {str(copy_error)[:120]}"
                    )
                    package_error_detail = ""
                    try:
                        self.root.after(0, lambda: self.status_text.set("원본 복사 실패, 클립보드 슬라이드 패키지 복원 중..."))
                        self.root.after(0, lambda: self.progress_var.set(25))
                        self._save_ppt_clipboard_package_copy(source_pres, save_path)
                        self._log_elapsed("PPT 클립보드 슬라이드 패키지 복원 시간", extract_start)
                        self.root.after(0, lambda: self.progress_var.set(100))
                        self.root.after(0, lambda: self.status_text.set("PPT 원본 구조 복원 완료!"))
                        self.root.after(0, lambda: messagebox.showinfo("완료",
                            f"PPT 원본 구조 복원 완료!\n{save_path}\n\n총 {total_slides}장"))
                        return
                    except Exception as package_error:
                        package_error_detail = str(package_error)[:200]
                        self.logger.log(
                            f"PPT 클립보드 슬라이드 패키지 복원 실패, 슬라이드 복제로 전환: {package_error_detail[:120]}"
                        )

                    try:
                        self.root.after(0, lambda: self.status_text.set("원본 복사 실패, PowerPoint 슬라이드 복제 중..."))
                        self.root.after(0, lambda: self.progress_var.set(30))
                        self._save_ppt_slide_clone(source_pres, save_path)
                        self._log_elapsed("PPT 슬라이드 복제 시간", extract_start)
                        self.root.after(0, lambda: self.progress_var.set(100))
                        self.root.after(0, lambda: self.status_text.set("PPT 슬라이드 복제 완료!"))
                        self.root.after(0, lambda: messagebox.showinfo("완료",
                            f"PPT 슬라이드 복제 완료!\n{save_path}\n\n총 {total_slides}장"))
                        return
                    except Exception as clone_error:
                        self.logger.log(
                            f"PPT 슬라이드 복제 실패, 이미지/하이브리드 자동 전환 안 함: {str(clone_error)[:120]}"
                        )
                        raise Exception(
                            "편집 가능한 PPT 원본 구조 복원에 실패했습니다. 이미지 변환은 비활성화되어 있습니다.\n"
                            f"- 원본 복사 실패: {str(copy_error)[:200]}\n"
                            f"- 클립보드 슬라이드 패키지 실패: {package_error_detail or '확인 불가'}\n"
                            f"- PowerPoint 슬라이드 복제 실패: {str(clone_error)[:200]}\n\n"
                            "원본 PPT를 PowerPoint에 연 상태로 두고 다시 시도하거나, DRM/보안 정책을 해제한 뒤 저장해야 합니다."
                        ) from clone_error

            if not HAS_PPTX:
                raise Exception("python-pptx 패키지가 필요합니다. pip install python-pptx")

            self.root.after(0, lambda: self.status_text.set("새 PPT 생성 중..."))
            self.root.after(0, lambda: self.progress_var.set(5))

            new_pres = Presentation()

            slide_width = source_pres.PageSetup.SlideWidth
            slide_height = source_pres.PageSetup.SlideHeight
            new_pres.slide_width = Emu(int(slide_width * 12700))
            new_pres.slide_height = Emu(int(slide_height * 12700))

            total_slides = source_pres.Slides.Count
            temp_dir = tempfile.mkdtemp()
            blank_layout = new_pres.slide_layouts[6]

            for i in range(1, total_slides + 1):
                slide_start = time.perf_counter()
                self.logger.log(f"--- 슬라이드 {i}/{total_slides} 처리 ---")
                progress = 5 + (i / total_slides) * 85
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                self.root.after(0, lambda n=i, t=total_slides: self.status_text.set(f"슬라이드 {n}/{t} 처리 중..."))

                source_slide = source_pres.Slides(i)
                new_slide = new_pres.slides.add_slide(blank_layout)

                # 슬라이드 배경색 복사
                self._copy_slide_background(source_slide, new_slide)

                if mode == "image_only":
                    self._export_slide_as_image(source_slide, new_slide, temp_dir, i, new_pres)
                elif mode == "text_only":
                    self._extract_text_with_object_images(source_slide, new_slide, temp_dir)
                else:
                    self._extract_hybrid(source_slide, new_slide, temp_dir, i, new_pres)
                self._log_elapsed(f"슬라이드 {i} 처리 시간", slide_start)

            self.root.after(0, lambda: self.status_text.set("파일 저장 중..."))
            self.root.after(0, lambda: self.progress_var.set(95))

            temp_ppt_path = self._make_local_temp_path(".pptx")
            try:
                new_pres.save(temp_ppt_path)
                self._publish_verified_file(temp_ppt_path, save_path, "PPT 재구성")
            finally:
                if os.path.exists(temp_ppt_path):
                    try:
                        os.remove(temp_ppt_path)
                    except Exception:
                        pass
            self.logger.log(f"저장 완료: {save_path}")
            self._log_elapsed("PPT 전체 추출 시간", extract_start)

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("PPT 추출 완료!"))
            self.root.after(0, lambda: messagebox.showinfo("완료",
                f"PPT 추출 완료!\n{save_path}\n\n총 {total_slides}장"))

        except Exception as e:
            error_message = str(e)
            self.logger.error("PPT 추출 오류", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"추출 중 오류:\n{error_message}"))

        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
            if ppt_app is not None and original_alerts is not None:
                try:
                    ppt_app.DisplayAlerts = original_alerts
                except Exception:
                    pass
            self.root.after(0, lambda: self.ppt_extract_button.config(state=tk.NORMAL))
            pythoncom.CoUninitialize()

    # ========== Excel 관련 메서드 ==========

    def browse_excel_save_path(self):
        """Excel 저장 경로 선택"""
        self.logger.log("Excel 저장 경로 선택")

        direct_source = self.excel_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("excel"):
            doc_name = os.path.basename(direct_source)
        else:
            doc_name = self.excel_doc_name.get()
        if doc_name and doc_name != "감지 중..." and doc_name != "열린 Excel 없음":
            src_ext = os.path.splitext(doc_name)[1] or ".xlsx"
            default_ext = src_ext if src_ext.lower() in [".xlsx", ".xlsm", ".xls", ".xlsb"] else ".xlsx"
            default_name = os.path.splitext(doc_name)[0] + "_복사본" + default_ext
        else:
            default_ext = ".xlsx"
            default_name = "새문서.xlsx"

        path = filedialog.asksaveasfilename(
            defaultextension=default_ext,
            filetypes=[("Excel 파일", "*.xlsx *.xlsm *.xls *.xlsb"), ("모든 파일", "*.*")],
            initialfile=default_name,
            title="저장할 위치 선택"
        )
        if path:
            self.excel_save_path.set(path)
            self.logger.log(f"Excel 저장 경로: {path}")

    def detect_open_excel(self, prefer_open=False):
        """열려있는 Excel 감지"""
        if prefer_open:
            self._use_open_document_input("excel", self.excel_source_path, "Excel")
        self.logger.log("Excel 감지 시작")
        if self._is_direct_file_input_active("excel"):
            self.status_text.set("Excel 감지 중... (파일 선택 유지)")
            self._show_direct_file_input("excel", self.excel_source_path.get().strip(), "Excel")
        else:
            self.status_text.set("Excel 감지 중...")
            self.excel_doc_name.set("감지 중...")
            self.excel_sheet_count.set("-")

        thread = threading.Thread(target=self._detect_excel)
        thread.daemon = True
        thread.start()

    def _detect_excel(self):
        """Excel 감지 (백그라운드)"""
        pythoncom.CoInitialize()

        try:
            excel, _ = self._get_excel_app(allow_dispatch=False)
            wb_count = excel.Workbooks.Count
            self.logger.log(f"Excel 연결 성공, 열린 통합문서 수: {wb_count}")

            if wb_count > 0:
                excel_names = []
                excel_info = []

                for i in range(1, wb_count + 1):
                    try:
                        workbook = excel.Workbooks(i)
                        name = workbook.Name
                        sheet_count = workbook.Sheets.Count
                        excel_names.append(f"{name} ({sheet_count}시트)")
                        excel_info.append((name, sheet_count, i))
                        self.logger.log(f"  Excel {i}: {name}, {sheet_count}시트")
                    except Exception as e:
                        self.logger.log(f"  Excel {i} 정보 실패: {str(e)}")

                self.excel_list = excel_info

                def update_combo():
                    self.excel_combo['values'] = excel_names
                    if self._is_direct_file_input_active("excel"):
                        self._show_direct_file_input("excel", self.excel_source_path.get().strip(), "Excel")
                        self.status_text.set(f"Excel {len(excel_names)}개 감지됨 (파일 선택 유지)")
                        return
                    if excel_names:
                        self.excel_combo.current(0)
                        self.selected_excel_index.set(1)
                        self.excel_doc_name.set(excel_info[0][0])
                        self.excel_sheet_count.set(f"{excel_info[0][1]}시트")
                    self.status_text.set(f"Excel {len(excel_names)}개 감지됨")

                self.root.after(0, update_combo)
            else:
                self.excel_list = []
                def clear_combo():
                    self.excel_combo.set("")
                    self.excel_combo['values'] = []
                    if self._is_direct_file_input_active("excel"):
                        self._show_direct_file_input("excel", self.excel_source_path.get().strip(), "Excel")
                        self.status_text.set("Excel 파일 선택됨")
                        return
                    self.excel_doc_name.set("열린 Excel 없음")
                    self.excel_sheet_count.set("-")
                    self.status_text.set("Excel을 먼저 열어주세요")
                self.root.after(0, clear_combo)

        except Exception as e:
            expected_not_running = self._is_expected_app_not_running(e, "Excel")
            if expected_not_running:
                self.logger.log(f"Excel 감지: Excel이 아직 실행 중이 아닙니다. ({str(e)[:80]})")
            else:
                self.logger.error("Excel 감지 실패", e)
            self.excel_list = []
            def show_error():
                self.excel_combo.set("")
                if self._is_direct_file_input_active("excel"):
                    self._show_direct_file_input("excel", self.excel_source_path.get().strip(), "Excel")
                    self.status_text.set("Excel 파일 선택됨")
                    return
                self.excel_doc_name.set("열린 Excel 없음")
                self.excel_sheet_count.set("-")
                if expected_not_running:
                    self.status_text.set("Excel을 먼저 열어주세요")
                else:
                    self.status_text.set("Excel 감지 실패")
            self.root.after(0, show_error)

        pythoncom.CoUninitialize()

    def on_excel_selected(self, event):
        """Excel 콤보박스 선택 이벤트"""
        selected_idx = self.excel_combo.current()
        if selected_idx >= 0 and selected_idx < len(self.excel_list):
            self._use_open_document_input("excel", self.excel_source_path, "Excel")
            name, sheet_count, excel_index = self.excel_list[selected_idx]
            self.selected_excel_index.set(excel_index)
            self.excel_doc_name.set(name)
            self.excel_sheet_count.set(f"{sheet_count}시트")
            self.logger.log(f"Excel 선택: {name}")

    def start_excel_extraction(self):
        """Excel 추출 시작"""
        self.logger.log("Excel 추출 시작")

        native_copy = self.excel_native_copy.get()

        if not native_copy and not HAS_OPENPYXL:
            messagebox.showerror("오류", "openpyxl 패키지가 필요합니다.\npip install openpyxl")
            return

        save_path = self.excel_save_path.get()
        include_format = self.excel_include_format.get()
        values_only = self.excel_include_formulas.get()
        excel_index = self.selected_excel_index.get()

        direct_source = self.excel_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("excel"):
            if self._start_direct_file_conversion(
                "excel", direct_source, save_path, self.excel_save_path, self.excel_extract_button, "Excel"
            ):
                return

        if not save_path:
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        if self.excel_doc_name.get() == "열린 Excel 없음" or not self.excel_list:
            messagebox.showwarning("경고", "열린 Excel이 없습니다.")
            return

        self.excel_extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(
            target=self._extract_excel,
            args=(save_path, include_format, values_only, excel_index, native_copy),
        )
        thread.daemon = True
        thread.start()

    def _extract_excel(self, save_path, include_format, values_only, excel_index, native_copy=False):
        """Excel 추출 (백그라운드)"""
        self.logger.log("=== Excel 추출 시작 ===")
        extract_start = time.perf_counter()
        pythoncom.CoInitialize()

        temp_dir = None

        try:
            self.root.after(0, lambda: self.status_text.set("Excel 연결 중..."))

            excel_app, _ = self._get_excel_app()
            wb_count = excel_app.Workbooks.Count
            if wb_count == 0:
                raise Exception("열린 Excel 문서가 없습니다. Excel에서 문서를 먼저 열어주세요.")

            if excel_index > 0 and excel_index <= wb_count:
                source_wb = excel_app.Workbooks(excel_index)
            else:
                source_wb = excel_app.ActiveWorkbook

            self.logger.log(f"원본 통합문서: {source_wb.Name}")

            if native_copy:
                try:
                    self.root.after(0, lambda: self.status_text.set("원본 Excel 그대로 복사 중..."))
                    self.root.after(0, lambda: self.progress_var.set(20))
                    total_sheets = source_wb.Sheets.Count
                    self._save_native_copy(source_wb, save_path, "Excel")
                    self._log_elapsed("Excel 원본 복사 시간", extract_start)
                    self.root.after(0, lambda: self.progress_var.set(100))
                    self.root.after(0, lambda: self.status_text.set("Excel 원본 복사 완료!"))
                    self.root.after(0, lambda: messagebox.showinfo("완료",
                        f"Excel 원본 복사 완료!\n{save_path}\n\n총 {total_sheets}시트"))
                    return
                except Exception as copy_error:
                    self.logger.log(
                        f"Excel 원본 복사 결과 검증 실패, 재구성으로 전환: {str(copy_error)[:120]}"
                    )
                    self.root.after(0, lambda: self.status_text.set("원본 복사 실패, Excel 재구성 중..."))

            if not HAS_OPENPYXL:
                raise Exception("openpyxl 패키지가 필요합니다. pip install openpyxl")

            # openpyxl로 새 통합문서 생성
            new_wb = Workbook()
            # 기본 시트 제거 (나중에 추가할 것이므로)
            default_sheet = new_wb.active

            total_sheets = source_wb.Sheets.Count
            temp_dir = tempfile.mkdtemp()
            rebuild_issues = []
            self.root.after(0, lambda: self.progress_var.set(5))

            for sheet_idx in range(1, total_sheets + 1):
                sheet_start = time.perf_counter()
                source_sheet = source_wb.Sheets(sheet_idx)
                sheet_name = source_sheet.Name
                self.logger.log(f"  시트 처리: {sheet_name}")

                progress = 5 + (sheet_idx / total_sheets) * 85
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                self.root.after(0, lambda n=sheet_name: self.status_text.set(f"시트 '{n}' 처리 중..."))

                # 새 시트 생성
                if sheet_idx == 1:
                    new_sheet = default_sheet
                    new_sheet.title = sheet_name
                else:
                    new_sheet = new_wb.create_sheet(title=sheet_name)

                copied_objects, visible_objects = self._copy_excel_sheet_objects(
                    source_sheet, new_sheet, temp_dir, sheet_name
                )
                if visible_objects and copied_objects < visible_objects:
                    rebuild_issues.append(
                        f"{sheet_name}: 삽입 객체 {copied_objects}/{visible_objects}개만 복사됨"
                    )

                # 사용 범위 가져오기
                try:
                    source_range, start_row, start_col, row_count, col_count = self._get_excel_effective_range(
                        source_sheet, sheet_name
                    )
                    if source_range is None:
                        continue

                    cell_count = row_count * col_count

                    self.logger.log(f"    범위: {row_count}행 x {col_count}열 ({cell_count:,}셀, 시작: {start_row},{start_col})")

                    if cell_count > self.EXCEL_VALUE_CELL_LIMIT:
                        raise Exception(
                            f"사용 범위가 너무 큽니다 ({cell_count:,}셀). "
                            "Excel에서 불필요한 빈 행/열 서식을 지운 뒤 다시 시도해주세요."
                        )

                    sheet_include_format = include_format
                    if sheet_include_format and cell_count > self.EXCEL_FORMAT_CELL_LIMIT:
                        sheet_include_format = False
                        self.logger.log(
                            f"    대용량 시트 서식 복사 자동 생략: {cell_count:,}셀 "
                            f"(한도 {self.EXCEL_FORMAT_CELL_LIMIT:,}셀)"
                        )
                        self.root.after(
                            0,
                            lambda n=sheet_name: self.status_text.set(
                                f"시트 '{n}' 값 복사 중... (대용량 서식 생략)"
                            ),
                        )

                    # 값/수식은 COM Range에서 한 번에 읽어 셀 단위 왕복을 줄인다.
                    try:
                        range_data = source_range.Value if values_only else source_range.Formula
                        data_rows = self._excel_range_to_rows(range_data, row_count, col_count)
                        for r, row_values in enumerate(data_rows):
                            for c, value in enumerate(row_values[:col_count]):
                                if value is not None:
                                    new_sheet.cell(row=start_row + r, column=start_col + c).value = value
                    except Exception as data_err:
                        self.logger.log(f"    범위 데이터 일괄 읽기 실패, 셀 단위로 폴백: {str(data_err)[:50]}")
                        data_rows = None

                    if data_rows is None or sheet_include_format:
                        for r in range(row_count):
                            if r % 100 == 0:
                                progress = 5 + ((sheet_idx - 1) / total_sheets) * 85
                                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                                self.root.after(
                                    0,
                                    lambda n=sheet_name, rr=r + 1, total=row_count:
                                        self.status_text.set(f"시트 '{n}' 행 {rr}/{total} 처리 중...")
                                )
                            for c in range(col_count):
                                try:
                                    cell = source_sheet.Cells(start_row + r, start_col + c)
                                    new_cell = new_sheet.cell(row=start_row + r, column=start_col + c)

                                    # 일괄 읽기에 실패한 경우에만 값/수식을 셀 단위로 읽는다.
                                    if data_rows is None:
                                        if values_only:
                                            new_cell.value = cell.Value
                                        else:
                                            try:
                                                formula = cell.Formula
                                                if formula and str(formula).startswith('='):
                                                    new_cell.value = formula
                                                else:
                                                    new_cell.value = cell.Value
                                            except Exception:
                                                new_cell.value = cell.Value

                                    # 서식 복사
                                    if sheet_include_format:
                                        try:
                                            # 폰트
                                            src_font = cell.Font
                                            new_cell.font = Font(
                                                name=src_font.Name if src_font.Name else 'Calibri',
                                                size=src_font.Size if src_font.Size else 11,
                                                bold=bool(src_font.Bold),
                                                italic=bool(src_font.Italic),
                                                color=self._excel_color_to_hex(src_font.Color)
                                            )
                                        except Exception as font_err:
                                            self.logger.log(f"    셀 서식(폰트) 복사 실패: {str(font_err)[:40]}")

                                        try:
                                            # 배경색
                                            interior_color = cell.Interior.Color
                                            if interior_color and interior_color != 16777215:  # 흰색이 아니면
                                                hex_color = self._excel_color_to_hex(interior_color)
                                                if hex_color:
                                                    new_cell.fill = PatternFill(start_color=hex_color,
                                                                                end_color=hex_color,
                                                                                fill_type='solid')
                                        except Exception as fill_err:
                                            self.logger.log(f"    셀 서식(배경색) 복사 실패: {str(fill_err)[:40]}")

                                        try:
                                            number_format = cell.NumberFormat
                                            if number_format:
                                                new_cell.number_format = number_format
                                        except Exception:
                                            pass

                                except Exception as cell_err:
                                    self.logger.log(f"    셀 처리 실패: {str(cell_err)[:40]}")

                    if row_count <= self.EXCEL_ROW_HEIGHT_COPY_LIMIT:
                        try:
                            for r in range(1, row_count + 1):
                                height = source_sheet.Rows(start_row + r - 1).RowHeight
                                if height:
                                    new_sheet.row_dimensions[start_row + r - 1].height = height
                        except Exception:
                            pass
                    else:
                        self.logger.log(f"    행 높이 복사 생략: {row_count:,}행")

                    if col_count <= self.EXCEL_COLUMN_WIDTH_COPY_LIMIT:
                        try:
                            for c in range(1, col_count + 1):
                                width = source_sheet.Columns(start_col + c - 1).ColumnWidth
                                new_sheet.column_dimensions[get_column_letter(start_col + c - 1)].width = width
                        except:
                            pass
                    else:
                        self.logger.log(f"    열 너비 복사 생략: {col_count:,}열")

                    self._copy_excel_merged_cells(
                        source_sheet,
                        new_sheet,
                        start_row,
                        start_col,
                        start_row + row_count - 1,
                        start_col + col_count - 1,
                    )

                except Exception as sheet_err:
                    self.logger.error(f"시트 '{sheet_name}' 처리 오류", sheet_err)
                    rebuild_issues.append(f"{sheet_name}: {str(sheet_err)[:160]}")
                finally:
                    self._log_elapsed(f"시트 '{sheet_name}' 처리 시간", sheet_start)

            # 저장
            self.root.after(0, lambda: self.status_text.set("파일 저장 중..."))
            self.root.after(0, lambda: self.progress_var.set(95))

            new_wb.save(save_path)
            self._validate_office_openxml(save_path, "Excel 재구성")
            if rebuild_issues:
                issue_text = self._format_excel_rebuild_issues(rebuild_issues)
                self.logger.log(f"저장 완료(부분 완료/확인 필요): {save_path}")
                self.logger.log(f"Excel 재구성 확인 필요 항목 {len(rebuild_issues)}개:\n{issue_text}")
            else:
                self.logger.log(f"저장 완료: {save_path}")
            self._log_elapsed("Excel 전체 추출 시간", extract_start)

            self.root.after(0, lambda: self.progress_var.set(100))
            if rebuild_issues:
                issue_text = self._format_excel_rebuild_issues(rebuild_issues)
                self.root.after(0, lambda: self.status_text.set("Excel 부분 완료 - 확인 필요"))
                self.root.after(0, lambda: messagebox.showwarning("부분 완료",
                    f"Excel 재구성 파일을 저장했지만 일부 항목 확인이 필요합니다.\n"
                    f"{save_path}\n\n"
                    f"확인 필요:\n{issue_text}"))
            else:
                self.root.after(0, lambda: self.status_text.set("Excel 추출 완료!"))
                self.root.after(0, lambda: messagebox.showinfo("완료",
                    f"Excel 추출 완료!\n{save_path}\n\n총 {total_sheets}시트"))

        except Exception as e:
            error_message = str(e)
            self.logger.error("Excel 추출 오류", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"추출 중 오류:\n{error_message}"))

        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
            self.root.after(0, lambda: self.excel_extract_button.config(state=tk.NORMAL))
            pythoncom.CoUninitialize()

    def _format_excel_rebuild_issues(self, issues, limit=8):
        shown = list(issues[:limit])
        lines = [f"- {issue}" for issue in shown]
        remaining = len(issues) - len(shown)
        if remaining > 0:
            lines.append(f"- 외 {remaining}개")
        return "\n".join(lines)

    def _excel_find_cell(self, source_sheet, search_order, search_direction, look_ins=None):
        """서식만 묻은 UsedRange 대신 실제 값/수식 셀을 찾는다."""
        cells = source_sheet.Cells
        if look_ins is None:
            look_ins = (-4163, -4123)  # xlValues, xlFormulas
        try:
            source_sheet.Application.FindFormat.Clear()
        except Exception:
            pass
        try:
            if search_direction == 2:  # xlPrevious
                after_cell = cells(1, 1)
            else:
                after_cell = cells(source_sheet.Rows.Count, source_sheet.Columns.Count)
        except Exception:
            after_cell = None

        for look_in in look_ins:
            try:
                find_kwargs = {
                    "What": "*",
                    "LookIn": look_in,
                    "LookAt": 2,  # xlPart
                    "SearchOrder": search_order,
                    "SearchDirection": search_direction,
                    "MatchCase": False,
                    "SearchFormat": False,
                }
                if after_cell is not None:
                    find_kwargs["After"] = after_cell
                cell = cells.Find(**find_kwargs)
                if cell is not None:
                    return cell
            except Exception:
                pass

            try:
                cell = cells.Find("*", after_cell, look_in, 2, search_order, search_direction, False, False, False)
                if cell is not None:
                    return cell
            except Exception:
                pass

        return None

    def _excel_bounds_cell_count(self, bounds):
        if bounds is None:
            return 0
        return max(1, bounds[2] - bounds[0] + 1) * max(1, bounds[3] - bounds[1] + 1)

    def _excel_content_bounds_by_look_in(self, source_sheet, look_in):
        """값 또는 수식 기준으로 실제 셀 범위를 반환한다."""
        try:
            look_ins = (look_in,)
            last_row_cell = self._excel_find_cell(source_sheet, 1, 2, look_ins)  # xlByRows, xlPrevious
            last_col_cell = self._excel_find_cell(source_sheet, 2, 2, look_ins)  # xlByColumns, xlPrevious
            if last_row_cell is None or last_col_cell is None:
                return None

            first_row_cell = self._excel_find_cell(source_sheet, 1, 1, look_ins) or last_row_cell  # xlNext
            first_col_cell = self._excel_find_cell(source_sheet, 2, 1, look_ins) or last_col_cell

            first_row = max(1, int(first_row_cell.Row))
            first_col = max(1, int(first_col_cell.Column))
            last_row = max(1, int(last_row_cell.Row))
            last_col = max(1, int(last_col_cell.Column))
            return (
                min(first_row, last_row),
                min(first_col, last_col),
                max(first_row, last_row),
                max(first_col, last_col),
            )
        except Exception:
            return None

    def _excel_content_bounds(self, source_sheet):
        """값/수식이 들어 있는 실제 셀 범위를 반환한다."""
        value_bounds = self._excel_content_bounds_by_look_in(source_sheet, -4163)  # xlValues
        formula_bounds = self._excel_content_bounds_by_look_in(source_sheet, -4123)  # xlFormulas

        if value_bounds is None:
            if self._excel_bounds_cell_count(formula_bounds) <= self.EXCEL_VALUE_CELL_LIMIT:
                return formula_bounds
            return None

        if formula_bounds is None:
            return value_bounds

        if self._excel_bounds_cell_count(formula_bounds) > self.EXCEL_VALUE_CELL_LIMIT:
            return value_bounds

        return self._union_excel_bounds(value_bounds, formula_bounds)

    def _excel_shape_bounds(self, source_sheet):
        """도형/그림이 배치된 셀 범위를 반환한다."""
        try:
            shapes = source_sheet.Shapes
            shape_count = shapes.Count
        except Exception:
            return None

        bounds = None
        for shape_idx in range(1, shape_count + 1):
            try:
                shape = shapes(shape_idx)
                try:
                    if not bool(shape.Visible):
                        continue
                except Exception:
                    pass

                top_left = shape.TopLeftCell
                bottom_right = shape.BottomRightCell
                top_row = max(1, int(top_left.Row))
                top_col = max(1, int(top_left.Column))
                bottom_row = max(1, int(bottom_right.Row))
                bottom_col = max(1, int(bottom_right.Column))
                current = (
                    min(top_row, bottom_row),
                    min(top_col, bottom_col),
                    max(top_row, bottom_row),
                    max(top_col, bottom_col),
                )
                bounds = self._union_excel_bounds(bounds, current)
            except Exception:
                continue

        return bounds

    def _excel_used_range_bounds(self, source_sheet):
        """Excel UsedRange 범위와 셀 수를 반환한다."""
        try:
            used_range = source_sheet.UsedRange
            if used_range is None:
                return None, 0

            start_row = int(used_range.Row)
            start_col = int(used_range.Column)
            row_count = int(used_range.Rows.Count)
            col_count = int(used_range.Columns.Count)
            bounds = (
                max(1, start_row),
                max(1, start_col),
                max(1, start_row + row_count - 1),
                max(1, start_col + col_count - 1),
            )
            return bounds, row_count * col_count
        except Exception:
            return None, 0

    def _union_excel_bounds(self, left, right):
        if left is None:
            return right
        if right is None:
            return left
        return (
            min(left[0], right[0]),
            min(left[1], right[1]),
            max(left[2], right[2]),
            max(left[3], right[3]),
        )

    def _intersects_excel_bounds(self, left, right):
        return not (
            left[2] < right[0]
            or right[2] < left[0]
            or left[3] < right[1]
            or right[3] < left[1]
        )

    def _expand_excel_bounds_for_merges(self, source_sheet, bounds):
        """범위 안의 병합 셀은 병합 영역 전체가 들어가도록 확장한다."""
        if bounds is None:
            return None

        expanded = bounds
        for _ in range(2):
            before = expanded
            try:
                check_range = source_sheet.Range(
                    source_sheet.Cells(expanded[0], expanded[1]),
                    source_sheet.Cells(expanded[2], expanded[3]),
                )
                merged_areas = check_range.MergeAreas
                for area_idx in range(1, merged_areas.Count + 1):
                    area = merged_areas(area_idx)
                    if area.Cells.Count <= 1:
                        continue
                    area_bounds = (
                        int(area.Row),
                        int(area.Column),
                        int(area.Row + area.Rows.Count - 1),
                        int(area.Column + area.Columns.Count - 1),
                    )
                    if self._intersects_excel_bounds(expanded, area_bounds):
                        expanded = self._union_excel_bounds(expanded, area_bounds)
            except Exception:
                break

            if expanded == before:
                break

        return expanded

    def _get_excel_effective_range(self, source_sheet, sheet_name):
        """서식만 있는 빈 행/열을 제외한 실제 복사 범위를 계산한다."""
        used_bounds, used_cell_count = self._excel_used_range_bounds(source_sheet)
        content_bounds = self._excel_content_bounds(source_sheet)

        if content_bounds is None:
            bounds = None
        elif used_bounds is not None and 0 < used_cell_count <= self.EXCEL_VALUE_CELL_LIMIT:
            bounds = used_bounds
        else:
            bounds = content_bounds
            if bounds is not None and used_bounds is not None:
                bounds = (
                    min(bounds[0], used_bounds[0]),
                    min(bounds[1], used_bounds[1]),
                    bounds[2],
                    bounds[3],
                )

        if bounds is None:
            if content_bounds is not None and used_bounds is not None and used_cell_count <= self.EXCEL_VALUE_CELL_LIMIT:
                bounds = used_bounds
            elif used_bounds is not None:
                bounds = (1, 1, 1, 1)
            else:
                bounds = (1, 1, 1, 1)

        bounds = self._expand_excel_bounds_for_merges(source_sheet, bounds)
        start_row, start_col, end_row, end_col = bounds
        row_count = max(1, end_row - start_row + 1)
        col_count = max(1, end_col - start_col + 1)
        effective_cell_count = row_count * col_count

        if used_cell_count and used_cell_count > max(effective_cell_count * 2, self.EXCEL_VALUE_CELL_LIMIT):
            self.logger.log(
                f"    UsedRange 보정: {used_cell_count:,}셀 -> {effective_cell_count:,}셀 "
                f"({sheet_name})"
            )

        source_range = source_sheet.Range(
            source_sheet.Cells(start_row, start_col),
            source_sheet.Cells(end_row, end_col),
        )
        return source_range, start_row, start_col, row_count, col_count

    def _copy_excel_merged_cells(self, source_sheet, new_sheet, start_row, start_col, end_row, end_col):
        """계산된 범위 안의 병합 셀만 복사한다."""
        bounds = (start_row, start_col, end_row, end_col)
        copied = 0
        seen = set()
        try:
            check_range = source_sheet.Range(
                source_sheet.Cells(start_row, start_col),
                source_sheet.Cells(end_row, end_col),
            )
            merged_areas = check_range.MergeAreas
            for area_idx in range(1, merged_areas.Count + 1):
                area = merged_areas(area_idx)
                if area.Cells.Count <= 1:
                    continue

                first_row = int(area.Row)
                first_col = int(area.Column)
                last_row = int(first_row + area.Rows.Count - 1)
                last_col = int(first_col + area.Columns.Count - 1)
                area_bounds = (first_row, first_col, last_row, last_col)
                if not self._intersects_excel_bounds(bounds, area_bounds):
                    continue

                key = (first_row, first_col, last_row, last_col)
                if key in seen:
                    continue
                seen.add(key)
                new_sheet.merge_cells(
                    start_row=first_row,
                    start_column=first_col,
                    end_row=last_row,
                    end_column=last_col,
                )
                copied += 1
        except Exception:
            return

        if copied:
            self.logger.log(f"    병합 셀 복사: {copied}개")

    def _excel_range_to_rows(self, value, row_count, col_count):
        """Excel COM Range 값을 항상 행 튜플 형태로 정규화한다."""
        if row_count <= 0 or col_count <= 0:
            return []

        if row_count == 1 and col_count == 1:
            return [(value,)]

        if row_count == 1:
            if isinstance(value, (tuple, list)):
                if len(value) == 1 and isinstance(value[0], (tuple, list)):
                    return [tuple(value[0])]
                return [tuple(value)]
            return [(value,)]

        rows = value if isinstance(value, (tuple, list)) else ((value,),)
        normalized = []
        for row in rows:
            if isinstance(row, (tuple, list)):
                normalized.append(tuple(row))
            else:
                normalized.append((row,))

        return normalized

    def _excel_color_to_hex(self, color):
        """Excel 색상을 hex 문자열로 변환"""
        try:
            # color가 None인 경우만 None 반환 (검정색 #000000은 유효한 색상)
            if color is None:
                return None
            # Excel 색상은 BGR 형식
            b = (color >> 16) & 0xFF
            g = (color >> 8) & 0xFF
            r = color & 0xFF
            return f"{r:02X}{g:02X}{b:02X}"
        except Exception as e:
            self.logger.error(f"색상 변환 실패: {color}", e)
            return None

    def _points_to_pixels(self, points):
        """Office point 단위를 96DPI 기준 픽셀로 근사 변환한다."""
        try:
            return max(1, int(round(float(points) * 96 / 72)))
        except Exception:
            return 1

    def _copy_excel_sheet_objects(self, source_sheet, new_sheet, temp_dir, sheet_name):
        """Excel 시트의 삽입 그림/도형/차트 객체를 이미지로 복사한다."""
        try:
            shapes = source_sheet.Shapes
            shape_count = shapes.Count
        except Exception as e:
            self.logger.log(f"    시트 객체 목록 읽기 실패: {str(e)[:50]}")
            return 0, 0

        if shape_count <= 0:
            return 0, 0

        self.logger.log(f"    삽입 객체 복사 시작: {shape_count}개")
        copied_count = 0
        visible_count = 0

        for shape_idx in range(1, shape_count + 1):
            try:
                shape = shapes(shape_idx)
                try:
                    if not bool(shape.Visible):
                        continue
                except Exception:
                    pass
                visible_count += 1

                img_path = self._excel_shape_to_image(shape, temp_dir, sheet_name, shape_idx)
                if not img_path:
                    self.logger.log(f"    객체 {shape_idx} 이미지 변환 실패")
                    continue

                if os.path.splitext(img_path)[1].lower() not in [".png", ".jpg", ".jpeg", ".bmp", ".gif"]:
                    self.logger.log(f"    객체 {shape_idx} 지원하지 않는 이미지 형식 생략: {img_path}")
                    continue

                image = OpenpyxlImage(img_path)
                try:
                    image.width = self._points_to_pixels(shape.Width)
                    image.height = self._points_to_pixels(shape.Height)
                except Exception:
                    pass

                try:
                    top_left = shape.TopLeftCell
                    anchor = f"{get_column_letter(top_left.Column)}{top_left.Row}"
                except Exception:
                    anchor = "A1"

                image.anchor = anchor
                new_sheet.add_image(image)
                copied_count += 1

            except Exception as e:
                self.logger.log(f"    객체 {shape_idx} 복사 실패: {str(e)[:60]}")

        self.logger.log(f"    삽입 객체 복사 완료: {copied_count}/{visible_count}개")
        return copied_count, visible_count

    def _excel_shape_to_image(self, source_shape, temp_dir, sheet_name, shape_idx):
        """Excel Shape를 클립보드 경유 이미지 파일로 저장한다."""
        safe_sheet = "".join(ch if ch.isalnum() else "_" for ch in str(sheet_name))[:30] or "sheet"

        for retry in range(self.EXCEL_OBJECT_RETRY_COUNT):
            try:
                try:
                    source_shape.CopyPicture(Appearance=1, Format=2)  # 1=screen, 2=bitmap
                except Exception:
                    try:
                        source_shape.CopyPicture(1, 2)
                    except Exception:
                        source_shape.Copy()

                time.sleep(self.EXCEL_OBJECT_RETRY_DELAY)
                img_path = self._get_image_from_clipboard(temp_dir)
                if img_path and os.path.splitext(img_path)[1].lower() in [".png", ".jpg", ".jpeg", ".bmp", ".gif"]:
                    ext = os.path.splitext(img_path)[1] or ".png"
                    final_path = os.path.join(
                        temp_dir,
                        f"excel_{safe_sheet}_{shape_idx}_{retry}{ext}"
                    )
                    try:
                        if os.path.abspath(img_path) != os.path.abspath(final_path):
                            shutil.copyfile(img_path, final_path)
                            try:
                                os.remove(img_path)
                            except Exception:
                                pass
                        return final_path
                    except Exception:
                        return img_path
                if img_path:
                    self.logger.log(
                        f"    객체 {shape_idx} 클립보드 이미지 형식 변환 필요: "
                        f"{os.path.splitext(img_path)[1].lower() or 'unknown'}"
                    )
                    try:
                        os.remove(img_path)
                    except Exception:
                        pass

                chart_path = self._excel_shape_to_chart_image(
                    source_shape, temp_dir, safe_sheet, shape_idx, retry
                )
                if chart_path:
                    return chart_path
            except Exception as e:
                self.logger.log(
                    f"    객체 {shape_idx} 클립보드 복사 실패 "
                    f"(시도 {retry+1}/{self.EXCEL_OBJECT_RETRY_COUNT}): {str(e)[:50]}"
                )
                time.sleep(self.EXCEL_OBJECT_RETRY_DELAY)

        return None

    def _excel_shape_to_chart_image(self, source_shape, temp_dir, safe_sheet, shape_idx, retry):
        """Excel 임시 ChartObject를 이용해 도형/그림/OLE 객체를 PNG로 내보낸다."""
        chart_obj = None
        try:
            worksheet = source_shape.Parent
            try:
                worksheet.Activate()
            except Exception:
                pass
            left = float(source_shape.Left)
            top = float(source_shape.Top)
            width = max(2.0, float(source_shape.Width))
            height = max(2.0, float(source_shape.Height))
            chart_obj = worksheet.ChartObjects().Add(left, top, width, height)
            chart = chart_obj.Chart

            for copy_mode in ("picture_bitmap", "picture_vector", "copy"):
                try:
                    try:
                        source_shape.Select(False)
                    except Exception:
                        pass
                    if copy_mode == "picture_bitmap":
                        source_shape.CopyPicture(Appearance=1, Format=2)
                    elif copy_mode == "picture_vector":
                        source_shape.CopyPicture(Appearance=1, Format=-4147)  # xlPicture
                    else:
                        source_shape.Copy()

                    time.sleep(self.EXCEL_OBJECT_RETRY_DELAY)
                    chart.Paste()
                    png_path = os.path.join(
                        temp_dir,
                        f"excel_{safe_sheet}_{shape_idx}_{retry}_chart.png",
                    )
                    if chart.Export(Filename=png_path, FilterName="PNG") and os.path.exists(png_path):
                        if os.path.getsize(png_path) > 0:
                            return png_path
                except Exception as mode_error:
                    self.logger.log(
                        f"    객체 {shape_idx} 차트 PNG {copy_mode} 실패: {str(mode_error)[:50]}"
                    )
                    continue

        except Exception as e:
            self.logger.log(f"    객체 {shape_idx} 차트 PNG 변환 실패: {str(e)[:60]}")
        finally:
            if chart_obj is not None:
                try:
                    chart_obj.Delete()
                except Exception:
                    pass

        return None

    # ========== PPT 슬라이드 처리 메서드 ==========

    def _copy_slide_background(self, source_slide, target_slide):
        """슬라이드 배경색 복사"""
        try:
            bg = source_slide.Background
            fill = bg.Fill
            if fill.Visible:
                fill_type = fill.Type  # 1=solid
                if fill_type == 1:
                    fill_rgb = fill.ForeColor.RGB
                    r = fill_rgb & 0xFF
                    g = (fill_rgb >> 8) & 0xFF
                    b = (fill_rgb >> 16) & 0xFF
                    background = target_slide.background
                    background.fill.solid()
                    background.fill.fore_color.rgb = RGBColor(r, g, b)
        except:
            pass

    def _export_slide_as_image(self, source_slide, target_slide, temp_dir, slide_num, new_pres):
        """슬라이드 전체를 이미지로 내보내기"""
        self.logger.log(f"  슬라이드 {slide_num}를 이미지로 내보내기...")

        try:
            img_path = os.path.join(temp_dir, f"slide_{slide_num}.png")
            width = int(source_slide.Parent.PageSetup.SlideWidth * 2)
            height = int(source_slide.Parent.PageSetup.SlideHeight * 2)

            source_slide.Export(img_path, "PNG", width, height)

            target_slide.shapes.add_picture(
                img_path, Emu(0), Emu(0),
                new_pres.slide_width, new_pres.slide_height
            )
            self.logger.log(f"  슬라이드 이미지 추가 성공")

        except Exception as e:
            self.logger.error(f"  슬라이드 이미지 내보내기 실패", e)
            self._extract_text_only(source_slide, target_slide)

    def _extract_text_only(self, source_slide, target_slide):
        """텍스트만 추출"""
        for shape in source_slide.Shapes:
            try:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text
                    if text.strip():
                        left = Emu(int(shape.Left * 12700))
                        top = Emu(int(shape.Top * 12700))
                        width = Emu(int(shape.Width * 12700))
                        height = Emu(int(shape.Height * 12700))

                        textbox = target_slide.shapes.add_textbox(left, top, width, height)
                        tf = textbox.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = text

                        try:
                            src_font = shape.TextFrame.TextRange.Font
                            p.font.size = Pt(src_font.Size)
                            p.font.bold = src_font.Bold
                            p.font.italic = src_font.Italic
                        except:
                            pass

            except Exception as e:
                pass

    def _extract_hybrid(self, source_slide, target_slide, temp_dir, slide_num, new_pres):
        """하이브리드 모드"""
        shapes_list = []
        for shape in source_slide.Shapes:
            try:
                z_order = shape.ZOrderPosition
                shapes_list.append((z_order, shape))
            except:
                shapes_list.append((0, shape))

        shapes_list.sort(key=lambda x: x[0])

        for z_order, shape in shapes_list:
            try:
                if not self._recreate_shape(shape, target_slide, temp_dir):
                    self._handle_unrecreated_shape(shape, target_slide, temp_dir)
            except Exception as e:
                self.logger.log(f"도형 처리 실패, 안전 폴백 시도: {str(e)[:50]}")
                self._handle_unrecreated_shape(shape, target_slide, temp_dir)

    def _extract_text_with_object_images(self, source_slide, target_slide, temp_dir):
        """텍스트 중심 모드에서도 도형/이미지는 그림으로 보존한다."""
        self._extract_hybrid(source_slide, target_slide, temp_dir, 0, None)

    def _shape_has_editable_text(self, source_shape):
        try:
            return bool(
                source_shape.HasTextFrame
                and source_shape.TextFrame.HasText
                and source_shape.TextFrame.TextRange.Text.strip()
            )
        except Exception:
            return False

    def _shape_contains_editable_text(self, source_shape):
        if self._shape_has_editable_text(source_shape):
            return True
        try:
            if source_shape.Type == 6:
                group_items = source_shape.GroupItems
                for idx in range(1, group_items.Count + 1):
                    if self._shape_contains_editable_text(group_items(idx)):
                        return True
        except Exception:
            pass
        return False

    def _copy_shape_text_as_textbox(self, source_shape, target_slide, left=None, top=None, width=None, height=None):
        if not self._shape_has_editable_text(source_shape):
            return False
        try:
            if left is None:
                left = Emu(int(source_shape.Left * 12700))
            if top is None:
                top = Emu(int(source_shape.Top * 12700))
            if width is None:
                width = Emu(int(source_shape.Width * 12700))
            if height is None:
                height = Emu(int(source_shape.Height * 12700))
            textbox = target_slide.shapes.add_textbox(left, top, width, height)
            self._copy_text_frame(source_shape, textbox)
            return True
        except Exception as e:
            self.logger.log(f"텍스트 도형 편집 가능 복원 실패: {str(e)[:60]}")
            return False

    def _handle_group_shape(self, source_shape, target_slide, temp_dir):
        """텍스트가 포함된 그룹은 그룹 전체 이미지화 대신 구성 요소별로 복원한다."""
        try:
            group_items = source_shape.GroupItems
        except Exception as e:
            self.logger.log(f"그룹 구성 요소 접근 실패: {str(e)[:60]}")
            return False

        handled = 0
        for idx in range(1, group_items.Count + 1):
            item = group_items(idx)
            try:
                if self._recreate_shape(item, target_slide, temp_dir):
                    handled += 1
                    continue
                if self._handle_unrecreated_shape(item, target_slide, temp_dir):
                    handled += 1
            except Exception as e:
                self.logger.log(f"그룹 내부 도형 {idx} 복원 실패: {str(e)[:60]}")
                if self._shape_contains_editable_text(item):
                    if self._copy_shape_text_as_textbox(item, target_slide):
                        handled += 1
        return handled > 0

    def _handle_unrecreated_shape(self, source_shape, target_slide, temp_dir):
        if self._shape_contains_editable_text(source_shape):
            self.logger.log("텍스트 포함 도형은 이미지 폴백 금지, 편집 가능 텍스트로 복원")
            if self._copy_shape_text_as_textbox(source_shape, target_slide):
                return True
            try:
                if source_shape.Type == 6:
                    return self._handle_group_shape(source_shape, target_slide, temp_dir)
            except Exception:
                pass
            return False
        return self._add_ppt_shape_snapshot(source_shape, target_slide, temp_dir)

    def _add_ppt_shape_snapshot(self, source_shape, target_slide, temp_dir, left=None, top=None, width=None, height=None):
        """재생성할 수 없는 PPT 도형을 이미지 스냅샷으로 보존한다."""
        if self._shape_contains_editable_text(source_shape):
            self.logger.log("텍스트 포함 도형 스냅샷 차단: 이미지 변환 대신 텍스트 복원")
            if self._copy_shape_text_as_textbox(source_shape, target_slide, left, top, width, height):
                return True
            try:
                if source_shape.Type == 6:
                    return self._handle_group_shape(source_shape, target_slide, temp_dir)
            except Exception:
                pass
            return False

        img_path = None
        try:
            if left is None:
                left = Emu(int(source_shape.Left * 12700))
            if top is None:
                top = Emu(int(source_shape.Top * 12700))
            if width is None:
                width = Emu(int(source_shape.Width * 12700))
            if height is None:
                height = Emu(int(source_shape.Height * 12700))

            for retry in range(self.PPT_CLIPBOARD_RETRY_COUNT):
                try:
                    try:
                        source_shape.CopyPicture()
                    except Exception:
                        source_shape.Copy()
                    time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
                    img_path = self._get_image_from_clipboard(temp_dir)
                    if img_path:
                        break
                except Exception as e:
                    if retry == self.PPT_CLIPBOARD_RETRY_COUNT - 1:
                        self.logger.log(f"도형 클립보드 보존 실패: {str(e)[:60]}")
                    time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)

            if not img_path or not os.path.exists(img_path):
                try:
                    img_path = os.path.join(temp_dir, f"shape_{id(source_shape)}.png")
                    source_shape.Export(img_path, 2)
                except Exception:
                    img_path = None

            if img_path and os.path.exists(img_path):
                target_slide.shapes.add_picture(img_path, left, top, width, height)
                return True

            return False
        except Exception as e:
            self.logger.log(f"도형 이미지 스냅샷 추가 실패: {str(e)[:60]}")
            return False

    def _recreate_shape(self, source_shape, target_slide, temp_dir):
        """도형 재생성"""
        shape_type = source_shape.Type

        try:
            left = Emu(int(source_shape.Left * 12700))
            top = Emu(int(source_shape.Top * 12700))
            width = Emu(int(source_shape.Width * 12700))
            height = Emu(int(source_shape.Height * 12700))

            # 이미지
            if shape_type in [13, 11]:
                return self._handle_image_shape(source_shape, target_slide, temp_dir, left, top, width, height)

            # 텍스트박스
            if shape_type == 17:
                textbox = target_slide.shapes.add_textbox(left, top, width, height)
                self._copy_text_frame(source_shape, textbox)
                return True

            # AutoShape
            if shape_type == 1:
                return self._handle_autoshape(source_shape, target_slide, left, top, width, height)

            # Placeholder
            if shape_type == 14:
                textbox = target_slide.shapes.add_textbox(left, top, width, height)
                self._copy_text_frame(source_shape, textbox)
                return True

            # Table
            if shape_type == 19:
                return self._handle_table(source_shape, target_slide, left, top, width, height)

            # Connector
            if shape_type == 9:
                return self._handle_connector(source_shape, target_slide, left, top)

            # Group
            if shape_type == 6:
                if self._shape_contains_editable_text(source_shape):
                    return self._handle_group_shape(source_shape, target_slide, temp_dir)
                return self._add_ppt_shape_snapshot(source_shape, target_slide, temp_dir, left, top, width, height)

            # Freeform
            if shape_type == 5:
                return self._handle_freeform(source_shape, target_slide, temp_dir, left, top, width, height)

            # 기타 - 텍스트만
            if source_shape.HasTextFrame and source_shape.TextFrame.HasText:
                text = source_shape.TextFrame.TextRange.Text
                if text.strip():
                    textbox = target_slide.shapes.add_textbox(left, top, width, height)
                    self._copy_text_frame(source_shape, textbox)
                    return True

            return False

        except Exception as e:
            self.logger.log(f"도형 재생성 실패 (타입 {shape_type}): {str(e)[:50]}")
            return False

    def _handle_image_shape(self, source_shape, target_slide, temp_dir, left, top, width, height):
        """이미지 도형 처리"""
        img_path = None
        export_error = None

        # 클립보드 우선: 일부 보안 PPT는 Shape.Export 호출 때 PowerPoint 경고창을 띄운다.
        for retry in range(self.PPT_CLIPBOARD_RETRY_COUNT):
            clipboard_img = None
            try:
                try:
                    source_shape.CopyPicture()
                except Exception:
                    source_shape.Copy()
                time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
                clipboard_img = self._get_image_from_clipboard(temp_dir)
                if clipboard_img:
                    target_slide.shapes.add_picture(clipboard_img, left, top, width, height)
                    return True
            except Exception as e:
                self.logger.log(
                    f"클립보드 이미지 추출 실패 "
                    f"(시도 {retry+1}/{self.PPT_CLIPBOARD_RETRY_COUNT}): {str(e)[:50]}"
                )
                time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
            finally:
                # 클립보드에서 생성된 임시 파일 정리
                if clipboard_img and os.path.exists(clipboard_img):
                    try:
                        os.remove(clipboard_img)
                    except Exception:
                        pass

        # Export 폴백
        try:
            img_path = os.path.join(temp_dir, f"img_{id(source_shape)}.png")
            source_shape.Export(img_path, 2)
            target_slide.shapes.add_picture(img_path, left, top, width, height)
            return True
        except Exception as e:
            export_error = e
            if img_path and os.path.exists(img_path):
                try:
                    os.remove(img_path)
                except Exception:
                    pass

        # Placeholder
        try:
            if export_error:
                self.logger.log(f"이미지 보존 폴백 실패, Placeholder 생성: {str(export_error)[:60]}")
            placeholder = target_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
            placeholder.text_frame.paragraphs[0].text = "[이미지]"
            return True
        except Exception as e:
            self.logger.error("이미지 Placeholder 생성 실패", e)
            return False

    def _handle_autoshape(self, source_shape, target_slide, left, top, width, height):
        """AutoShape 처리"""
        try:
            auto_shape_type = source_shape.AutoShapeType
            pptx_shape_type = AUTOSHAPE_MAPPING.get(auto_shape_type, MSO_SHAPE.RECTANGLE)
            new_shape = target_slide.shapes.add_shape(pptx_shape_type, left, top, width, height)

            # 회전
            try:
                rotation = source_shape.Rotation
                if rotation and rotation != 0:
                    new_shape.rotation = rotation
            except:
                pass

            # 채우기
            try:
                fill = source_shape.Fill
                if fill.Visible:
                    fill_type = fill.Type  # 1=solid,2=gradient,3=texture,4=pattern,5=background
                    if fill_type == 1:  # solid
                        fill_rgb = fill.ForeColor.RGB
                        r, g, b = fill_rgb & 0xFF, (fill_rgb >> 8) & 0xFF, (fill_rgb >> 16) & 0xFF
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = RGBColor(r, g, b)
                        # 투명도
                        try:
                            transparency = fill.Transparency
                            if transparency and transparency > 0:
                                new_shape.fill.fore_color.theme_color  # 접근 확인
                                # XML로 직접 투명도 설정
                                from lxml import etree
                                from pptx.oxml.ns import qn as _qn
                                spPr = new_shape._element.spPr
                                solidFill = spPr.find('.//' + _qn('a:solidFill'))
                                if solidFill is not None:
                                    srgbClr = solidFill.find(_qn('a:srgbClr'))
                                    if srgbClr is not None:
                                        alpha = etree.SubElement(srgbClr, _qn('a:alpha'))
                                        alpha.set('val', str(int((1 - transparency) * 100000)))
                        except:
                            pass
                    else:
                        # gradient/pattern → 단색 근사
                        try:
                            fill_rgb = fill.ForeColor.RGB
                            r, g, b = fill_rgb & 0xFF, (fill_rgb >> 8) & 0xFF, (fill_rgb >> 16) & 0xFF
                            new_shape.fill.solid()
                            new_shape.fill.fore_color.rgb = RGBColor(r, g, b)
                        except:
                            pass
                else:
                    new_shape.fill.background()  # 투명
            except:
                pass

            # 테두리
            try:
                line = source_shape.Line
                if line.Visible:
                    line_rgb = line.ForeColor.RGB
                    r, g, b = line_rgb & 0xFF, (line_rgb >> 8) & 0xFF, (line_rgb >> 16) & 0xFF
                    new_shape.line.color.rgb = RGBColor(r, g, b)
                    new_shape.line.width = Pt(line.Weight)
                else:
                    new_shape.line.fill.background()  # 테두리 없음
            except:
                pass

            self._copy_text_frame(source_shape, new_shape)
            return True
        except Exception as e:
            self.logger.log(f"AutoShape 처리 실패: {str(e)[:60]}")
            return False

    def _handle_table(self, source_shape, target_slide, left, top, width, height):
        """테이블 처리 — 셀 텍스트·배경색·폰트 복사"""
        try:
            src_table = source_shape.Table
            rows = src_table.Rows.Count
            cols = src_table.Columns.Count
            tbl_shape = target_slide.shapes.add_table(rows, cols, left, top, width, height)
            table = tbl_shape.table

            # 열 너비 복사
            try:
                for c in range(1, cols + 1):
                    col_width = src_table.Columns(c).Width
                    table.columns[c-1].width = Emu(int(col_width * 12700))
            except:
                pass

            # 행 높이 복사
            try:
                for r in range(1, rows + 1):
                    row_height = src_table.Rows(r).Height
                    table.rows[r-1].height = Emu(int(row_height * 12700))
            except:
                pass

            for r in range(1, rows + 1):
                for c in range(1, cols + 1):
                    try:
                        src_cell = src_table.Cell(r, c)
                        dst_cell = table.cell(r-1, c-1)
                        src_cell_shape = src_cell.Shape

                        # 텍스트+서식 복사
                        self._copy_text_frame(src_cell_shape, dst_cell)

                        # 셀 배경색
                        try:
                            fill = src_cell_shape.Fill
                            if fill.Visible and fill.Type == 1:
                                fill_rgb = fill.ForeColor.RGB
                                r_c = fill_rgb & 0xFF
                                g_c = (fill_rgb >> 8) & 0xFF
                                b_c = (fill_rgb >> 16) & 0xFF
                                dst_cell.fill.solid()
                                dst_cell.fill.fore_color.rgb = RGBColor(r_c, g_c, b_c)
                        except:
                            pass

                    except Exception as cell_err:
                        self.logger.log(f"    테이블 셀({r},{c}) 처리 실패: {str(cell_err)[:60]}")
            return True
        except Exception as e:
            self.logger.log(f"테이블 처리 실패: {str(e)[:60]}")
            return False

    def _handle_connector(self, source_shape, target_slide, left, top):
        """연결선 처리"""
        try:
            end_x = Emu(int((source_shape.Left + source_shape.Width) * 12700))
            end_y = Emu(int((source_shape.Top + source_shape.Height) * 12700))

            connector = target_slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, left, top, end_x, end_y
            )

            try:
                line_rgb = source_shape.Line.ForeColor.RGB
                r, g, b = line_rgb & 0xFF, (line_rgb >> 8) & 0xFF, (line_rgb >> 16) & 0xFF
                connector.line.color.rgb = RGBColor(r, g, b)
                connector.line.width = Pt(source_shape.Line.Weight)
            except:
                pass
            return True
        except Exception as e:
            self.logger.log(f"연결선 처리 실패: {str(e)[:60]}")
            return False

    def _handle_freeform(self, source_shape, target_slide, temp_dir, left, top, width, height):
        """Freeform 처리"""
        if self._shape_has_editable_text(source_shape):
            return self._copy_shape_text_as_textbox(source_shape, target_slide, left, top, width, height)

        img_path = None
        for retry in range(self.PPT_CLIPBOARD_RETRY_COUNT):
            try:
                source_shape.Copy()
                time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
                img_path = self._get_image_from_clipboard(temp_dir)
                if img_path:
                    target_slide.shapes.add_picture(img_path, left, top, width, height)
                    return True
            except Exception as e:
                self.logger.log(
                    f"Freeform 클립보드 추출 실패 "
                    f"(시도 {retry+1}/{self.PPT_CLIPBOARD_RETRY_COUNT}): {str(e)[:50]}"
                )
                time.sleep(self.PPT_CLIPBOARD_RETRY_DELAY)
            finally:
                # 임시 파일 정리
                if img_path and os.path.exists(img_path):
                    try:
                        os.remove(img_path)
                    except Exception:
                        pass

        # 텍스트만
        if source_shape.HasTextFrame and source_shape.TextFrame.HasText:
            textbox = target_slide.shapes.add_textbox(left, top, width, height)
            self._copy_text_frame(source_shape, textbox)
            return True
        return False

    def _copy_text_frame(self, source_shape, target_shape):
        """텍스트 프레임 복사 — 단락/런 단위 서식 보존"""
        try:
            if not source_shape.HasTextFrame or not source_shape.TextFrame.HasText:
                return

            src_tf = source_shape.TextFrame
            dst_tf = target_shape.text_frame
            dst_tf.word_wrap = True

            # 수직 정렬
            try:
                va = src_tf.TextRange.ParagraphFormat.Alignment  # COM 수직은 별도
            except:
                pass
            try:
                anchor_map = {1: MSO_ANCHOR.TOP, 3: MSO_ANCHOR.MIDDLE, 4: MSO_ANCHOR.BOTTOM}
                va_val = src_tf.VerticalAnchor  # 1=top,2=middle(없음),3=middle,4=bottom
                if va_val in anchor_map:
                    dst_tf.vertical_anchor = anchor_map[va_val]
            except:
                pass

            para_count = src_tf.TextRange.Paragraphs().Count

            # 기존 단락 초기화
            from pptx.oxml.ns import qn as _qn
            from lxml import etree
            txBody = dst_tf._txBody
            for old_p in txBody.findall(_qn('a:p')):
                txBody.remove(old_p)

            for pi in range(1, para_count + 1):
                src_para = src_tf.TextRange.Paragraphs(pi)

                # 새 단락 XML 엘리먼트 생성
                new_p = etree.SubElement(txBody, _qn('a:p'))

                # 단락 서식 (pPr)
                try:
                    pPr = etree.SubElement(new_p, _qn('a:pPr'))
                    align_map = {1: 'l', 2: 'ctr', 3: 'r', 4: 'dist', 5: 'just'}
                    align_val = src_para.ParagraphFormat.Alignment
                    if align_val in align_map:
                        pPr.set('algn', align_map[align_val])
                    indent = src_para.ParagraphFormat.Indent
                    if indent and indent != 0:
                        pPr.set('indent', str(int(indent * 12700)))
                    space_before = src_para.ParagraphFormat.SpaceBefore
                    if space_before and space_before > 0:
                        spcBef = etree.SubElement(pPr, _qn('a:spcBef'))
                        spcPts = etree.SubElement(spcBef, _qn('a:spcPts'))
                        spcPts.set('val', str(int(space_before * 100)))
                except:
                    pass

                # 런(Run) 단위 텍스트+서식
                try:
                    run_count = src_para.Runs.Count
                    if run_count == 0:
                        # 빈 단락
                        etree.SubElement(new_p, _qn('a:endParaRPr'), attrib={'lang': 'ko-KR'})
                        continue

                    for ri in range(1, run_count + 1):
                        src_run = src_para.Runs(ri)
                        run_text = src_run.Text
                        src_font = src_run.Font

                        new_r = etree.SubElement(new_p, _qn('a:r'))
                        rPr = etree.SubElement(new_r, _qn('a:rPr'), attrib={'lang': 'ko-KR', 'dirty': '0'})

                        # 굵게/기울임/밑줄
                        try:
                            if src_font.Bold:
                                rPr.set('b', '1')
                        except: pass
                        try:
                            if src_font.Italic:
                                rPr.set('i', '1')
                        except: pass
                        try:
                            if src_font.Underline:
                                rPr.set('u', 'sng')
                        except: pass
                        # 글꼴 크기
                        try:
                            if src_font.Size:
                                rPr.set('sz', str(int(src_font.Size * 100)))
                        except: pass
                        # 글꼴 이름
                        try:
                            if src_font.Name:
                                latin = etree.SubElement(rPr, _qn('a:latin'))
                                latin.set('typeface', src_font.Name)
                        except: pass
                        # 글자 색상
                        try:
                            rgb_val = src_font.Color.RGB
                            r = rgb_val & 0xFF
                            g = (rgb_val >> 8) & 0xFF
                            b = (rgb_val >> 16) & 0xFF
                            solidFill = etree.SubElement(rPr, _qn('a:solidFill'))
                            srgbClr = etree.SubElement(solidFill, _qn('a:srgbClr'))
                            srgbClr.set('val', f'{r:02X}{g:02X}{b:02X}')
                        except: pass

                        # 텍스트
                        t_elem = etree.SubElement(new_r, _qn('a:t'))
                        t_elem.text = run_text if run_text else ''

                except Exception as re:
                    # 런 읽기 실패 시 전체 단락 텍스트로 폴백
                    try:
                        para_text = src_para.Text
                        new_r = etree.SubElement(new_p, _qn('a:r'))
                        rPr = etree.SubElement(new_r, _qn('a:rPr'), attrib={'lang': 'ko-KR'})
                        try:
                            sz = src_tf.TextRange.Font.Size
                            if sz:
                                rPr.set('sz', str(int(sz * 100)))
                        except: pass
                        t_elem = etree.SubElement(new_r, _qn('a:t'))
                        t_elem.text = para_text if para_text else ''
                    except:
                        pass

        except Exception as e:
            self.logger.log(f"텍스트 프레임 복사 실패: {str(e)[:80]}")

    def _get_image_from_clipboard(self, temp_dir):
        """클립보드에서 이미지 추출"""
        try:
            import win32clipboard
            from PIL import Image
            import struct

            win32clipboard.OpenClipboard()
            try:
                # DIB
                if win32clipboard.IsClipboardFormatAvailable(8):
                    data = win32clipboard.GetClipboardData(8)
                    img = self._dib_to_image(data)
                    if img:
                        img_path = os.path.join(temp_dir, f"clipboard_{id(data)}.png")
                        img.save(img_path, "PNG")
                        return img_path

                # PNG
                png_format = win32clipboard.RegisterClipboardFormat("PNG")
                if win32clipboard.IsClipboardFormatAvailable(png_format):
                    data = win32clipboard.GetClipboardData(png_format)
                    img_path = os.path.join(temp_dir, f"clipboard_{id(data)}.png")
                    with open(img_path, 'wb') as f:
                        f.write(data)
                    return img_path

                # EMF
                if win32clipboard.IsClipboardFormatAvailable(14):
                    import ctypes
                    hemf = win32clipboard.GetClipboardData(14)
                    gdi32 = ctypes.windll.gdi32
                    size = gdi32.GetEnhMetaFileBits(hemf, 0, None)
                    if size > 0:
                        buffer = ctypes.create_string_buffer(size)
                        gdi32.GetEnhMetaFileBits(hemf, size, buffer)
                        emf_path = os.path.join(temp_dir, f"clipboard_{id(hemf)}.emf")
                        with open(emf_path, 'wb') as f:
                            f.write(buffer.raw)
                        return emf_path

            finally:
                win32clipboard.CloseClipboard()
        except:
            pass

        try:
            from PIL import Image, ImageGrab

            grabbed = ImageGrab.grabclipboard()
            if isinstance(grabbed, Image.Image):
                img_path = os.path.join(temp_dir, f"clipboard_grab_{int(time.time() * 1000)}.png")
                if grabbed.mode not in ("RGB", "RGBA"):
                    grabbed = grabbed.convert("RGBA")
                grabbed.save(img_path, "PNG")
                return img_path

            if isinstance(grabbed, list):
                for item in grabbed:
                    if not item or not os.path.exists(item):
                        continue
                    ext = os.path.splitext(item)[1].lower()
                    if ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif"]:
                        img_path = os.path.join(
                            temp_dir,
                            f"clipboard_file_{int(time.time() * 1000)}{ext}",
                        )
                        shutil.copyfile(item, img_path)
                        return img_path
        except:
            pass
        return None

    def _dib_to_image(self, dib_data):
        """DIB 데이터를 PIL Image로 변환"""
        try:
            from PIL import Image
            import struct

            header_size = struct.unpack('<I', dib_data[0:4])[0]
            width = struct.unpack('<i', dib_data[4:8])[0]
            height = struct.unpack('<i', dib_data[8:12])[0]
            bit_count = struct.unpack('<H', dib_data[14:16])[0]

            if bit_count == 32:
                pixel_data = dib_data[header_size:]
                img = Image.frombytes('RGBA', (width, abs(height)), pixel_data, 'raw', 'BGRA')
                if height > 0:
                    img = img.transpose(Image.FLIP_TOP_BOTTOM)
                return img

            elif bit_count == 24:
                row_size = ((width * 3 + 3) // 4) * 4
                pixel_data = dib_data[header_size:]
                orientation = -1 if height > 0 else 1
                return Image.frombytes(
                    'RGB',
                    (width, abs(height)),
                    pixel_data,
                    'raw',
                    'BGR',
                    row_size,
                    orientation,
                )
        except:
            pass
        return None

    # ========== 한글 관련 메서드 ==========

    def browse_hwp_source_path(self):
        """원본 한글 파일 직접 선택 (보안 PC 권장 경로)"""
        self.logger.log("한글 원본 파일 선택 대화상자 열기")
        path = filedialog.askopenfilename(
            title="원본 한글 파일 선택",
            filetypes=[("한글 문서", "*.hwp *.hwpx"), ("모든 파일", "*.*")],
        )
        if not path:
            return
        self.hwp_source_path.set(path)
        self.logger.log(f"한글 원본 파일 선택됨: {path}")
        if not self.hwp_save_path.get().strip():
            stem = os.path.splitext(path)[0]
            ext = ".hwpx" if self.hwp_save_format.get() == "hwpx" else ".hwp"
            self.hwp_save_path.set(f"{stem}_복사본{ext}")

    def browse_hwp_save_path(self):
        """한글 저장 경로 선택"""
        self.logger.log("한글 저장 경로 선택 대화상자 열기")

        doc_name = self.hwp_doc_name.get()
        save_format = self.hwp_save_format.get()

        if doc_name and doc_name != "감지 중..." and doc_name != "열린 한글 없음":
            default_name = os.path.splitext(doc_name)[0] + "_복사본"
        else:
            default_name = "새문서"

        if save_format == "hwpx":
            ext = ".hwpx"
            filetypes = [("한글 2014+ 파일", "*.hwpx")]
        else:
            ext = ".hwp"
            filetypes = [("한글 파일", "*.hwp")]

        path = filedialog.asksaveasfilename(
            defaultextension=ext,
            filetypes=filetypes,
            initialfile=default_name,
            title="저장할 위치 선택"
        )
        if path:
            self.hwp_save_path.set(path)
            self.logger.log(f"한글 저장 경로 선택됨: {path}")

    def detect_open_hwp(self):
        """열려있는 한글 감지"""
        if self._hwp_detecting:
            self.logger.log("한글 감지 요청 생략: 이전 감지가 아직 진행 중")
            return

        self.logger.log("한글 감지 시작")
        self.status_text.set("한글 감지 중...")
        self.hwp_doc_name.set("감지 중...")
        self._hwp_detecting = True

        thread = threading.Thread(target=self._detect_hwp)
        thread.daemon = True
        thread.start()

    def _detect_hwp(self):
        """한글 감지 (백그라운드)"""
        self.logger.log("백그라운드 한글 감지 스레드 시작")
        pythoncom.CoInitialize()

        try:
            try:
                hwp, _ = self._get_hwp_app(allow_dispatch=False)
            except Exception as connect_err:
                self.logger.log(
                    f"한글 활성 COM 연결 실패, 창 제목 감지로 대체: {str(connect_err)[:80]}"
                )
                hwp_windows = self._list_hwp_windows()
                if hwp_windows:
                    self.hwp_list = [
                        (title, "", idx + 1, hwnd)
                        for idx, (hwnd, title) in enumerate(hwp_windows)
                    ]
                    hwp_titles = [title for _, title in hwp_windows]
                    self.logger.log(f"한글 창 감지: {len(hwp_titles)}개")

                    def update_window_combo():
                        self.hwp_combo['values'] = hwp_titles
                        self.hwp_combo.current(0)
                        self.selected_hwp_index.set(1)
                        self.hwp_doc_name.set(hwp_titles[0])
                        self.status_text.set("한글 창 감지됨")

                    self.root.after(0, update_window_combo)
                else:
                    self.hwp_list = []

                    def clear_window_combo():
                        self.hwp_combo.set("")
                        self.hwp_combo['values'] = []
                        self.hwp_doc_name.set("열린 한글 없음")
                        self.status_text.set("한글을 먼저 열어주세요")

                    self.root.after(0, clear_window_combo)
                return

            # 열린 문서 확인
            try:
                path = hwp.Path
                if path:
                    name = os.path.basename(path)
                    self.hwp_list = [(name, path, 1, 0)]
                    self.logger.log(f"한글 문서 감지: {name}")

                    def update_combo():
                        self.hwp_combo['values'] = [name]
                        self.hwp_combo.current(0)
                        self.selected_hwp_index.set(1)
                        self.hwp_doc_name.set(name)
                        self.status_text.set("한글 문서 감지됨")

                    self.root.after(0, update_combo)
                else:
                    # 제목 없는 문서
                    self.hwp_list = [("제목 없음", "", 1, 0)]
                    self.logger.log("한글 제목 없는 문서 감지")

                    def update_combo():
                        self.hwp_combo['values'] = ["제목 없음"]
                        self.hwp_combo.current(0)
                        self.selected_hwp_index.set(1)
                        self.hwp_doc_name.set("제목 없음")
                        self.status_text.set("한글 문서 감지됨")

                    self.root.after(0, update_combo)
            except Exception as e:
                self.logger.log(f"한글 문서 정보 가져오기 실패: {str(e)}")
                self.hwp_list = []

                def clear_combo():
                    self.hwp_combo.set("")
                    self.hwp_combo['values'] = []
                    self.hwp_doc_name.set("열린 한글 없음")
                    self.status_text.set("한글을 먼저 열어주세요")

                self.root.after(0, clear_combo)

        except Exception as e:
            self.logger.error("한글 감지 실패", e)
            self.hwp_list = []
            err_msg = str(e)[:30]

            def show_error():
                self.hwp_combo.set("")
                self.hwp_doc_name.set("열린 한글 없음")
                self.status_text.set(f"한글 감지 실패: {err_msg}")

            self.root.after(0, show_error)

        finally:
            self._hwp_detecting = False
            pythoncom.CoUninitialize()

    def on_hwp_selected(self, event):
        """한글 콤보박스 선택 이벤트"""
        selected_idx = self.hwp_combo.current()
        if selected_idx >= 0 and selected_idx < len(self.hwp_list):
            name, path, hwp_index, hwnd = self._unpack_hwp_item(self.hwp_list[selected_idx])
            self.selected_hwp_index.set(hwp_index)
            self.hwp_doc_name.set(name)
            self.logger.log(f"한글 선택: {name} (인덱스 {hwp_index}, hwnd={hwnd})")

    def start_hwp_extraction(self):
        """한글 추출 시작"""
        self.logger.log("한글 추출 시작 버튼 클릭")

        source_path = self.hwp_source_path.get().strip()

        # 원본 미선택 시 원본 파일 선택을 먼저 유도한다.
        # (보안 PC에선 열린 문서 COM 연결이 막혀 직접 Open이 사실상 유일한 경로다.)
        if not source_path:
            self.logger.log("원본 미선택 → 원본 파일 선택 대화상자 자동 표시")
            self.browse_hwp_source_path()
            source_path = self.hwp_source_path.get().strip()

        save_path = self.hwp_save_path.get().strip()
        save_format = self.hwp_save_format.get()

        # 원본 직접 선택 모드 (권장: 새 한글 인스턴스로 Open → 메모리 추출)
        if source_path:
            if not os.path.exists(source_path):
                messagebox.showwarning("경고", "원본 파일을 찾을 수 없습니다.")
                return
            if not save_path:
                messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
                return
            self.hwp_extract_button.config(state=tk.DISABLED)
            self.progress_var.set(0)
            thread = threading.Thread(
                target=self._extract_hwp, args=(save_path, save_format, None, source_path))
            thread.daemon = True
            thread.start()
            return

        # 원본을 선택하지 않으면(취소) 종료 — 회사 보안 PC에선 원본 직접 Open만 가능하다.
        messagebox.showwarning("경고", "원본 한글 파일을 선택해주세요.")

    def _extract_hwp(self, save_path, save_format, hwp_item=None, source_path=None):
        """한글 추출 (백그라운드)"""
        self.logger.log("=== 한글 추출 프로세스 시작 ===")
        extract_start = time.perf_counter()
        pythoncom.CoInitialize()

        try:
            # 원본 파일 직접 변환 모드 (보안 PC 권장: 열린 문서 COM 연결을 건너뜀)
            if source_path:
                self.logger.log(f"원본 파일 직접 변환 모드: {source_path}")
                self._extract_hwp_from_file(source_path, save_path, save_format, extract_start)
                return

            self.root.after(0, lambda: self.status_text.set("원본 한글 연결 중..."))

            try:
                hwp, _ = self._get_hwp_app_for_extraction()
            except Exception as connect_error:
                hwnd = self._find_hwp_window_for_item(hwp_item)
                if not hwnd:
                    raise

                self.logger.log(f"한글 COM 연결 실패, UI 저장 대체 경로로 전환: {str(connect_error)[:120]}")
                self.root.after(0, lambda: self.status_text.set("한글 창에서 직접 저장 중..."))
                self.root.after(0, lambda: self.progress_var.set(30))
                self._save_hwp_via_window(hwnd, save_path, save_format)

                self.root.after(0, lambda: self.progress_var.set(100))
                self.root.after(0, lambda: self.status_text.set("한글 추출 완료!"))
                self.root.after(0, lambda: messagebox.showinfo("완료",
                    f"한글 추출 완료 (창 저장 방식)!\n{save_path}"))
                self.logger.log(f"한글 UI 저장 완료: {save_path}")
                self._log_elapsed("한글 전체 추출 시간", extract_start)
                return

            try:
                current_path = hwp.Path
            except Exception:
                current_path = ""

            self.logger.log(f"원본 한글 문서 연결 성공")
            self.root.after(0, lambda: self.status_text.set("새 문서로 저장 중..."))
            self.root.after(0, lambda: self.progress_var.set(30))

            # 열린 문서 연결 성공 → 3단계 폴백 저장 (직접 SaveAs → 메모리 재구성 → HWPML 직접)
            self._hwp_save_with_fallbacks(hwp, save_path, save_format, extract_start)

        except Exception as e:
            error_message = str(e)
            self.logger.error("한글 추출 오류", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"추출 중 오류:\n{error_message}"))

        finally:
            self.root.after(0, lambda: self.hwp_extract_button.config(state=tk.NORMAL))
            pythoncom.CoUninitialize()

    # ========== Word 관련 메서드 ==========

    def _collect_word_runs(self, source_range, full_text):
        """Word 단락 Range에서 서식 동질 런(run) 목록을 추출.

        어절(Words) 단위로 Font 속성을 한 번씩만 읽어 서식 변경점을 감지.
        문자 단위 순회 대비 COM 왕복 횟수를 평균 5~10배 줄인다.

        반환: [(text, font_name, font_size, bold, italic, underline, color), ...]
              어절 접근 실패 시 빈 리스트.
        """
        runs_data = []
        try:
            words_col = source_range.Words
            words_count = words_col.Count
        except Exception:
            return runs_data

        if words_count <= 0:
            return runs_data

        current_text = ""
        prev_key = None
        prev_attrs = None

        for w_idx in range(1, words_count + 1):
            try:
                word_range = words_col(w_idx)
                word_text = word_range.Text
                if not word_text:
                    continue

                font = word_range.Font
                try:
                    font_name = font.Name
                except Exception:
                    font_name = None
                try:
                    font_size = font.Size
                except Exception:
                    font_size = None
                try:
                    bold = font.Bold
                except Exception:
                    bold = None
                try:
                    italic = font.Italic
                except Exception:
                    italic = None
                try:
                    underline = font.Underline
                except Exception:
                    underline = None
                try:
                    color = font.Color
                except Exception:
                    color = 0

                key = (font_name, font_size, bold, italic, underline, color)

                if prev_key is None:
                    current_text = word_text
                    prev_key = key
                    prev_attrs = (font_name, font_size, bold, italic, underline, color)
                elif key == prev_key:
                    current_text += word_text
                else:
                    runs_data.append((current_text,) + prev_attrs)
                    current_text = word_text
                    prev_key = key
                    prev_attrs = (font_name, font_size, bold, italic, underline, color)
            except Exception:
                continue

        if current_text and prev_attrs is not None:
            runs_data.append((current_text,) + prev_attrs)

        # 런 추출이 아예 실패했는데 full_text는 있는 경우의 안전망
        if not runs_data and full_text:
            runs_data.append((full_text, None, None, None, None, None, 0))

        return runs_data

    def _get_word_app(self, allow_dispatch=True):
        """Word 애플리케이션 연결 (3단계 폴백 공통화).

        반환: (word_app, created_new) — created_new는 Dispatch로 새 인스턴스를 만든 경우 True.
        실패 시 예외.
        """
        return self._connect_com_app("Word.Application", "Word", allow_dispatch=allow_dispatch)

    def _setup_word_tab(self):
        """Word 탭 설정"""
        tab = self.word_tab

        # 문서 정보 프레임
        info_frame = self._create_section(tab, "Word 입력 선택")

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=2)
        ttk.Label(source_inner, text="파일 선택:", width=12).pack(side=tk.LEFT)
        self.word_source_entry = ttk.Entry(source_inner, textvariable=self.word_source_path, width=45, state="readonly")
        self.word_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_word_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # Word 선택 콤보박스
        select_frame = ttk.Frame(info_frame)
        select_frame.pack(fill=tk.X, pady=2)
        ttk.Label(select_frame, text="Word 선택:", width=12).pack(side=tk.LEFT)
        self.word_combo = ttk.Combobox(select_frame, state="readonly", width=40)
        self.word_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.word_combo.bind("<<ComboboxSelected>>", self.on_word_selected)

        # 페이지 수
        page_frame = ttk.Frame(info_frame)
        page_frame.pack(fill=tk.X, pady=2)
        ttk.Label(page_frame, text="페이지 수:", width=12).pack(side=tk.LEFT)
        ttk.Label(page_frame, textvariable=self.word_page_count,
                  font=("맑은 고딕", 10, "bold")).pack(side=tk.LEFT)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=lambda: self.detect_open_word(prefer_open=True),
                   style="Secondary.TButton").pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = self._create_section(tab, "새 파일 저장 위치")

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.word_save_path, width=45).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_word_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # 추출 옵션 프레임
        option_frame = self._create_section(tab, "추출 옵션")

        self.word_include_format = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="서식 포함 (글꼴, 색상, 정렬, 들여쓰기)",
                        variable=self.word_include_format).pack(anchor=tk.W)
        self.word_use_saveas = tk.BooleanVar(value=True)
        ttk.Checkbutton(option_frame, text="원본 파일 복사 우선 (원본 상태 변경 없음)",
                        variable=self.word_use_saveas).pack(anchor=tk.W)

        # 추출 버튼
        self.word_extract_button = ttk.Button(tab, text="새 Word로 내보내기",
                                               command=self.start_word_extraction,
                                               style="Accent.TButton")
        self.word_extract_button.pack(pady=10)

    def browse_word_save_path(self):
        """Word 저장 경로 선택"""
        self.logger.log("Word 저장 경로 선택")

        direct_source = self.word_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("word"):
            doc_name = os.path.basename(direct_source)
        else:
            doc_name = self.word_doc_name.get()
        if doc_name and doc_name != "감지 중..." and doc_name != "열린 Word 없음":
            src_ext = os.path.splitext(doc_name)[1] or ".docx"
            default_ext = src_ext if src_ext.lower() in [".docx", ".docm", ".doc", ".rtf"] else ".docx"
            default_name = os.path.splitext(doc_name)[0] + "_복사본" + default_ext
        else:
            default_ext = ".docx"
            default_name = "새문서.docx"

        path = filedialog.asksaveasfilename(
            defaultextension=default_ext,
            filetypes=[("Word 파일", "*.docx *.docm *.doc *.rtf"), ("모든 파일", "*.*")],
            initialfile=default_name,
            title="저장할 위치 선택"
        )
        if path:
            self.word_save_path.set(path)
            self.logger.log(f"Word 저장 경로: {path}")

    def detect_open_word(self, prefer_open=False):
        """열려있는 Word 감지"""
        if prefer_open:
            self._use_open_document_input("word", self.word_source_path, "Word")
        self.logger.log("Word 감지 시작")
        if self._is_direct_file_input_active("word"):
            self.status_text.set("Word 감지 중... (파일 선택 유지)")
            self._show_direct_file_input("word", self.word_source_path.get().strip(), "Word")
        else:
            self.status_text.set("Word 감지 중...")
            self.word_doc_name.set("감지 중...")
            self.word_page_count.set("-")

        thread = threading.Thread(target=self._detect_word)
        thread.daemon = True
        thread.start()

    def _detect_word(self):
        """Word 감지 (백그라운드)"""
        self.logger.log("백그라운드 Word 감지 스레드 시작")
        pythoncom.CoInitialize()

        try:
            word, _ = self._get_word_app(allow_dispatch=False)

            doc_count = word.Documents.Count
            self.logger.log(f"Word 연결 성공, 열린 문서 수: {doc_count}")

            if doc_count > 0:
                word_names = []
                word_info = []

                for i in range(1, doc_count + 1):
                    try:
                        doc = word.Documents(i)
                        name = doc.Name
                        try:
                            page_count = doc.ComputeStatistics(2)  # wdStatisticPages = 2
                        except Exception:
                            page_count = 0
                        word_names.append(f"{name} ({page_count}페이지)")
                        word_info.append((name, page_count, i))
                        self.logger.log(f"  Word {i}: {name}, {page_count}페이지")
                    except Exception as e:
                        self.logger.log(f"  Word {i} 정보 가져오기 실패: {str(e)}")

                self.word_list = word_info

                def update_combo():
                    self.word_combo['values'] = word_names
                    if self._is_direct_file_input_active("word"):
                        self._show_direct_file_input("word", self.word_source_path.get().strip(), "Word")
                        self.status_text.set(f"Word {len(word_names)}개 감지됨 (파일 선택 유지)")
                        return
                    if word_names:
                        self.word_combo.current(0)
                        self.selected_word_index.set(1)
                        self.word_doc_name.set(word_info[0][0])
                        self.word_page_count.set(f"{word_info[0][1]}페이지")
                    self.status_text.set(f"Word {len(word_names)}개 감지됨")

                self.root.after(0, update_combo)
            else:
                self.word_list = []
                def clear_combo():
                    self.word_combo.set("")
                    self.word_combo['values'] = []
                    if self._is_direct_file_input_active("word"):
                        self._show_direct_file_input("word", self.word_source_path.get().strip(), "Word")
                        self.status_text.set("Word 파일 선택됨")
                        return
                    self.word_doc_name.set("열린 Word 없음")
                    self.word_page_count.set("-")
                    self.status_text.set("Word를 먼저 열어주세요")
                self.root.after(0, clear_combo)

        except Exception as e:
            expected_not_running = self._is_expected_app_not_running(e, "Word")
            if expected_not_running:
                self.logger.log(f"Word 감지: Word가 아직 실행 중이 아닙니다. ({str(e)[:80]})")
            else:
                self.logger.error("Word 감지 실패", e)
            self.word_list = []
            err_msg = str(e)[:30]
            def show_error():
                self.word_combo.set("")
                if self._is_direct_file_input_active("word"):
                    self._show_direct_file_input("word", self.word_source_path.get().strip(), "Word")
                    self.status_text.set("Word 파일 선택됨")
                    return
                self.word_doc_name.set("열린 Word 없음")
                self.word_page_count.set("-")
                if expected_not_running:
                    self.status_text.set("Word를 먼저 열어주세요")
                else:
                    self.status_text.set(f"Word 감지 실패: {err_msg}")
            self.root.after(0, show_error)

        pythoncom.CoUninitialize()

    def on_word_selected(self, event):
        """Word 콤보박스 선택 이벤트"""
        selected_idx = self.word_combo.current()
        if selected_idx >= 0 and selected_idx < len(self.word_list):
            self._use_open_document_input("word", self.word_source_path, "Word")
            name, page_count, word_index = self.word_list[selected_idx]
            self.selected_word_index.set(word_index)
            self.word_doc_name.set(name)
            self.word_page_count.set(f"{page_count}페이지")
            self.logger.log(f"Word 선택: {name}")

    def start_word_extraction(self):
        """Word 추출 시작"""
        self.logger.log("Word 추출 시작")

        save_path = self.word_save_path.get()
        include_format = self.word_include_format.get()
        use_saveas = self.word_use_saveas.get()
        word_index = self.selected_word_index.get()

        direct_source = self.word_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("word"):
            if self._start_direct_file_conversion(
                "word", direct_source, save_path, self.word_save_path, self.word_extract_button, "Word"
            ):
                return

        if not use_saveas and not HAS_DOCX:
            messagebox.showerror("오류", "python-docx 패키지가 필요합니다.\npip install python-docx")
            return

        if not save_path:
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        if self.word_doc_name.get() == "열린 Word 없음" or not self.word_list:
            messagebox.showwarning("경고", "열린 Word가 없습니다.")
            return

        self.word_extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(
            target=self._extract_word,
            args=(save_path, include_format, use_saveas, word_index),
        )
        thread.daemon = True
        thread.start()

    def _extract_word(self, save_path, include_format, use_saveas, word_index):
        """Word 추출 (백그라운드)"""
        self.logger.log("=== Word 추출 시작 ===")
        extract_start = time.perf_counter()
        pythoncom.CoInitialize()

        try:
            self.root.after(0, lambda: self.status_text.set("Word 연결 중..."))

            # Word 연결 (공통 헬퍼)
            word_app, _ = self._get_word_app()

            doc_count = word_app.Documents.Count
            if doc_count == 0:
                raise Exception("열린 Word 문서가 없습니다. Word에서 문서를 먼저 열어주세요.")

            if word_index > 0 and word_index <= doc_count:
                source_doc = word_app.Documents(word_index)
            else:
                source_doc = word_app.ActiveDocument

            # 대상 문서 유효성 검증 (빈 문서/이름 없는 문서에 SaveAs2 방지)
            try:
                doc_char_count = source_doc.Characters.Count
            except Exception:
                doc_char_count = 0
            if not source_doc.Name or doc_char_count <= 1:
                raise Exception(f"선택된 문서가 비어 있거나 유효하지 않습니다: '{source_doc.Name}'")

            self.logger.log(f"원본 문서: {source_doc.Name} ({doc_char_count}자)")

            if use_saveas:
                try:
                    self.root.after(0, lambda: self.status_text.set("원본 Word 파일 복사 중..."))
                    self.root.after(0, lambda: self.progress_var.set(10))
                    self._copy_word_document_file(source_doc, save_path)
                    self.root.after(0, lambda: self.progress_var.set(100))
                    self.root.after(0, lambda: self.status_text.set("Word 추출 완료! (원본 파일 복사)"))
                    self.root.after(0, lambda: messagebox.showinfo("완료",
                        f"Word 추출 완료 (원본 파일 복사)\n\n{save_path}"))
                    self.logger.log(f"Word 원본 파일 복사 성공: {save_path}")
                    self._log_elapsed("Word 전체 추출 시간", extract_start)
                    return
                except Exception as e:
                    self.logger.log(f"Word 원본 파일 복사 실패: {str(e)[:120]}")
                    openxml_error = None
                    if os.path.splitext(save_path)[1].lower() == ".docx":
                        try:
                            self.root.after(0, lambda: self.status_text.set("Word 원본 구조 복원 중..."))
                            self.root.after(0, lambda: self.progress_var.set(30))
                            self._save_word_openxml_copy(source_doc, save_path)
                            self.root.after(0, lambda: self.progress_var.set(100))
                            self.root.after(0, lambda: self.status_text.set("Word 추출 완료! (원본 구조 복원)"))
                            self.root.after(0, lambda: messagebox.showinfo("완료",
                                f"Word 추출 완료 (원본 구조 복원)\n\n{save_path}"))
                            self._log_elapsed("Word 전체 추출 시간", extract_start)
                            return
                        except Exception as restore_error:
                            openxml_error = restore_error
                            self.logger.log(f"Word WordOpenXML 구조 복원 실패: {str(restore_error)[:120]}")

                    restore_hint = ""
                    if openxml_error is not None:
                        restore_hint = f"\nWordOpenXML 복원 원인: {str(openxml_error)}"
                    raise Exception(
                        "Word 원본 파일 복사/구조 복원에 실패해 중단했습니다.\n"
                        "이미지, 표, 머리글/바닥글, 도형을 보존하려면 원본 파일 복사가 필요합니다.\n\n"
                        "해결 방법:\n"
                        "1. Word에서 문서를 먼저 저장하세요.\n"
                        "2. 원본과 같은 확장자로 저장 경로를 선택하세요.\n"
                        "3. 텍스트만 재구성해도 되는 경우에만 '원본 파일 복사 우선' 체크를 해제하고 .docx로 저장하세요.\n\n"
                        f"원본 복사 원인: {str(e)}{restore_hint}"
                    )

            if not HAS_DOCX:
                raise Exception("python-docx 패키지가 필요합니다. pip install python-docx")

            target_ext = os.path.splitext(save_path)[1].lower()
            if target_ext != ".docx":
                raise Exception(
                    "Word 텍스트 재구성 방식은 .docx 저장만 지원합니다.\n"
                    "원본 그대로 복사하려면 원본과 같은 확장자를 선택하고, "
                    "텍스트 재구성을 사용할 때는 저장 경로를 .docx로 지정하세요."
                )

            # 방법 2: python-docx로 직접 재구성
            self.root.after(0, lambda: self.status_text.set("텍스트 추출 방식으로 진행 중..."))
            self.root.after(0, lambda: self.progress_var.set(20))

            new_doc = DocxDocument()

            total_paragraphs = source_doc.Paragraphs.Count
            self.logger.log(f"총 단락 수: {total_paragraphs}")

            for p_idx in range(1, total_paragraphs + 1):
                progress = 20 + (p_idx / total_paragraphs) * 70
                if p_idx % 50 == 0 or p_idx == total_paragraphs:
                    self.root.after(0, lambda p=progress: self.progress_var.set(p))
                    self.root.after(0, lambda n=p_idx, t=total_paragraphs:
                        self.status_text.set(f"단락 {n}/{t} 처리 중..."))

                try:
                    source_para = source_doc.Paragraphs(p_idx)
                    source_range = source_para.Range

                    # 새 단락 추가
                    new_para = new_doc.add_paragraph()

                    # 단락 서식 복사
                    if include_format:
                        try:
                            # 정렬
                            alignment = source_para.Alignment
                            # 0=Left, 1=Center, 2=Right, 3=Justify
                            align_map = {0: WD_ALIGN_PARAGRAPH.LEFT, 1: WD_ALIGN_PARAGRAPH.CENTER,
                                         2: WD_ALIGN_PARAGRAPH.RIGHT, 3: WD_ALIGN_PARAGRAPH.JUSTIFY}
                            new_para.alignment = align_map.get(alignment, WD_ALIGN_PARAGRAPH.LEFT)
                        except Exception as align_err:
                            self.logger.log(f"  단락 {p_idx} 정렬 복사 실패: {str(align_err)[:40]}")

                        try:
                            # 들여쓰기 (포인트 → EMU 변환)
                            left_indent = source_para.Format.LeftIndent
                            if left_indent and left_indent > 0:
                                new_para.paragraph_format.left_indent = DocxPt(left_indent)
                            first_indent = source_para.Format.FirstLineIndent
                            if first_indent and first_indent > 0:
                                new_para.paragraph_format.first_line_indent = DocxPt(first_indent)
                        except Exception as indent_err:
                            self.logger.log(f"  단락 {p_idx} 들여쓰기 복사 실패: {str(indent_err)[:40]}")

                        try:
                            # 단락 앞/뒤 간격
                            space_before = source_para.Format.SpaceBefore
                            if space_before and space_before > 0:
                                new_para.paragraph_format.space_before = DocxPt(space_before)
                            space_after = source_para.Format.SpaceAfter
                            if space_after and space_after > 0:
                                new_para.paragraph_format.space_after = DocxPt(space_after)
                        except Exception as space_err:
                            self.logger.log(f"  단락 {p_idx} 간격 복사 실패: {str(space_err)[:40]}")

                        try:
                            # 줄 간격
                            line_spacing = source_para.Format.LineSpacing
                            if line_spacing and line_spacing > 0:
                                new_para.paragraph_format.line_spacing = DocxPt(line_spacing)
                        except Exception as ls_err:
                            self.logger.log(f"  단락 {p_idx} 줄간격 복사 실패: {str(ls_err)[:40]}")

                    # 런 단위로 텍스트 복사 — Words(어절) 단위 배치 처리로 COM 왕복 최소화
                    # (기존 문자 단위 순회는 O(n^2) COM 호출로 대용량 문서에서 UI 동결 발생)
                    try:
                        full_text = source_range.Text
                        if full_text and full_text.strip():
                            if not include_format:
                                text = self._clean_xml_text(full_text.rstrip('\r\n\x07\x0d'))
                                if text:
                                    new_para.add_run(text)
                            else:
                                runs_data = self._collect_word_runs(source_range, full_text)
                                if runs_data:
                                    for run_text, fn, fs, b, it, ul, clr in runs_data:
                                        run_text = self._clean_xml_text(run_text.rstrip('\r\n\x07\x0d'))
                                        if not run_text:
                                            continue
                                        run = new_para.add_run(run_text)
                                        try:
                                            if fn:
                                                run.font.name = fn
                                            if fs and fs > 0 and fs < 1000:
                                                run.font.size = DocxPt(fs)
                                            if b is not None and b != 9999999:
                                                run.font.bold = bool(b)
                                            if it is not None and it != 9999999:
                                                run.font.italic = bool(it)
                                            if ul and ul != 0 and ul != 9999999:
                                                run.font.underline = True
                                            if clr and clr != 0 and clr != 9999999:
                                                r_val = clr & 0xFF
                                                g_val = (clr >> 8) & 0xFF
                                                b_val = (clr >> 16) & 0xFF
                                                run.font.color.rgb = DocxRGBColor(r_val, g_val, b_val)
                                        except Exception as fmt_err:
                                            self.logger.log(f"  런 서식 적용 실패: {str(fmt_err)[:50]}")
                                else:
                                    # 런 분석 실패 시 전체 텍스트로 폴백 (서식 없음)
                                    text = self._clean_xml_text(full_text.rstrip('\r\n\x07\x0d'))
                                    if text:
                                        new_para.add_run(text)
                    except Exception as range_err:
                        self.logger.log(f"  단락 {p_idx} Range 접근 실패: {str(range_err)[:50]}")
                        try:
                            text = self._clean_xml_text(source_range.Text.rstrip('\r\n\x07\x0d'))
                            if text:
                                new_para.add_run(text)
                        except Exception:
                            pass

                except Exception as para_err:
                    self.logger.log(f"  단락 {p_idx} 처리 실패: {str(para_err)[:50]}")

            # 저장
            self.root.after(0, lambda: self.status_text.set("파일 저장 중..."))
            self.root.after(0, lambda: self.progress_var.set(95))

            target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
            os.makedirs(target_dir, exist_ok=True)
            new_doc.save(save_path)
            self._validate_office_openxml(save_path, "Word 재구성")
            self.logger.log(f"저장 완료: {save_path}")
            self._log_elapsed("Word 전체 추출 시간", extract_start)

            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("Word 추출 완료! (텍스트 재구성)"))
            self.root.after(0, lambda: messagebox.showinfo("완료",
                f"Word 추출 완료 (텍스트 재구성 방식)\n\n"
                f"{save_path}\n\n"
                f"총 {total_paragraphs}단락\n\n"
                f"⚠️ 이 방식은 서식이 일부 달라질 수 있습니다.\n"
                f"결과를 열어 원본과 비교해 주세요."))

        except Exception as e:
            error_message = str(e)
            self.logger.error("Word 추출 오류", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"추출 중 오류:\n{error_message}"))

        finally:
            self.root.after(0, lambda: self.word_extract_button.config(state=tk.NORMAL))
            pythoncom.CoUninitialize()

    # ========== 메모장 관련 메서드 ==========

    def _setup_notepad_tab(self):
        """메모장 탭 설정"""
        tab = self.notepad_tab

        # 안내 프레임
        info_frame = self._create_section(tab, "메모장 텍스트 추출")

        ttk.Label(info_frame,
                  text="현재 열려있는 메모장 창의 텍스트를 추출합니다.\n"
                       "메모장을 먼저 열고 '감지' 버튼을 누르세요.",
                  font=("맑은 고딕", 9), justify=tk.LEFT).pack(anchor=tk.W, pady=2)

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=5)
        ttk.Label(source_inner, text="파일 선택:", width=12).pack(side=tk.LEFT)
        self.notepad_source_entry = ttk.Entry(source_inner, textvariable=self.notepad_source_path, width=45, state="readonly")
        self.notepad_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_notepad_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # 메모장 선택 콤보박스
        select_frame = ttk.Frame(info_frame)
        select_frame.pack(fill=tk.X, pady=5)
        ttk.Label(select_frame, text="메모장 선택:", width=12).pack(side=tk.LEFT)
        self.notepad_combo = ttk.Combobox(select_frame, state="readonly", width=40)
        self.notepad_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.notepad_combo.bind("<<ComboboxSelected>>", self.on_notepad_selected)

        # 새로고침 버튼
        ttk.Button(info_frame, text="다시 감지", command=lambda: self.detect_open_notepad(prefer_open=True),
                   style="Secondary.TButton").pack(pady=(10, 0))

        # 저장 경로 프레임
        path_frame = self._create_section(tab, "새 파일 저장 위치")

        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.notepad_save_path, width=45).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_notepad_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        # 저장 형식 프레임
        format_frame = self._create_section(tab, "저장 형식")

        self.notepad_save_format = tk.StringVar(value="txt")
        ttk.Radiobutton(format_frame, text="TXT (텍스트 파일)",
                        variable=self.notepad_save_format, value="txt").pack(anchor=tk.W)
        ttk.Radiobutton(format_frame, text="DOCX (Word 파일로 변환)",
                        variable=self.notepad_save_format, value="docx").pack(anchor=tk.W)

        # 추출 버튼
        self.notepad_extract_button = ttk.Button(tab, text="메모장 텍스트 추출",
                                                  command=self.start_notepad_extraction,
                                                  style="Accent.TButton")
        self.notepad_extract_button.pack(pady=10)

    def _setup_pdf_tab(self):
        """PDF 보안 해제 탭 설정"""
        tab = self.pdf_tab

        info_frame = self._create_section(tab, "PDF 보안 해제")
        ttk.Label(
            info_frame,
            text="암호·편집제한이 걸린 PDF를 제한 없는 일반 PDF로 해제해 새 파일로 저장합니다.\n"
                 "정상 PDF는 그대로 복사됩니다. (AES 암호화 PDF는 보안PC 빌드에서 지원되지 않습니다.)",
            font=("맑은 고딕", 9), justify=tk.LEFT,
        ).pack(anchor=tk.W, pady=2)

        source_inner = ttk.Frame(info_frame, style="Card.TFrame")
        source_inner.pack(fill=tk.X, pady=5)
        ttk.Label(source_inner, text="PDF 선택:", width=12).pack(side=tk.LEFT)
        self.pdf_source_entry = ttk.Entry(source_inner, textvariable=self.pdf_source_path,
                                          width=45, state="readonly")
        self.pdf_source_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(source_inner, text="찾아보기", command=self.browse_pdf_source_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        path_frame = self._create_section(tab, "새 파일 저장 위치")
        path_inner = ttk.Frame(path_frame)
        path_inner.pack(fill=tk.X)
        ttk.Entry(path_inner, textvariable=self.pdf_save_path, width=45).pack(
            side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(path_inner, text="찾아보기", command=self.browse_pdf_save_path,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        self.pdf_extract_button = ttk.Button(tab, text="PDF 보안 해제",
                                             command=self.start_pdf_conversion,
                                             style="Accent.TButton")
        self.pdf_extract_button.pack(pady=10)

    def browse_pdf_source_path(self):
        """보안 해제할 PDF 원본 선택"""
        path = filedialog.askopenfilename(
            title="보안 해제할 PDF 선택",
            filetypes=[("PDF 파일", "*.pdf"), ("모든 파일", "*.*")],
        )
        if path:
            self.pdf_source_path.set(path)
            if not self.pdf_save_path.get().strip():
                self.pdf_save_path.set(self._default_direct_save_path(path, "pdf"))
            self.status_text.set("PDF 파일 선택됨")
            self.logger.log(f"PDF 원본 선택: {path}")

    def browse_pdf_save_path(self):
        """PDF 저장 위치 선택"""
        source = self.pdf_source_path.get().strip()
        if source:
            initial = os.path.splitext(os.path.basename(source))[0] + "_복사본.pdf"
        else:
            initial = "해제.pdf"
        path = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF 파일", "*.pdf")],
            initialfile=initial,
            title="저장할 위치 선택",
        )
        if path:
            self.pdf_save_path.set(path)
            self.logger.log(f"PDF 저장 경로: {path}")

    def start_pdf_conversion(self):
        """PDF 보안 해제 실행"""
        source = self.pdf_source_path.get().strip()
        if not source:
            messagebox.showwarning("경고", "보안 해제할 PDF 파일을 선택하세요.")
            return
        if not HAS_PYPDF:
            messagebox.showerror("오류", "PDF 처리 모듈(pypdf)이 없습니다.")
            return
        self._start_direct_file_conversion(
            "pdf", source, self.pdf_save_path.get().strip(),
            self.pdf_save_path, self.pdf_extract_button, "PDF",
        )

    def _setup_batch_tab(self):
        """파일 일괄 변환 탭 설정"""
        tab = self.batch_tab

        file_frame = self._create_section(tab, "변환할 파일")
        ttk.Label(
            file_frame,
            text="PPT, Excel, Word, TXT 파일을 여러 개 추가할 수 있습니다. HWP는 제외됩니다.",
        ).pack(anchor=tk.W, pady=(0, 6))

        list_frame = ttk.Frame(file_frame, style="Card.TFrame")
        list_frame.pack(fill=tk.BOTH, expand=True)
        self.batch_file_listbox = tk.Listbox(
            list_frame,
            height=6,
            selectmode=tk.EXTENDED,
            bg="#ffffff",
            fg=self.ui_colors["text"],
            selectbackground=self.ui_colors["nav_selected_bg"],
            selectforeground=self.ui_colors["nav_selected_fg"],
            relief=tk.FLAT,
            highlightthickness=1,
            highlightbackground=self.ui_colors["border"],
            font=("맑은 고딕", 9),
        )
        self.batch_file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.batch_file_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.batch_file_listbox.configure(yscrollcommand=scrollbar.set)

        button_row = ttk.Frame(file_frame, style="Card.TFrame")
        button_row.pack(fill=tk.X, pady=(8, 0))
        ttk.Button(button_row, text="파일 추가", command=self.add_batch_files,
                   style="Secondary.TButton").pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(button_row, text="폴더 추가", command=self.add_batch_folder,
                   style="Secondary.TButton").pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(button_row, text="선택 제거", command=self.remove_selected_batch_files,
                   style="Secondary.TButton").pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(button_row, text="목록 비우기", command=self.clear_batch_files,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        output_frame = self._create_section(tab, "출력 폴더")
        output_inner = ttk.Frame(output_frame, style="Card.TFrame")
        output_inner.pack(fill=tk.X)
        ttk.Entry(output_inner, textvariable=self.batch_output_dir, width=45).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        ttk.Button(output_inner, text="찾아보기", command=self.browse_batch_output_dir,
                   style="Secondary.TButton").pack(side=tk.LEFT)

        status_frame = self._create_section(tab, "처리 상태")
        status_inner = ttk.Frame(status_frame, style="Card.TFrame")
        status_inner.pack(fill=tk.X)
        ttk.Label(status_inner, textvariable=self.batch_status_text, wraplength=390).pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.batch_extract_button = ttk.Button(status_inner, text="일괄 변환 시작",
                                               command=self.start_batch_conversion,
                                               style="Accent.TButton")
        self.batch_extract_button.pack(side=tk.RIGHT, padx=(10, 0))

    def _batch_file_kind(self, path):
        ext = os.path.splitext(path)[1].lower()
        if ext in {".ppt", ".pptx", ".pptm", ".ppsx", ".potx"}:
            return "ppt"
        if ext in {".xls", ".xlsx", ".xlsm", ".xlsb"}:
            return "excel"
        if ext in {".doc", ".docx", ".docm"}:
            return "word"
        if ext in {".pdf"}:
            return "pdf"
        if ext in {".hwp", ".hwpx"}:
            return "hwp"
        if ext in {".txt"}:
            return "text"
        return None

    def _batch_target_extension(self, path, kind):
        source_ext = os.path.splitext(path)[1].lower()
        if kind == "ppt":
            return source_ext if source_ext in {".pptx", ".pptm", ".ppsx", ".potx"} else ".pptx"
        if kind == "excel":
            return source_ext if source_ext in {".xlsx", ".xlsm", ".xlsb"} else ".xlsx"
        if kind == "word":
            return ".docx"
        if kind == "pdf":
            return ".pdf"
        if kind == "hwp":
            return source_ext if source_ext in {".hwp", ".hwpx"} else ".hwp"
        return ".txt"

    def _make_unique_output_path(self, output_dir, source_path, kind):
        stem = os.path.splitext(os.path.basename(source_path))[0]
        ext = self._batch_target_extension(source_path, kind)
        candidate = os.path.join(output_dir, f"{stem}_복사본{ext}")
        if not os.path.exists(candidate):
            return candidate
        for index in range(2, 1000):
            candidate = os.path.join(output_dir, f"{stem}_복사본_{index}{ext}")
            if not os.path.exists(candidate):
                return candidate
        raise Exception(f"출력 파일명을 만들 수 없습니다: {source_path}")

    def _add_batch_paths(self, paths):
        added = 0
        added_paths = []
        seen = {os.path.abspath(path).lower() for path in self.batch_files}
        for path in self._expand_supported_drop_paths(paths):
            key = os.path.abspath(path).lower()
            if key in seen:
                continue
            self.batch_files.append(path)
            seen.add(key)
            added += 1
            added_paths.append(path)
        self._refresh_batch_file_list()
        self.batch_status_text.set(f"{added}개 파일 추가됨, 총 {len(self.batch_files)}개")
        if added_paths:
            # 첫 변환의 콜드 스타트까지 없애기 위해 필요한 Office를 백그라운드로 미리 띄운다.
            self._prewarm_batch_office(added_paths)
        return added

    def _refresh_batch_file_list(self):
        self.batch_file_listbox.delete(0, tk.END)
        for index, path in enumerate(self.batch_files, start=1):
            kind = self._batch_file_kind(path) or "skip"
            self.batch_file_listbox.insert(tk.END, f"{index}. [{kind.upper()}] {path}")

    # 일괄 변환 Office 형식별 prog_id 매핑 (예열·재사용 공통)
    _OFFICE_PROG_IDS = {
        "ppt": ("PowerPoint.Application", "PowerPoint"),
        "excel": ("Excel.Application", "Excel"),
        "word": ("Word.Application", "Word"),
    }

    def _prewarm_batch_office(self, paths):
        """추가된 파일에 맞춰 Office를 백그라운드로 미리 띄운다. 실제 판단·실행은
        워커 스레드에서 수행해 GUI가 멈추지 않게 한다."""
        if not HAS_WIN32COM:
            return
        self._ensure_office_worker()
        self._office_job_queue.put(lambda ps=list(paths): self._prewarm_office_job(ps))

    def _prewarm_office_job(self, paths):
        """워커 스레드에서 실행: 정상 OpenXML은 건너뛰고, Office 복원이 필요한
        형식(zip이 아닌 DRM/보안 컨테이너 등)만 골라 예열 인스턴스를 띄운다."""
        needed = set()
        for path in paths:
            kind = self._batch_file_kind(path)
            if kind not in self._OFFICE_PROG_IDS or kind in needed:
                continue
            try:
                # 이미 정상 OpenXML이면 Office 없이 빠른 복사가 가능하므로 예열 불필요.
                self._validate_office_openxml(path, "예열 점검", deep=False)
            except Exception:
                needed.add(kind)
            if len(needed) == len(self._OFFICE_PROG_IDS):
                break
        for kind in sorted(needed):
            prog_id, display_name = self._OFFICE_PROG_IDS[kind]
            try:
                self._acquire_warm_office_app(kind, prog_id, display_name)
            except Exception as exc:
                self.logger.log(f"Office 예열 실패({kind}): {str(exc)[:80]}")
        if needed:
            self.logger.log(f"일괄 변환 Office 예열 완료: {', '.join(sorted(needed))}")

    def add_batch_files(self):
        paths = filedialog.askopenfilenames(
            title="일괄 변환할 파일 선택",
            filetypes=[
                ("지원 문서", "*.ppt;*.pptx;*.pptm;*.ppsx;*.potx;*.xls;*.xlsx;*.xlsm;*.xlsb;*.doc;*.docx;*.docm;*.pdf;*.txt"),
                ("PowerPoint", "*.ppt;*.pptx;*.pptm;*.ppsx;*.potx"),
                ("Excel", "*.xls;*.xlsx;*.xlsm;*.xlsb"),
                ("Word", "*.doc;*.docx;*.docm"),
                ("PDF", "*.pdf"),
                ("텍스트", "*.txt"),
                ("모든 파일", "*.*"),
            ],
        )
        if paths:
            self._add_batch_paths(paths)

    def add_batch_folder(self):
        folder = filedialog.askdirectory(title="일괄 변환할 폴더 선택")
        if not folder:
            return
        paths = []
        for root_dir, _dirs, files in os.walk(folder):
            for filename in files:
                path = os.path.join(root_dir, filename)
                if self._batch_file_kind(path):
                    paths.append(path)
        self._add_batch_paths(sorted(paths))

    def remove_selected_batch_files(self):
        selected = set(self.batch_file_listbox.curselection())
        if not selected:
            return
        self.batch_files = [path for index, path in enumerate(self.batch_files) if index not in selected]
        self._refresh_batch_file_list()
        self.batch_status_text.set(f"선택 파일 제거 완료, 총 {len(self.batch_files)}개")

    def clear_batch_files(self):
        self.batch_files = []
        self._refresh_batch_file_list()
        self.batch_status_text.set("파일 목록을 비웠습니다.")

    def browse_batch_output_dir(self):
        folder = filedialog.askdirectory(title="일괄 변환 결과를 저장할 폴더 선택")
        if folder:
            self.batch_output_dir.set(folder)
            self.batch_status_text.set(f"출력 폴더: {folder}")

    def start_batch_conversion(self):
        self.logger.log("일괄 변환 시작 버튼 클릭")
        if not self.batch_files:
            messagebox.showwarning("경고", "일괄 변환할 파일을 추가해주세요.")
            return

        output_dir = self.batch_output_dir.get().strip()
        if not output_dir:
            messagebox.showwarning("경고", "출력 폴더를 선택해주세요.")
            return

        os.makedirs(output_dir, exist_ok=True)
        files = list(self.batch_files)
        self.batch_extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)
        # 예열·재사용 워커에 작업을 맡겨, 변환할 때마다 Office를 새로 켜는 비용을 없앤다.
        self._ensure_office_worker()
        self._office_job_queue.put(lambda: self._extract_batch(files, output_dir))

    def _ensure_office_worker(self):
        """Office 변환을 처리하는 영속 워커 스레드를 보장한다.
        한 스레드가 COM 아파트먼트와 모든 Office 인스턴스를 소유해야
        변환 사이에 인스턴스를 안전하게 재사용할 수 있다."""
        worker = self._office_worker
        if worker is not None and worker.is_alive():
            return
        worker = threading.Thread(
            target=self._office_worker_loop, name="OfficeWorker", daemon=True
        )
        self._office_worker = worker
        worker.start()

    def _office_worker_loop(self):
        """큐에 들어온 변환 작업을 순서대로 처리하고, 예열 Office 인스턴스를 보관한다."""
        com_ready = False
        if HAS_WIN32COM:
            try:
                pythoncom.CoInitialize()
                com_ready = True
            except Exception as exc:
                self.logger.log(f"Office 워커 COM 초기화 실패: {str(exc)[:80]}")
        try:
            while True:
                job = self._office_job_queue.get()
                try:
                    if job is None:  # 종료 신호
                        break
                    job()
                except Exception as exc:
                    self.logger.error("Office 워커 작업 오류", exc)
                finally:
                    self._office_job_queue.task_done()
        finally:
            for key in list(self._warm_office_apps.keys()):
                app = self._warm_office_apps.pop(key, None)
                if app is None:
                    continue
                try:
                    app.Quit()
                except Exception:
                    pass
            if com_ready:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

    def _acquire_warm_office_app(self, key, prog_id, display_name):
        """예열된 격리 Office 인스턴스를 재사용한다. 죽었으면 다시 만든다.
        반드시 Office 워커 스레드 안에서만 호출해야 COM 아파트먼트가 일치한다."""
        app = self._warm_office_apps.get(key)
        if app is not None:
            try:
                _ = app.Name  # 살아있는지 가벼운 속성 접근으로 확인
                return app
            except Exception as exc:
                self.logger.log(
                    f"일괄 변환 {display_name} 예열 인스턴스 만료, 재생성: {str(exc)[:60]}"
                )
                self._warm_office_apps.pop(key, None)

        app, _created = self._create_isolated_com_app(prog_id, display_name)
        # 변환 전용 인스턴스의 불필요한 렌더링/경고를 줄여 속도를 높인다.
        if key == "ppt":
            try:
                app.DisplayAlerts = 1
            except Exception:
                pass
            try:
                app.WindowState = 2  # ppWindowMinimized: 창 렌더링 최소화
            except Exception as exc:
                self.logger.log(f"PowerPoint 창 최소화 실패(무시): {str(exc)[:60]}")
        elif key == "excel":
            try:
                app.DisplayAlerts = False
            except Exception:
                pass
        self._warm_office_apps[key] = app
        self.logger.log(f"일괄 변환 {display_name} 예열 인스턴스 준비 완료")
        return app

    def _shutdown_office_worker(self):
        """프로그램 종료 시 예열 Office 인스턴스를 정리하고 워커를 종료한다."""
        worker = self._office_worker
        if worker is None or not worker.is_alive():
            return
        try:
            self._office_job_queue.put(None)
            worker.join(timeout=10)
        except Exception:
            pass
        self._office_worker = None

    def _batch_convert_ppt_file(self, ppt_app, source_path, target_path, skip_direct=False):
        source_pres = None
        try:
            if not skip_direct:
                try:
                    self._publish_existing_verified_file(source_path, target_path, "PPT 일괄")
                    return
                except Exception as direct_copy_error:
                    self.logger.log(
                        f"  PPT 파일 직접 복사 불가, PowerPoint 내부 복원 시도: {str(direct_copy_error)[:120]}"
                    )

            self._close_office_modal_dialogs(ppt_app, "PowerPoint 일괄")
            source_pres = ppt_app.Presentations.Open(source_path, True, False, False)
            self.logger.log(f"  PPT 열기 완료: {source_pres.Name}")
            try:
                self._save_ppt_clipboard_package_copy(source_pres, target_path)
                return
            except Exception as package_error:
                self.logger.log(f"  PPT 일괄 클립보드 패키지 실패, 슬라이드 복제 시도: {str(package_error)[:120]}")
            self._save_ppt_slide_clone(source_pres, target_path)
        finally:
            try:
                self._close_office_modal_dialogs(ppt_app, "PowerPoint 일괄")
            except Exception:
                pass
            if source_pres is not None:
                try:
                    source_pres.Close()
                except Exception:
                    pass
            try:
                self._close_office_modal_dialogs(ppt_app, "PowerPoint 일괄")
            except Exception:
                pass

    def _save_excel_as_openxml_copy(self, source_wb, target_path):
        target_ext = os.path.splitext(target_path)[1].lower()
        file_format = {
            ".xlsx": 51,
            ".xlsm": 52,
            ".xlsb": 50,
        }.get(target_ext, 51)
        temp_path = self._make_local_temp_path(target_ext if target_ext else ".xlsx")
        try:
            self._run_with_heartbeat(
                "Excel 일괄 SaveAs",
                lambda: source_wb.SaveAs(temp_path, FileFormat=file_format),
            )
            self._publish_verified_file(temp_path, target_path, "Excel 일괄 SaveAs")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _set_excel_conversion_options(self, excel_app):
        """일괄 변환용 Excel 인스턴스에서 불필요한 지연 요소를 끈다."""
        options = {
            "AskToUpdateLinks": False,
            "EnableEvents": False,
            "ScreenUpdating": False,
            "DisplayAlerts": False,
            "AutomationSecurity": 3,  # msoAutomationSecurityForceDisable
            "Calculation": -4135,  # xlCalculationManual
        }
        for name, value in options.items():
            try:
                setattr(excel_app, name, value)
            except Exception as exc:
                self.logger.log(f"Excel 변환 옵션 설정 실패({name}): {str(exc)[:60]}")

    def _open_excel_workbook_for_conversion(self, excel_app, source_path):
        self._set_excel_conversion_options(excel_app)
        try:
            return excel_app.Workbooks.Open(
                Filename=source_path,
                UpdateLinks=0,
                ReadOnly=True,
                IgnoreReadOnlyRecommended=True,
                Notify=False,
                AddToMru=False,
                Local=True,
                CorruptLoad=0,
            )
        except Exception as exc:
            message = str(exc)
            if "매개 변수" in message or "parameter" in message.lower():
                return excel_app.Workbooks.Open(source_path, 0, True)
            raise

    def _batch_convert_excel_file(self, excel_app, source_path, target_path, skip_direct=False):
        source_wb = None
        try:
            if not skip_direct:
                try:
                    self._try_existing_office_file_copy(source_path, target_path, "Excel 일괄")
                    return
                except Exception as direct_copy_error:
                    self.logger.log(
                        f"  Excel 일괄 직접 복사 불가, Excel 열기 시도: {str(direct_copy_error)[:120]}"
                    )
            source_wb = self._run_with_heartbeat(
                "Excel 파일 열기",
                lambda: self._open_excel_workbook_for_conversion(excel_app, source_path),
            )
            self.logger.log(f"  Excel 열기 완료: {source_wb.Name}")
            try:
                self._save_native_copy(source_wb, target_path, "Excel 일괄")
                return
            except Exception as copy_error:
                self.logger.log(f"  Excel 일괄 원본 복사 실패, SaveAs 변환 시도: {str(copy_error)[:120]}")
            self._save_excel_as_openxml_copy(source_wb, target_path)
        finally:
            if source_wb is not None:
                try:
                    source_wb.Close(False)
                except Exception:
                    pass

    def _save_word_as_docx_copy(self, source_doc, target_path):
        temp_path = self._make_local_temp_path(".docx")
        try:
            self._run_with_heartbeat(
                "Word 일괄 SaveAs2",
                lambda: source_doc.SaveAs2(temp_path, FileFormat=16),
            )
            self._publish_verified_file(temp_path, target_path, "Word 일괄 SaveAs2")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _batch_convert_word_file(self, word_app, source_path, target_path):
        source_doc = None
        try:
            source_doc = word_app.Documents.Open(source_path, ReadOnly=True, AddToRecentFiles=False, Visible=False)
            self.logger.log(f"  Word 열기 완료: {source_doc.Name}")
            try:
                if os.path.splitext(source_path)[1].lower() == os.path.splitext(target_path)[1].lower():
                    self._copy_word_document_file(source_doc, target_path)
                    return
            except Exception as copy_error:
                self.logger.log(f"  Word 일괄 원본 복사 실패, 구조 복원 시도: {str(copy_error)[:120]}")
            try:
                self._save_word_openxml_copy(source_doc, target_path)
                return
            except Exception as openxml_error:
                self.logger.log(f"  Word 일괄 WordOpenXML 실패, SaveAs2 시도: {str(openxml_error)[:120]}")
            self._save_word_as_docx_copy(source_doc, target_path)
        finally:
            if source_doc is not None:
                try:
                    source_doc.Close(False)
                except Exception:
                    pass

    def _batch_convert_text_file(self, source_path, target_path):
        target_dir = os.path.dirname(os.path.abspath(target_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        shutil.copy2(source_path, target_path)
        if not os.path.exists(target_path) or os.path.getsize(target_path) <= 0:
            raise Exception("TXT 복사 결과 파일이 없거나 비어 있습니다.")

    def _looks_like_pdf(self, path):
        """파일 헤더로 정상 PDF 여부를 확인한다(%PDF-)."""
        try:
            with open(path, "rb") as handle:
                return handle.read(5).startswith(b"%PDF-")
        except Exception:
            return False

    def _validate_pdf(self, path, label):
        if not os.path.exists(path) or os.path.getsize(path) <= 0:
            raise Exception(f"{label} 결과 파일이 없거나 비어 있습니다.")
        if not self._looks_like_pdf(path):
            raise Exception(f"{label} 결과가 정상 PDF 파일이 아닙니다.")

    def _publish_pdf_copy(self, source_path, save_path, label):
        """이미 해제된(암호화 안 된) 정상 PDF는 그대로 복사한다."""
        if os.path.abspath(source_path).lower() == os.path.abspath(save_path).lower():
            raise Exception("원본과 같은 경로로는 복사할 수 없습니다.")
        target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
        os.makedirs(target_dir, exist_ok=True)
        fd, stage_path = tempfile.mkstemp(prefix=".docextract_", suffix=".pdf", dir=target_dir)
        os.close(fd)
        try:
            shutil.copyfile(source_path, stage_path)
            os.replace(stage_path, save_path)
            self._validate_copied_file(source_path, save_path, f"{label} 최종 파일")
        except Exception:
            if os.path.exists(stage_path):
                try:
                    os.remove(stage_path)
                except Exception:
                    pass
            raise

    def _release_pdf_file(self, source_path, save_path, label):
        """암호/편집제한 PDF를 해제(복호화·제한제거)해 일반 PDF로 저장한다."""
        reader = PdfReader(source_path)
        if reader.is_encrypted:
            decrypted = False
            try:
                decrypted = bool(reader.decrypt(""))  # 빈 사용자 암호(편집제한 PDF 대부분)
            except Exception as exc:
                self.logger.log(f"  PDF 복호화 시도 실패: {str(exc)[:80]}")
            if not decrypted:
                raise Exception(
                    "암호로 보호된 PDF라 자동 해제가 불가합니다(사용자 암호 필요)."
                )
        writer = PdfWriter()
        writer.append(reader)
        try:
            if reader.metadata:
                writer.add_metadata(reader.metadata)
        except Exception:
            pass

        temp_path = self._make_local_temp_path(".pdf")
        try:
            with open(temp_path, "wb") as handle:
                writer.write(handle)
            self._validate_pdf(temp_path, f"{label} 해제 PDF")

            target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
            os.makedirs(target_dir, exist_ok=True)
            fd, stage_path = tempfile.mkstemp(prefix=".docextract_", suffix=".pdf", dir=target_dir)
            os.close(fd)
            try:
                shutil.copyfile(temp_path, stage_path)
                os.replace(stage_path, save_path)
            except Exception:
                if os.path.exists(stage_path):
                    try:
                        os.remove(stage_path)
                    except Exception:
                        pass
                raise
            self._validate_pdf(save_path, f"{label} 최종 PDF")
            self.logger.log(f"{label} 보안 해제 저장 완료: {save_path}")
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def _convert_pdf_file(self, source_path, save_path):
        """PDF 보안 해제: 정상 PDF는 복사, 암호/편집제한 PDF는 해제 후 저장."""
        label = "PDF"
        if not self._looks_like_pdf(source_path):
            raise Exception(
                "정상 PDF 파일이 아닙니다(DRM 보안 컨테이너일 수 있음). "
                "원본을 인가된 뷰어로 연 뒤 'Microsoft Print to PDF'로 다시 저장하세요."
            )
        if not HAS_PYPDF:
            raise Exception("PDF 처리에는 pypdf 패키지가 필요합니다(pip install pypdf).")

        try:
            reader = PdfReader(source_path)
            encrypted = bool(reader.is_encrypted)
        except Exception as exc:
            raise Exception(f"PDF를 열 수 없습니다: {str(exc)[:120]}")

        if not encrypted:
            try:
                self._publish_pdf_copy(source_path, save_path, label)
                return
            except Exception as copy_error:
                self.logger.log(f"  PDF 직접 복사 실패, 재구성 저장 시도: {str(copy_error)[:120]}")

        self._release_pdf_file(source_path, save_path, label)

    def _extract_batch(self, files, output_dir):
        self.logger.log("=== 일괄 변환 시작 ===")
        extract_start = time.perf_counter()
        successes = []
        failures = []

        try:
            total = len(files)
            for index, source_path in enumerate(files, start=1):
                kind = self._batch_file_kind(source_path)
                progress = ((index - 1) / total) * 95
                self.root.after(0, lambda p=progress: self.progress_var.set(p))
                self.root.after(0, lambda i=index, t=total: self.status_text.set(f"일괄 변환 중... {i}/{t}"))

                if not kind:
                    failures.append((source_path, "지원하지 않는 확장자"))
                    continue

                target_path = self._make_unique_output_path(output_dir, source_path, kind)
                self.logger.log(f"[{index}/{total}] 일괄 변환: {source_path} -> {target_path}")
                self.root.after(0, lambda p=source_path: self.batch_status_text.set(f"처리 중: {os.path.basename(p)}"))

                try:
                    if kind == "ppt":
                        try:
                            self._publish_existing_verified_file(source_path, target_path, "PPT 일괄")
                            successes.append(target_path)
                            self.logger.log(f"  일괄 변환 완료: {target_path}")
                            continue
                        except Exception as direct_copy_error:
                            self.logger.log(
                                f"  PPT 일괄 직접 복사 불가, PowerPoint 내부 복원 시도: {str(direct_copy_error)[:120]}"
                            )
                        if not HAS_WIN32COM:
                            raise Exception("PPT 내부 복원에는 pywin32/win32com이 필요합니다.")
                        ppt_app = self._acquire_warm_office_app("ppt", "PowerPoint.Application", "PowerPoint")
                        self._batch_convert_ppt_file(ppt_app, source_path, target_path, skip_direct=True)
                    elif kind == "excel":
                        try:
                            self._try_existing_office_file_copy(source_path, target_path, "Excel 일괄")
                            successes.append(target_path)
                            self.logger.log(f"  일괄 변환 완료: {target_path}")
                            continue
                        except Exception as direct_copy_error:
                            self.logger.log(
                                f"  Excel 일괄 직접 복사 불가, Excel 내부 복원 시도: {str(direct_copy_error)[:120]}"
                            )
                        if not HAS_WIN32COM:
                            raise Exception("Excel 일괄 변환에는 pywin32/win32com이 필요합니다.")
                        excel_app = self._acquire_warm_office_app("excel", "Excel.Application", "Excel")
                        self._batch_convert_excel_file(excel_app, source_path, target_path, skip_direct=True)
                    elif kind == "word":
                        try:
                            self._try_existing_office_file_copy(source_path, target_path, "Word 일괄")
                            successes.append(target_path)
                            self.logger.log(f"  일괄 변환 완료: {target_path}")
                            continue
                        except Exception as direct_copy_error:
                            self.logger.log(
                                f"  Word 일괄 직접 복사 불가, Word 내부 복원 시도: {str(direct_copy_error)[:120]}"
                            )
                        if not HAS_WIN32COM:
                            raise Exception("Word 일괄 변환에는 pywin32/win32com이 필요합니다.")
                        word_app = self._acquire_warm_office_app("word", "Word.Application", "Word")
                        self._batch_convert_word_file(word_app, source_path, target_path)
                    elif kind == "pdf":
                        self._convert_pdf_file(source_path, target_path)
                    elif kind == "hwp":
                        self._convert_hwp_file(source_path, target_path)
                    elif kind == "text":
                        self._batch_convert_text_file(source_path, target_path)
                    successes.append(target_path)
                    self.logger.log(f"  일괄 변환 완료: {target_path}")
                except Exception as item_error:
                    failures.append((source_path, str(item_error)))
                    self.logger.error(f"일괄 변환 실패: {source_path}", item_error)

            self.root.after(0, lambda: self.progress_var.set(100))
            self._log_elapsed("일괄 변환 전체 시간", extract_start)
            summary = f"일괄 변환 완료: 성공 {len(successes)}개, 실패 {len(failures)}개"
            if failures:
                summary += f"\n첫 실패: {os.path.basename(failures[0][0])} - {failures[0][1][:80]}"
            self.logger.log(summary.replace("\n", " / "))
            self.root.after(0, lambda: self.status_text.set(summary.split("\n")[0]))
            self.root.after(0, lambda: self.batch_status_text.set(summary))
            self.root.after(0, lambda: messagebox.showinfo("일괄 변환 완료", summary))
        except Exception as error:
            error_message = str(error)
            self.logger.error("일괄 변환 오류", error)
            self.root.after(0, lambda: self.status_text.set(f"일괄 변환 오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"일괄 변환 중 오류:\n{error_message}"))
        finally:
            # 예열 인스턴스는 종료하지 않고 살려둬 다음 변환에서 즉시 재사용한다.
            # (실제 종료는 _shutdown_office_worker / 프로그램 종료 시 수행)
            self.root.after(0, lambda: self.batch_extract_button.config(state=tk.NORMAL))

    def browse_notepad_save_path(self):
        """메모장 저장 경로 선택"""
        self.logger.log("메모장 저장 경로 선택")
        save_format = self.notepad_save_format.get()

        direct_source = self.notepad_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("text"):
            doc_name = os.path.basename(direct_source)
        else:
            doc_name = self.notepad_doc_name.get()
        if doc_name and doc_name != "감지 중..." and doc_name != "열린 메모장 없음":
            base_name = os.path.splitext(doc_name)[0] + "_복사본"
        else:
            base_name = "새문서"

        if save_format == "docx":
            ext = ".docx"
            filetypes = [("Word 파일", "*.docx")]
        else:
            ext = ".txt"
            filetypes = [("텍스트 파일", "*.txt")]

        path = filedialog.asksaveasfilename(
            defaultextension=ext,
            filetypes=filetypes,
            initialfile=base_name + ext,
            title="저장할 위치 선택"
        )
        if path:
            self.notepad_save_path.set(path)
            self.logger.log(f"메모장 저장 경로: {path}")

    def detect_open_notepad(self, prefer_open=False):
        """열려있는 메모장 감지"""
        if prefer_open:
            self._use_open_document_input("text", self.notepad_source_path, "메모장")
        self.logger.log("메모장 감지 시작")
        if self._is_direct_file_input_active("text"):
            self.status_text.set("메모장 감지 중... (파일 선택 유지)")
            self._show_direct_file_input("text", self.notepad_source_path.get().strip(), "TXT")
        else:
            self.status_text.set("메모장 감지 중...")
            self.notepad_doc_name.set("감지 중...")

        thread = threading.Thread(target=self._detect_notepad)
        thread.daemon = True
        thread.start()

    def _detect_notepad(self):
        """메모장 감지 (백그라운드) - Win32 API로 메모장 창 찾기.

        순수 Win32 ctypes만 사용하므로 COM 초기화 불필요.
        """
        self.logger.log("백그라운드 메모장 감지 스레드 시작")

        try:
            user32 = ctypes.windll.user32

            notepad_windows = []

            # EnumWindows 콜백
            def enum_callback(hwnd, lparam):
                if user32.IsWindowVisible(hwnd):
                    if self._is_notepad_window(hwnd):
                        title = self._get_window_title(hwnd)
                        class_name = self._get_window_class_name(hwnd)
                        notepad_windows.append((hwnd, title))
                        self.logger.log(
                            f"  메모장 발견: hwnd={hwnd}, class='{class_name}', 제목='{title}'"
                        )
                return True

            WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, wintypes.HWND, wintypes.LPARAM)
            user32.EnumWindows(WNDENUMPROC(enum_callback), 0)

            if notepad_windows:
                self.notepad_list = notepad_windows
                names = [title if title else "제목 없음" for _, title in notepad_windows]

                def update_combo():
                    self.notepad_combo['values'] = names
                    if self._is_direct_file_input_active("text"):
                        self._show_direct_file_input("text", self.notepad_source_path.get().strip(), "TXT")
                        self.status_text.set(f"메모장 {len(names)}개 감지됨 (파일 선택 유지)")
                        return
                    self.notepad_combo.current(0)
                    self.notepad_doc_name.set(names[0])
                    self.status_text.set(f"메모장 {len(names)}개 감지됨")

                self.root.after(0, update_combo)
            else:
                self.notepad_list = []
                def clear_combo():
                    self.notepad_combo.set("")
                    self.notepad_combo['values'] = []
                    if self._is_direct_file_input_active("text"):
                        self._show_direct_file_input("text", self.notepad_source_path.get().strip(), "TXT")
                        self.status_text.set("TXT 파일 선택됨")
                        return
                    self.notepad_doc_name.set("열린 메모장 없음")
                    self.status_text.set("메모장을 먼저 열어주세요")
                self.root.after(0, clear_combo)

        except Exception as e:
            self.logger.error("메모장 감지 실패", e)
            self.notepad_list = []
            err_msg = str(e)[:30]
            def show_error():
                self.notepad_combo.set("")
                if self._is_direct_file_input_active("text"):
                    self._show_direct_file_input("text", self.notepad_source_path.get().strip(), "TXT")
                    self.status_text.set("TXT 파일 선택됨")
                    return
                self.notepad_doc_name.set("열린 메모장 없음")
                self.status_text.set(f"메모장 감지 실패: {err_msg}")
            self.root.after(0, show_error)

    def on_notepad_selected(self, event):
        """메모장 콤보박스 선택 이벤트"""
        selected_idx = self.notepad_combo.current()
        if selected_idx >= 0 and selected_idx < len(self.notepad_list):
            self._use_open_document_input("text", self.notepad_source_path, "메모장")
            hwnd, title = self.notepad_list[selected_idx]
            self.notepad_doc_name.set(title if title else "제목 없음")
            self.logger.log(f"메모장 선택: {title} (hwnd={hwnd})")

    def start_notepad_extraction(self):
        """메모장 추출 시작"""
        self.logger.log("메모장 추출 시작")

        save_path = self.notepad_save_path.get()
        save_format = self.notepad_save_format.get()
        selected_idx = self.notepad_combo.current()

        direct_source = self.notepad_source_path.get().strip()
        if direct_source and self._is_direct_file_input_active("text"):
            if not save_path:
                preferred_ext = ".docx" if save_format == "docx" else ".txt"
                save_path = self._default_direct_save_path(direct_source, "text", preferred_ext)
                self.notepad_save_path.set(save_path)
            if self._start_direct_file_conversion(
                "text", direct_source, save_path, self.notepad_save_path, self.notepad_extract_button, "TXT"
            ):
                return

        if not save_path:
            messagebox.showwarning("경고", "저장 경로를 선택해주세요.")
            return

        if self.notepad_doc_name.get() == "열린 메모장 없음" or not self.notepad_list:
            messagebox.showwarning("경고", "열린 메모장이 없습니다.")
            return

        if selected_idx < 0 or selected_idx >= len(self.notepad_list):
            messagebox.showwarning("경고", "메모장을 선택해주세요.")
            return

        hwnd, title = self.notepad_list[selected_idx]

        self.notepad_extract_button.config(state=tk.DISABLED)
        self.progress_var.set(0)

        thread = threading.Thread(target=self._extract_notepad, args=(save_path, save_format, hwnd, title))
        thread.daemon = True
        thread.start()

    def _extract_notepad(self, save_path, save_format, hwnd, title):
        """메모장 추출 (백그라운드) - Win32 API로 텍스트 읽기.

        순수 Win32 ctypes만 사용하므로 COM 초기화 불필요.
        """
        self.logger.log("=== 메모장 추출 시작 ===")
        extract_start = time.perf_counter()

        try:
            user32 = ctypes.windll.user32

            self.logger.log(f"대상 메모장: {title} (hwnd={hwnd})")
            if not user32.IsWindow(hwnd):
                raise Exception("선택한 메모장 창이 닫혔습니다. 다시 감지 후 선택해주세요.")

            self.root.after(0, lambda: self.status_text.set("메모장 텍스트 읽는 중..."))
            self.root.after(0, lambda: self.progress_var.set(20))

            # 메모장 내부 텍스트 컨트롤 찾기
            edit_hwnd = self._find_child_window_by_classes(
                hwnd,
                ["Edit", "RichEditD2DPT", "RICHEDIT50W"],
            )
            if not edit_hwnd:
                # Windows 11 새 메모장(UWP/WinUI)은 이 방식으로 접근할 수 없음 — 사용자 친화 안내
                raise Exception(
                    "Windows 11 기본 메모장(새 버전)은 자동 추출이 지원되지 않습니다.\n\n"
                    "해결 방법:\n"
                    "1. 메모장에서 Ctrl+A 로 전체 선택\n"
                    "2. Ctrl+C 로 복사\n"
                    "3. 새 메모장(또는 Word)에 Ctrl+V 로 붙여넣고 저장\n\n"
                    "또는 메모장에서 '파일 → 다른 이름으로 저장'을 직접 사용하세요."
                )

            self.logger.log(f"Edit 컨트롤 hwnd: {edit_hwnd}")

            # WM_GETTEXTLENGTH, WM_GETTEXT 로 텍스트 가져오기
            WM_GETTEXTLENGTH = 0x000E
            WM_GETTEXT = 0x000D

            text_length = user32.SendMessageW(edit_hwnd, WM_GETTEXTLENGTH, 0, 0)
            self.logger.log(f"텍스트 길이: {text_length}")

            if text_length <= 0:
                raise Exception("메모장에 텍스트가 없습니다.")

            self.root.after(0, lambda: self.progress_var.set(50))

            # 사용자가 읽는 사이 타이핑해 텍스트가 늘어날 수 있으므로 여유 버퍼(+16) 할당
            buffer_size = text_length + 16
            buffer = ctypes.create_unicode_buffer(buffer_size)
            copied = user32.SendMessageW(edit_hwnd, WM_GETTEXT, buffer_size, buffer)
            self.logger.log(f"WM_GETTEXT 반환값(실제 복사 문자 수): {copied}")

            # SendMessageW 반환값 검증 — 창이 응답 거부/행잉이면 0 또는 음수
            if copied <= 0:
                raise Exception(
                    "메모장 창에서 텍스트를 읽어올 수 없습니다.\n"
                    "창이 응답하지 않거나 닫혔을 수 있습니다. 다시 시도해 주세요."
                )

            text = buffer.value
            # 버퍼 안 읽힘 경계 검증 (비정상 절단 감지)
            if len(text) < copied:
                self.logger.log(f"  경고: 버퍼 길이({len(text)}) < 반환값({copied}) — 데이터 절단 가능")

            self.logger.log(f"텍스트 추출 성공: {len(text)}글자")
            self.root.after(0, lambda: self.progress_var.set(70))

            # 저장
            target_dir = os.path.dirname(os.path.abspath(save_path)) or os.getcwd()
            os.makedirs(target_dir, exist_ok=True)

            if save_format == "docx":
                if not HAS_DOCX:
                    raise Exception("DOCX 저장에는 python-docx 패키지가 필요합니다.")
                self.root.after(0, lambda: self.status_text.set("DOCX 파일로 저장 중..."))
                new_doc = DocxDocument()
                for line in text.split('\n'):
                    line = self._clean_xml_text(line.rstrip('\r'))
                    new_doc.add_paragraph(line)
                new_doc.save(save_path)
                self._validate_office_openxml(save_path, "메모장 DOCX")
            else:
                self.root.after(0, lambda: self.status_text.set("TXT 파일로 저장 중..."))
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                if not os.path.exists(save_path) or os.path.getsize(save_path) <= 0:
                    raise Exception("메모장 TXT 저장 결과 파일이 없거나 비어 있습니다.")

            self.logger.log(f"저장 완료: {save_path}")
            self._log_elapsed("메모장 전체 추출 시간", extract_start)
            self.root.after(0, lambda: self.progress_var.set(100))
            self.root.after(0, lambda: self.status_text.set("메모장 추출 완료!"))
            self.root.after(0, lambda: messagebox.showinfo("완료",
                f"메모장 추출 완료!\n{save_path}\n\n{len(text)}글자"))

        except Exception as e:
            error_message = str(e)
            self.logger.error("메모장 추출 오류", e)
            self.root.after(0, lambda: self.status_text.set(f"오류: {error_message[:50]}"))
            self.root.after(0, lambda: messagebox.showerror("오류", f"추출 중 오류:\n{error_message}"))

        finally:
            self.root.after(0, lambda: self.notepad_extract_button.config(state=tk.NORMAL))

    def run(self):
        """프로그램 실행"""
        self.logger.log("메인 루프 시작")
        self.root.mainloop()
        self.logger.log("메인 루프 종료")
        self._shutdown_office_worker()
        self.logger.close()


def check_dependencies():
    """의존성 확인"""
    errors = []

    if not HAS_WIN32COM:
        errors.append("pywin32 패키지 필요 (pip install pywin32)")

    if not HAS_PPTX:
        errors.append("python-pptx 패키지 필요 (pip install python-pptx)")

    if not HAS_OPENPYXL:
        errors.append("openpyxl 패키지 필요 (pip install openpyxl)")

    if not HAS_DOCX:
        errors.append("python-docx 패키지 필요 (pip install python-docx)")

    if not HAS_TKINTERDND:
        errors.append("tkinterdnd2 패키지 필요 (드래그앤드롭용, pip install tkinterdnd2)")

    if errors:
        root = tk.Tk()
        root.withdraw()
        messagebox.showwarning("패키지 필요",
            "다음 패키지를 설치하면 더 많은 기능을 사용할 수 있습니다:\n\n" + "\n".join(errors))

    # pywin32는 필수
    if not HAS_WIN32COM:
        messagebox.showerror("필수 패키지 누락",
            "pywin32 패키지가 필요합니다.\npip install pywin32")
        sys.exit(1)


def write_startup_error(exc):
    """GUI가 뜨기 전 실패도 사용자가 전달할 수 있게 파일로 남긴다."""
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    candidates = [
        os.path.join(os.path.expanduser("~"), "Desktop"),
        os.path.join(os.path.expanduser("~"), "Documents"),
        tempfile.gettempdir(),
    ]
    lines = [
        "=== DocumentExtractor v3 startup error ===",
        f"time: {datetime.datetime.now()}",
        f"executable: {sys.executable}",
        f"argv: {sys.argv}",
        f"cwd: {os.getcwd()}",
        f"frozen: {getattr(sys, 'frozen', False)}",
        f"exception: {type(exc).__name__}: {exc}",
        "",
        traceback.format_exc(),
    ]
    for folder in candidates:
        try:
            if not folder:
                continue
            os.makedirs(folder, exist_ok=True)
            path = os.path.join(folder, f"DocExtractor_Startup_Error_{timestamp}.txt")
            with open(path, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
            return path
        except Exception:
            pass
    return None


if __name__ == "__main__":
    try:
        check_dependencies()
        app = DocumentExtractorV3()
        app.run()
    except Exception as exc:
        error_path = write_startup_error(exc)
        try:
            root = tk.Tk()
            root.withdraw()
            messagebox.showerror(
                "DocumentExtractor 실행 오류",
                "프로그램 시작 중 오류가 발생했습니다.\n\n"
                f"{type(exc).__name__}: {exc}\n\n"
                f"오류 로그: {error_path or '생성 실패'}"
            )
            root.destroy()
        except Exception:
            pass
        sys.exit(1)
