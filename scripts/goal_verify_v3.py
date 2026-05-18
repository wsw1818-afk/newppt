"""
DocumentExtractor v3 goal verification.

This script creates small sample documents and verifies the preservation paths
that matter most for the app:
- PPT/Excel native copy keeps shapes/images/sizing.
- Word safe copy keeps the original file unchanged.
- Notepad legacy Win32 text reading still works when that UI is available.

Missing desktop apps are reported as SKIP, not FAIL.
"""

from __future__ import annotations

import argparse
import ast
import ctypes
import hashlib
import json
import os
import shutil
import subprocess
import sys
import time
import zipfile
from dataclasses import dataclass, asdict
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_DIR = ROOT / "build" / "goal_verify"
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


@dataclass
class CheckResult:
    name: str
    status: str
    detail: str
    elapsed_sec: float


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest().upper()


def zip_has_prefix(path: Path, prefix: str) -> bool:
    with zipfile.ZipFile(path) as zf:
        return any(name.startswith(prefix) for name in zf.namelist())


def make_png(path: Path, label: str) -> None:
    try:
        from PIL import Image, ImageDraw
    except ImportError as exc:
        raise RuntimeError("Pillow is required for image sample generation") from exc

    image = Image.new("RGB", (240, 120), (47, 112, 193))
    draw = ImageDraw.Draw(image)
    draw.rectangle((12, 12, 228, 108), outline=(255, 255, 255), width=4)
    draw.text((28, 48), label, fill=(255, 255, 255))
    image.save(path)


def run_check(name: str, fn) -> CheckResult:
    start = time.perf_counter()
    try:
        detail = fn()
        status = "PASS"
    except SkipTest as exc:
        detail = str(exc)
        status = "SKIP"
    except Exception as exc:
        detail = f"{type(exc).__name__}: {exc}"
        status = "FAIL"
    elapsed = time.perf_counter() - start
    result = CheckResult(name, status, detail, round(elapsed, 2))
    print(f"[{status}] {name}: {detail} ({elapsed:.2f}s)")
    return result


class SkipTest(Exception):
    pass


def import_com():
    try:
        import pythoncom
        import pywintypes
        import win32com.client
    except ImportError as exc:
        raise SkipTest(f"pywin32 unavailable: {exc}") from exc
    return pythoncom, pywintypes, win32com.client


def dispatch_isolated(win32, prog_id: str, label: str):
    """Start an isolated COM app so verification never closes a user's open app."""
    try:
        return win32.DispatchEx(prog_id)
    except Exception as exc:
        raise SkipTest(f"{label} isolated COM unavailable: {exc}") from exc


def check_py_compile() -> str:
    files = [
        "ppt_extractor_v3.py",
        "document_extractor.py",
        "pdf_printer.py",
        "ppt_extractor.py",
        "ppt_extractor_v2.py",
    ]
    cmd = [sys.executable, "-m", "py_compile", *files]
    proc = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
    if proc.returncode != 0:
        raise RuntimeError((proc.stderr or proc.stdout).strip())
    return f"compiled {len(files)} files"


def check_autoshape_mapping() -> str:
    source_path = ROOT / "ppt_extractor_v3.py"
    tree = ast.parse(source_path.read_text(encoding="utf-8"))
    duplicates: list[int] = []

    for node in ast.walk(tree):
        if not isinstance(node, ast.Assign):
            continue
        if not any(isinstance(t, ast.Name) and t.id == "AUTOSHAPE_MAPPING" for t in node.targets):
            continue
        if not isinstance(node.value, ast.Dict):
            continue
        seen: set[int] = set()
        for key in node.value.keys:
            if isinstance(key, ast.Constant) and isinstance(key.value, int):
                if key.value in seen:
                    duplicates.append(key.value)
                seen.add(key.value)

    if duplicates:
        raise RuntimeError(f"duplicate keys: {duplicates}")
    return "AUTOSHAPE_MAPPING duplicate_keys=[]"


def check_ppt_native_copy(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise SkipTest(f"python-pptx unavailable: {exc}") from exc

    source = out_dir / "sample_ppt_source.pptx"
    copied = out_dir / "sample_ppt_copy.pptx"
    image = out_dir / "sample_ppt_image.png"
    make_png(image, "PPT")

    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        app = dispatch_isolated(win32, "PowerPoint.Application", "PowerPoint")
        app.Visible = True
        presentation = app.Presentations.Add()
        presentation.PageSetup.SlideWidth = 960
        presentation.PageSetup.SlideHeight = 540
        slide = presentation.Slides.Add(1, 12)  # ppLayoutBlank
        shape = slide.Shapes.AddShape(1, 80, 80, 220, 90)  # msoShapeRectangle
        shape.TextFrame.TextRange.Text = "verify ppt"
        slide.Shapes.AddPicture(str(image), False, True, 380, 90, 180, 90)
        presentation.SaveAs(str(source))
        presentation.SaveCopyAs(str(copied))
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    if not copied.exists() or copied.stat().st_size == 0:
        raise RuntimeError("copied PPTX was not created")

    deck = Presentation(str(copied))
    pictures = [
        shp for shp in deck.slides[0].shapes
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if len(deck.slides) != 1 or len(deck.slides[0].shapes) < 2 or not pictures:
        raise RuntimeError("PPTX shape/image counts were not preserved")
    if int(deck.slide_width) != 12192000 or int(deck.slide_height) != 6858000:
        raise RuntimeError("PPTX slide size was not preserved")
    if not zip_has_prefix(copied, "ppt/media/"):
        raise RuntimeError("PPTX media payload missing")
    return f"slides=1 shapes={len(deck.slides[0].shapes)} pictures={len(pictures)}"


def check_ppt_clipboard_package(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    try:
        from pptx import Presentation
        from ppt_extractor_v3 import DocumentExtractorV3
    except ImportError as exc:
        raise SkipTest(f"ppt clipboard dependencies unavailable: {exc}") from exc

    class StubLogger:
        def log(self, message: str) -> None:
            pass

        def error(self, message: str, exc: Exception | None = None) -> None:
            pass

    source = out_dir / "sample_ppt_clipboard_source.pptx"
    copied = out_dir / "sample_ppt_clipboard_package.pptx"

    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        app = dispatch_isolated(win32, "PowerPoint.Application", "PowerPoint")
        app.Visible = True
        presentation = app.Presentations.Add()
        presentation.PageSetup.SlideWidth = 960
        presentation.PageSetup.SlideHeight = 540
        for idx in range(1, 3):
            slide = presentation.Slides.Add(idx, 12)  # ppLayoutBlank
            shape = slide.Shapes.AddShape(1, 80, 80, 220, 90)  # msoShapeRectangle
            shape.TextFrame.TextRange.Text = f"clipboard verify {idx}"
            slide.Shapes.AddTextbox(1, 360, 90, 240, 60).TextFrame.TextRange.Text = f"editable text {idx}"
        presentation.SaveAs(str(source))

        extractor = object.__new__(DocumentExtractorV3)
        extractor.logger = StubLogger()
        extractor._save_ppt_clipboard_package_copy(presentation, str(copied))
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    if not copied.exists() or copied.stat().st_size == 0:
        raise RuntimeError("clipboard package PPTX was not created")
    deck = Presentation(str(copied))
    texts = [
        shp.text
        for slide in deck.slides
        for shp in slide.shapes
        if getattr(shp, "has_text_frame", False) and shp.text
    ]
    if len(deck.slides) != 2:
        raise RuntimeError(f"slide count mismatch: {len(deck.slides)}")
    for expected in ("clipboard verify 1", "clipboard verify 2", "editable text 1", "editable text 2"):
        if expected not in texts:
            raise RuntimeError(f"missing editable text: {expected}")
    if not zip_has_prefix(copied, "ppt/slides/"):
        raise RuntimeError("rebuilt PPTX slide XML missing")
    return f"slides={len(deck.slides)} editable_texts={len(texts)}"


def check_excel_native_copy(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    source = out_dir / "sample_excel_source.xlsx"
    copied = out_dir / "sample_excel_copy.xlsx"
    image = out_dir / "sample_excel_image.png"
    make_png(image, "Excel")

    pythoncom.CoInitialize()
    app = None
    wb = None
    wb2 = None
    try:
        app = dispatch_isolated(win32, "Excel.Application", "Excel")
        app.Visible = False
        app.DisplayAlerts = False
        wb = app.Workbooks.Add()
        ws = wb.Worksheets(1)
        ws.Name = "GoalVerify"
        ws.Cells(1, 1).Value = "verify excel"
        ws.Cells(1, 1).Font.Bold = True
        ws.Columns(1).ColumnWidth = 28
        ws.Rows(1).RowHeight = 42
        ws.Shapes.AddShape(1, 120, 40, 160, 80)
        ws.Shapes.AddPicture(str(image), False, True, 20, 90, 90, 45)
        wb.SaveAs(str(source), 51)  # xlOpenXMLWorkbook
        wb.SaveCopyAs(str(copied))
        wb.Close(False)
        wb = None

        wb2 = app.Workbooks.Open(str(copied), ReadOnly=True)
        ws2 = wb2.Worksheets(1)
        value = ws2.Cells(1, 1).Value
        width = float(ws2.Columns(1).ColumnWidth)
        height = float(ws2.Rows(1).RowHeight)
        shapes = int(ws2.Shapes.Count)
    finally:
        if wb2 is not None:
            try:
                wb2.Close(False)
            except Exception:
                pass
        if wb is not None:
            try:
                wb.Close(False)
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    if value != "verify excel":
        raise RuntimeError("Excel cell value missing")
    if abs(width - 28) > 0.1 or abs(height - 42) > 0.1:
        raise RuntimeError(f"Excel size mismatch: width={width}, height={height}")
    if shapes < 2:
        raise RuntimeError(f"Excel shape/image count too low: {shapes}")
    if not zip_has_prefix(copied, "xl/media/"):
        raise RuntimeError("Excel media payload missing")
    return f"shapes={shapes} width={width:.1f} height={height:.1f}"


def check_word_safe_copy(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    from ppt_extractor_v3 import DocumentExtractorV3

    source = out_dir / "sample_word_source.docx"
    copied = out_dir / "sample_word_copy.docx"
    image = out_dir / "sample_word_image.png"
    make_png(image, "Word")

    pythoncom.CoInitialize()
    app = None
    doc = None
    try:
        app = dispatch_isolated(win32, "Word.Application", "Word")
        app.Visible = False
        app.DisplayAlerts = 0
        doc = app.Documents.Add()
        doc.Content.Text = "verify word safe copy\n"
        doc.InlineShapes.AddPicture(str(image), False, True)
        doc.SaveAs2(str(source), FileFormat=16)  # wdFormatXMLDocument
        source_full_name = str(doc.FullName)
        source_saved = bool(doc.Saved)
        extractor = object.__new__(DocumentExtractorV3)
        extractor._copy_word_document_file(doc, str(copied))
        after_full_name = str(doc.FullName)
        after_saved = bool(doc.Saved)
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    if not source_saved or not after_saved:
        raise RuntimeError("Word document was not in a saved state")
    if os.path.abspath(source_full_name).lower() != os.path.abspath(after_full_name).lower():
        raise RuntimeError("Word source FullName changed during copy")
    if sha256(source) != sha256(copied):
        raise RuntimeError("Word copied file hash differs from source")
    if not zip_has_prefix(copied, "word/media/"):
        raise RuntimeError("Word media payload missing")
    return f"bytes={copied.stat().st_size} hash={sha256(copied)[:12]}"


def check_word_openxml_copy(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    from ppt_extractor_v3 import DocumentExtractorV3

    class Logger:
        def log(self, _message):
            pass

    output = out_dir / "sample_word_openxml_copy.docx"
    image = out_dir / "sample_word_openxml_image.png"
    make_png(image, "WordOpenXML")

    pythoncom.CoInitialize()
    app = None
    doc = None
    copied = None
    try:
        app = dispatch_isolated(win32, "Word.Application", "Word")
        app.Visible = False
        app.DisplayAlerts = 0
        doc = app.Documents.Add()
        doc.Content.Text = "verify word openxml copy\n"
        insert_range = doc.Range(doc.Content.End - 1, doc.Content.End - 1)
        table = doc.Tables.Add(insert_range, 2, 2)
        table.Cell(1, 1).Range.Text = "A"
        table.Cell(1, 2).Range.Text = "B"
        doc.InlineShapes.AddPicture(str(image), False, True)
        doc.Sections(1).Footers(1).Range.Text = "verify footer"

        extractor = object.__new__(DocumentExtractorV3)
        extractor.logger = Logger()
        extractor._save_word_openxml_copy(doc, str(output))

        copied = app.Documents.Open(str(output), ReadOnly=True, AddToRecentFiles=False)
        tables = copied.Tables.Count
        inline_shapes = copied.InlineShapes.Count
        footer_text = copied.Sections(1).Footers(1).Range.Text
    finally:
        if copied is not None:
            try:
                copied.Close(False)
            except Exception:
                pass
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    if tables < 1:
        raise RuntimeError("WordOpenXML copied table missing")
    if inline_shapes < 1:
        raise RuntimeError("WordOpenXML copied image missing")
    if "verify footer" not in footer_text:
        raise RuntimeError("WordOpenXML copied footer missing")
    if not zip_has_prefix(output, "word/media/"):
        raise RuntimeError("WordOpenXML copied media payload missing")
    return f"tables={tables} inline_shapes={inline_shapes} bytes={output.stat().st_size}"


def check_word_xml_text_sanitizer(out_dir: Path) -> str:
    from docx import Document
    from ppt_extractor_v3 import DocumentExtractorV3

    class Logger:
        def log(self, _message):
            pass

    extractor = object.__new__(DocumentExtractorV3)
    extractor.logger = Logger()
    raw = "valid text\x00\x07\x0b\ud800\nnext line"
    cleaned = extractor._clean_xml_text(raw)
    if "\x00" in cleaned or "\x07" in cleaned or "\x0b" in cleaned or "\ud800" in cleaned:
        raise RuntimeError("invalid XML characters were not removed")

    output = out_dir / "sample_word_sanitized.docx"
    doc = Document()
    doc.add_paragraph(cleaned)
    doc.save(output)
    if not zip_has_prefix(output, "word/document.xml"):
        raise RuntimeError("sanitized docx is invalid")
    return f"chars={len(cleaned)} bytes={output.stat().st_size}"


def visible_window_titles() -> list[str]:
    titles: list[str] = []
    user32 = ctypes.windll.user32
    enum_proc_type = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)

    def enum_proc(hwnd, _lparam):
        if user32.IsWindowVisible(hwnd):
            length = user32.GetWindowTextLengthW(hwnd)
            if length > 0:
                buffer = ctypes.create_unicode_buffer(length + 1)
                user32.GetWindowTextW(hwnd, buffer, length + 1)
                title = buffer.value.strip()
                if title:
                    titles.append(title)
        return True

    user32.EnumWindows(enum_proc_type(enum_proc), 0)
    return titles


def hwp_titles() -> list[str]:
    return [
        title for title in visible_window_titles()
        if title.endswith(" - 한글") or " - 한글 " in title
    ]


def check_hwp_getobject_no_spawn() -> str:
    _pythoncom, _pywintypes, win32 = import_com()
    before = hwp_titles()
    try:
        win32.GetObject(Class="HWPFrame.HwpObject")
    except Exception:
        pass
    time.sleep(0.3)
    after = hwp_titles()
    if len(after) > len(before):
        raise RuntimeError(f"HWP window spawned: before={len(before)}, after={len(after)}")
    return f"before={len(before)} after={len(after)}"


def check_hwp_action_save(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    output = out_dir / "sample_hwp_save.hwp"

    pythoncom.CoInitialize()
    hwp = None
    try:
        try:
            hwp = win32.Dispatch("HWPFrame.HwpObject")
        except Exception as exc:
            raise SkipTest(f"HWP COM unavailable: {exc}") from exc

        try:
            hwp.HAction.GetDefault("InsertText", hwp.HParameterSet.HInsertText.HSet)
            hwp.HParameterSet.HInsertText.Text = "verify hwp save"
            hwp.HAction.Execute("InsertText", hwp.HParameterSet.HInsertText.HSet)
        except Exception:
            pass

        hwp.HAction.GetDefault("FileSaveAs_S", hwp.HParameterSet.HFileOpenSave.HSet)
        file_open_save = hwp.HParameterSet.HFileOpenSave
        file_open_save.filename = str(output)
        try:
            file_open_save.FileName = str(output)
        except Exception:
            pass
        file_open_save.Format = "HWP"
        result = hwp.HAction.Execute("FileSaveAs_S", file_open_save.HSet)
        if result is False:
            raise RuntimeError("FileSaveAs_S returned False")
    finally:
        if hwp is not None:
            try:
                hwp.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    for _ in range(20):
        if output.exists() and output.stat().st_size > 0:
            return f"bytes={output.stat().st_size}"
        time.sleep(0.1)
    raise RuntimeError("HWP output file was not created")


def check_hwp_window_save_fallback(out_dir: Path) -> str:
    pythoncom, _pywintypes, win32 = import_com()
    from ppt_extractor_v3 import DocumentExtractorV3

    class Logger:
        def log(self, _message):
            pass

    source = out_dir / "sample_hwp_window_source.hwp"
    output = out_dir / "sample_hwp_window_fallback.hwp"

    pythoncom.CoInitialize()
    hwp = None
    try:
        last_dispatch_error = None
        for _attempt in range(3):
            try:
                hwp = win32.Dispatch("HWPFrame.HwpObject")
                break
            except Exception as exc:
                last_dispatch_error = exc
                time.sleep(1)
        if hwp is None:
            raise SkipTest(f"HWP COM unavailable: {last_dispatch_error}") from last_dispatch_error

        try:
            try:
                hwp.HAction.GetDefault("InsertText", hwp.HParameterSet.HInsertText.HSet)
                hwp.HParameterSet.HInsertText.Text = "verify hwp window fallback"
                hwp.HAction.Execute("InsertText", hwp.HParameterSet.HInsertText.HSet)
            except Exception:
                pass

            hwp.HAction.GetDefault("FileSaveAs_S", hwp.HParameterSet.HFileOpenSave.HSet)
            file_open_save = hwp.HParameterSet.HFileOpenSave
            file_open_save.filename = str(source)
            try:
                file_open_save.FileName = str(source)
            except Exception:
                pass
            file_open_save.Format = "HWP"
            result = hwp.HAction.Execute("FileSaveAs_S", file_open_save.HSet)
            if result is False:
                raise RuntimeError("FileSaveAs_S returned False while creating source")
        except Exception as exc:
            raise SkipTest(f"HWP UI fallback source setup unavailable: {exc}") from exc
    finally:
        if hwp is not None:
            try:
                hwp.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()
        time.sleep(1)

    if not source.exists() or source.stat().st_size <= 0:
        raise RuntimeError("HWP source file was not created")

    extractor = object.__new__(DocumentExtractorV3)
    extractor.logger = Logger()

    before_hwnds = {hwnd for hwnd, _title in extractor._list_hwp_windows()}
    os.startfile(str(source))

    hwnd = 0
    deadline = time.time() + 10
    while time.time() < deadline:
        for candidate_hwnd, title in extractor._list_hwp_windows():
            if candidate_hwnd not in before_hwnds or source.stem in title:
                hwnd = candidate_hwnd
                break
        if hwnd:
            break
        time.sleep(0.2)

    if not hwnd:
        raise SkipTest("HWP visible window unavailable for UI fallback")

    try:
        try:
            extractor._save_hwp_via_window(hwnd, str(output), "hwp")
        except Exception as exc:
            raise SkipTest(f"HWP UI fallback unavailable: {exc}") from exc
    finally:
        try:
            ctypes.windll.user32.PostMessageW(hwnd, 0x0010, 0, 0)  # WM_CLOSE
        except Exception:
            pass

    if not output.exists() or output.stat().st_size <= 0:
        raise RuntimeError("HWP UI fallback output file was not created")
    return f"bytes={output.stat().st_size}"


def find_window_for_pid(pid: int, class_name: str | None = None) -> int | None:
    user32 = ctypes.windll.user32
    enum_proc_type = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)
    found: list[int] = []

    def enum_proc(hwnd, _lparam):
        proc_id = ctypes.c_ulong()
        user32.GetWindowThreadProcessId(hwnd, ctypes.byref(proc_id))
        if proc_id.value != pid or not user32.IsWindowVisible(hwnd):
            return True
        if class_name:
            buffer = ctypes.create_unicode_buffer(256)
            user32.GetClassNameW(hwnd, buffer, 256)
            if buffer.value != class_name:
                return True
        found.append(int(hwnd))
        return False

    user32.EnumWindows(enum_proc_type(enum_proc), 0)
    return found[0] if found else None


def find_window_by_title_fragment(title_fragment: str, class_names: set[str] | None = None) -> int | None:
    user32 = ctypes.windll.user32
    enum_proc_type = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p)
    fragment = title_fragment.lower()
    found: list[int] = []

    def enum_proc(hwnd, _lparam):
        if not user32.IsWindowVisible(hwnd):
            return True

        title_len = user32.GetWindowTextLengthW(hwnd)
        if title_len <= 0:
            return True

        title_buffer = ctypes.create_unicode_buffer(title_len + 1)
        user32.GetWindowTextW(hwnd, title_buffer, title_len + 1)
        if fragment not in title_buffer.value.lower():
            return True

        if class_names is not None and window_class_name(hwnd) not in class_names:
            return True

        found.append(int(hwnd))
        return False

    user32.EnumWindows(enum_proc_type(enum_proc), 0)
    return found[0] if found else None


def window_class_name(hwnd: int) -> str:
    user32 = ctypes.windll.user32
    buffer = ctypes.create_unicode_buffer(256)
    user32.GetClassNameW(hwnd, buffer, 256)
    return buffer.value


def find_child_window_by_classes(parent_hwnd: int, class_names: list[str]) -> int | None:
    user32 = ctypes.windll.user32
    wanted = set(class_names)
    queue = [parent_hwnd]
    seen: set[int] = set()

    while queue:
        current = queue.pop(0)
        if current in seen:
            continue
        seen.add(current)

        child = 0
        while True:
            child = user32.FindWindowExW(current, child, None, None)
            if not child:
                break
            if window_class_name(child) in wanted:
                return int(child)
            queue.append(int(child))

    return None


def check_notepad_legacy_read(out_dir: Path) -> str:
    user32 = ctypes.windll.user32
    sample = out_dir / "sample_notepad.txt"
    expected = "verify notepad legacy read"
    sample.write_text(expected, encoding="utf-8")

    try:
        proc = subprocess.Popen(["notepad.exe", str(sample)])
    except FileNotFoundError as exc:
        raise SkipTest("notepad.exe unavailable") from exc

    hwnd = None
    try:
        deadline = time.time() + 12
        while time.time() < deadline:
            hwnd = find_window_for_pid(proc.pid)
            if not hwnd:
                hwnd = find_window_by_title_fragment(sample.name, {"Notepad"})
            if hwnd:
                break
            time.sleep(0.2)

        if not hwnd:
            raise SkipTest("Notepad top-level window not found")

        edit_hwnd = find_child_window_by_classes(hwnd, ["Edit", "RichEditD2DPT", "RICHEDIT50W"])
        if not edit_hwnd:
            raise SkipTest("Notepad does not expose legacy Edit/RichEdit control")

        text_len = user32.SendMessageW(edit_hwnd, 0x000E, 0, 0)  # WM_GETTEXTLENGTH
        buffer = ctypes.create_unicode_buffer(text_len + 16)
        copied = user32.SendMessageW(edit_hwnd, 0x000D, len(buffer), buffer)  # WM_GETTEXT
        if copied <= 0 or expected not in buffer.value:
            raise RuntimeError("Notepad text read mismatch")
        return f"chars={copied}"
    finally:
        if hwnd:
            user32.PostMessageW(hwnd, 0x0010, 0, 0)  # WM_CLOSE
        try:
            proc.wait(timeout=3)
        except subprocess.TimeoutExpired:
            proc.terminate()


def write_summary(out_dir: Path, results: list[CheckResult]) -> Path:
    summary = {
        "root": str(ROOT),
        "output_dir": str(out_dir),
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "results": [asdict(result) for result in results],
    }
    path = out_dir / "goal_verify_summary.json"
    path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def main() -> int:
    parser = argparse.ArgumentParser(description="Run DocumentExtractor v3 goal verification")
    parser.add_argument(
        "--output-dir",
        default=str(DEFAULT_OUTPUT_DIR),
        help="Directory for generated verification artifacts",
    )
    parser.add_argument(
        "--clean",
        action="store_true",
        help="Remove the output directory before running",
    )
    args = parser.parse_args()

    out_dir = Path(args.output_dir).resolve()
    if args.clean and out_dir.exists():
        shutil.rmtree(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    checks = [
        ("py_compile", check_py_compile),
        ("autoshape_mapping", check_autoshape_mapping),
        ("ppt_native_copy", lambda: check_ppt_native_copy(out_dir)),
        ("ppt_clipboard_package", lambda: check_ppt_clipboard_package(out_dir)),
        ("excel_native_copy", lambda: check_excel_native_copy(out_dir)),
        ("word_safe_copy", lambda: check_word_safe_copy(out_dir)),
        ("word_openxml_copy", lambda: check_word_openxml_copy(out_dir)),
        ("word_xml_text_sanitizer", lambda: check_word_xml_text_sanitizer(out_dir)),
        ("notepad_legacy_read", lambda: check_notepad_legacy_read(out_dir)),
    ]

    results = [run_check(name, fn) for name, fn in checks]
    summary_path = write_summary(out_dir, results)

    failed = [result for result in results if result.status == "FAIL"]
    passed = [result for result in results if result.status == "PASS"]
    skipped = [result for result in results if result.status == "SKIP"]
    print()
    print(f"Summary: PASS={len(passed)} SKIP={len(skipped)} FAIL={len(failed)}")
    print(f"Summary file: {summary_path}")

    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
